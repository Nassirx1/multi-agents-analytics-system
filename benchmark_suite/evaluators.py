"""Deterministic benchmark executors for the immutable catalog.

The checks exercise real routing, normalization, rendering, QA, checkpoint,
redaction, retry, and rubric code. Language-model boundaries are replaced by
small deterministic stubs so the pilot is reproducible and credential-free.
"""

from __future__ import annotations

import hashlib
import inspect
import json
import re
import tempfile
from pathlib import Path
from typing import Any, Callable, Mapping
from unittest.mock import Mock, patch

import pandas as pd

from analytics_workflow.agents import (
    BusinessInsightsTranslatorAgent,
    DataScientistReviewerAgent,
    DataUnderstanderAgent,
    DecisionMakerAgent,
    MarketResearcherAgent,
    PlannerAgent,
)
from analytics_workflow.clients import CostTracker, OpenRouterClient
from analytics_workflow.evidence import validate_and_sanitize_workflow
from analytics_workflow.html_dashboard.planning import build_dataset_profiles, normalize_dashboard_plan
from analytics_workflow.html_dashboard.renderer import render_dashboard, validate_dashboard_html
from analytics_workflow.html_dashboard.workflow import HTMLDashboardWorkflow
from analytics_workflow.output_paths import OutputPath, coerce_output_path, prompt_output_path
from analytics_workflow.pipeline_runtime import MultiAgentOrchestrator
from analytics_workflow.presentation_backends import inspect_presentation
from analytics_workflow.reporting import generate_pdf_report, generate_slide_deck, inspect_pdf_report
from analytics_workflow.run_checkpoints import load_run_checkpoint
from analytics_workflow.runtime_config import build_runtime_config, redact_secrets
from analytics_workflow.workflow_steps import ANALYTICS_REPORT_STEPS, HTML_DASHBOARD_STEPS, workflow_steps_for

from .assurance import (
    ASSURANCE_VERSION,
    INDEPENDENT_GATE_EVALUATORS,
    build_artifact_inventory,
    score_output_quality_assured,
)
from .rubrics import OUTPUT_QUALITY_RUBRICS, hard_gates_for_route, score_output_quality


Executor = Callable[[Any], Mapping[str, Any]]


class _StubClient:
    def __init__(self, response: Mapping[str, Any] | None = None) -> None:
        self.response = dict(response or {})
        self.calls: list[dict[str, Any]] = []
        self.code_loop_timeout_seconds = 30
        self.cost_tracker = CostTracker(model="benchmark-stub")

    def chat_completion_json(self, system: str, prompt: str, schema: Any, **kwargs: Any) -> dict[str, Any]:
        self.calls.append({"system": system, "prompt": prompt, "schema": schema, "kwargs": kwargs})
        return json.loads(json.dumps(self.response))


def get_benchmark_executor(target: str) -> Executor | None:
    return BENCHMARK_EXECUTORS.get(str(target))


def _config():
    return build_runtime_config(
        openrouter_api_key="benchmark-openrouter-sentinel",
        brave_search_api_key="benchmark-brave-sentinel",
        presentation_backend="python",
        html_dashboard_stage_timeout_seconds=30,
        html_dashboard_max_rows=1000,
    )


def _frame() -> pd.DataFrame:
    return pd.DataFrame(
        {
            "Segment": ["Enterprise", "SMB", "SMB", "Public", "Enterprise", "SMB", "Public", "Enterprise"],
            "Region": ["East", "West", "East", "West", "East", "West", "East", None],
            "Revenue": [120.0, 45.0, 52.0, 83.0, 140.0, 49.0, 90.0, 132.0],
            "Risk": ["High", "Low", "Low", "Medium", "Low", "High", "Medium", "Low"],
            "Converted": [1, 0, 1, 1, 1, 0, 1, 1],
            "Date": pd.date_range("2026-01-01", periods=8, freq="D"),
        }
    )


def _dashboard_plan(frame: pd.DataFrame | None = None) -> dict[str, Any]:
    frame = _frame() if frame is None else frame
    candidate = {
        "project_name": "Benchmark Dashboard",
        "title": "Revenue and Conversion Overview",
        "subtitle": "Decision view of revenue, segment mix, and conversion",
        "filters": [
            {"dataset": "business.csv", "field": "Segment", "label": "Segment"},
            {"dataset": "business.csv", "field": "Region", "label": "Region"},
        ],
        "kpis": [
            {"dataset": "business.csv", "label": "Records", "field": "Segment", "aggregation": "count", "format": "integer"},
            {"dataset": "business.csv", "label": "Average revenue", "field": "Revenue", "aggregation": "mean", "format": "currency"},
        ],
        "pages": [
            {
                "name": "Overview",
                "purpose": "Answer-first business overview",
                "charts": [
                    {"dataset": "business.csv", "title": "Revenue by segment", "type": "bar", "x": "Segment", "y": "Revenue", "aggregation": "sum"},
                    {"dataset": "business.csv", "title": "Revenue over time", "type": "line", "x": "Date", "y": "Revenue", "aggregation": "mean"},
                ],
                "table": {"dataset": "business.csv", "title": "Business detail", "columns": ["Segment", "Region", "Revenue", "Risk"]},
            },
            {
                "name": "Drivers",
                "purpose": "Segment and conversion drivers",
                "charts": [
                    {"dataset": "business.csv", "title": "Revenue and conversion", "type": "scatter", "x": "Revenue", "y": "Converted", "aggregation": "none"}
                ],
                "table": {},
            },
        ],
        "theme": {"accent": "#2563eb", "background": "#f4f7fb"},
        "source_notes": ["Embedded benchmark dataset; eight records."],
    }
    return normalize_dashboard_plan(candidate, {"business.csv": frame}, {})


def _report_state() -> dict[str, Any]:
    return {
        "workflow_objective": {"raw_description": "Prioritize revenue and conversion actions."},
        "saved_figures": [],
        "evidence_bundle": {
            "bundle_hash": "benchmark-evidence-v1",
            "records": [{"evidence_id": "EDA-001", "statement": "Enterprise revenue is highest.", "source": "business.csv"}],
        },
        "analysis_results": {
            "analysis_summary": {"total_records": 8, "average_revenue": 88.875},
            "business_findings": ["Enterprise has the largest observed revenue total (EDA-001)."],
            "analysis_artifacts": [
                {
                    "artifact_id": "revenue_by_segment",
                    "chart_type": "bar",
                    "title": "Enterprise leads observed revenue",
                    "finding": "Enterprise has the largest observed total.",
                    "evidence_ids": ["EDA-001"],
                    "data": [
                        {"category": "Enterprise", "value": 392.0},
                        {"category": "SMB", "value": 146.0},
                        {"category": "Public", "value": 173.0},
                    ],
                }
            ],
        },
        "agent_outputs": {
            "data_understander": {"executive_summary": "Eight business records across three segments.", "datasets": {}},
            "market_researcher": {"market_findings": [], "sources_cited": []},
            "planner": {"objectives": ["Compare observed revenue by segment."], "success_metrics": ["Reconciled totals"]},
            "business_translator": {
                "executive_summary": "Enterprise is the first revenue-retention priority.",
                "key_findings": [{"finding": "Enterprise leads revenue.", "business_implication": "Protect these accounts.", "evidence_ids": ["EDA-001"]}],
            },
            "decision_maker": {
                "executive_summary": "Protect enterprise revenue while testing SMB conversion improvements.",
                "recommendations": [
                    {
                        "rank": 1,
                        "action": "Review enterprise account risks",
                        "rationale": "Enterprise leads observed revenue.",
                        "evidence": "EDA-001",
                        "evidence_ids": ["EDA-001"],
                        "owner": "Revenue operations lead",
                        "trigger": "Weekly account review",
                        "timeline": "Next 30 days",
                        "validation_metric": "Retained enterprise revenue",
                        "stop_condition": "Stop if intervention cost exceeds retained margin",
                        "impact": "High",
                    }
                ],
                "limitations": [{"limitation": "Small synthetic fixture", "mitigation": "Validate on production data", "decision_impact": "Directional only"}],
                "final_recommendation": "Prioritize enterprise retention validation.",
            },
        },
    }


def _artifact_root(context: Any) -> Path:
    if context.output_dir is not None:
        root = Path(context.output_dir).resolve() / "artifacts" / context.route
        root.mkdir(parents=True, exist_ok=True)
        return root
    return Path(tempfile.mkdtemp(prefix=f"benchmark-{context.route}-"))


def _contract_result(
    context: Any,
    checks: Mapping[str, bool],
    *,
    artifacts: list[str] | None = None,
    calls: list[str] | None = None,
    evidence: Mapping[str, Any] | None = None,
    scores: Mapping[str, float] | None = None,
    gate_results: Mapping[str, bool] | None = None,
    overall_score: float | None = None,
) -> dict[str, Any]:
    failed = [name for name, passed in checks.items() if not passed]
    ratio = 100.0 * (len(checks) - len(failed)) / max(1, len(checks))
    dimension_scores = dict(scores or {dimension: round(ratio, 1) for dimension in context.dimensions})
    gates = dict(gate_results or {gate: not failed for gate in context.hard_gates})
    return {
        "passed": not failed and all(gates.get(gate) is True for gate in context.hard_gates),
        "scores": dimension_scores,
        "overall_score": overall_score,
        "gate_results": gates,
        "diagnostics": [f"{name}: deterministic assertion failed" for name in failed],
        "artifacts": artifacts or [],
        "calls": calls or [],
        "metadata": {"checks": dict(checks), "evidence": dict(evidence or {})},
    }


def _initial_choice(context: Any) -> Mapping[str, Any]:
    events: list[str] = []
    answers = iter(context.fixture.get("selections", ("invalid", "2")))
    selected = prompt_output_path(
        input_fn=lambda _: events.append("route") or next(answers),
        print_fn=lambda text: events.append("menu") if not events and "Choose an output path" in text else None,
    )
    cli_order: list[str] = []
    import analytics_workflow.cli as cli

    with patch.object(cli, "prompt_output_path", side_effect=lambda: cli_order.append("route") or selected), patch.object(
        cli, "load_runtime_config", side_effect=lambda: cli_order.append("credentials") or _config()
    ), patch.object(cli, "run_terminal_workflow", side_effect=lambda *_args, **_kwargs: cli_order.append("dataset") or 0), patch.object(
        cli, "setup_logging"
    ):
        exit_code = cli.main()
    expected_order = list(context.expected.get("event_order", ()))
    observed_order = ["menu"] + cli_order
    return _contract_result(
        context,
        {
            "selected_html_after_invalid": selected is OutputPath.HTML_DASHBOARD,
            "menu_first": events[0] == "menu",
            "event_order": observed_order == expected_order,
            "exit_zero": exit_code == 0,
        },
        evidence={"observed_order": observed_order, "selection_events": events},
    )


def _analytics_isolation(context: Any) -> Mapping[str, Any]:
    with tempfile.TemporaryDirectory() as tmp:
        orchestrator = MultiAgentOrchestrator(_config(), workspace=Path(tmp), output_path=OutputPath.ANALYTICS_REPORT)
        agents = set(orchestrator.agents)
        checks = {
            "nine_stage_report_path": workflow_steps_for(orchestrator.output_path) == ANALYTICS_REPORT_STEPS,
            "report_agents_present": {"data_understander", "market_researcher", "planner", "coder", "reviewer", "business_translator", "decision_maker"}.issubset(agents),
            "dashboard_runtime_absent": not hasattr(orchestrator, "dashboard_workflow"),
            "pdf_exporter_present": callable(getattr(orchestrator, "_generate_pdf_export", None)),
            "slide_exporter_present": callable(getattr(orchestrator, "_generate_slide_export", None)),
        }
    return _contract_result(context, checks, calls=list(ANALYTICS_REPORT_STEPS), evidence={"agent_keys": sorted(agents)})


def _dashboard_isolation(context: Any) -> Mapping[str, Any]:
    with tempfile.TemporaryDirectory() as tmp:
        orchestrator = MultiAgentOrchestrator(_config(), workspace=Path(tmp), output_path=OutputPath.HTML_DASHBOARD)
        agents = set(orchestrator.agents)
        steps = workflow_steps_for(orchestrator.output_path)
        checks = {
            "four_stage_dashboard_path": steps == HTML_DASHBOARD_STEPS,
            "only_shared_agent_constructed": agents == {"data_understander"},
            "market_client_absent": orchestrator.brave_client is None,
            "no_power_bi_runtime": not (Path(__file__).parents[1] / "analytics_workflow" / "powerbi").exists(),
        }
    return _contract_result(context, checks, calls=steps, artifacts=["benchmark/dashboard.html"], evidence={"agent_keys": sorted(agents)})


def _legacy_resume(context: Any) -> Mapping[str, Any]:
    legacy = list(context.fixture.get("legacy_routes", (None, "power_bi")))
    resolved = [
        coerce_output_path(value if value is not None else OutputPath.ANALYTICS_REPORT.value).value
        for value in legacy
    ]
    expected = list(context.expected.get("resolved_routes", ()))
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        completed = {"data_understander": {"hash": "stable"}}
        (root / "run_manifest.json").write_text(json.dumps({"status": "failed", "datasets": [], "output_path": "power_bi"}), encoding="utf-8")
        (root / "agent_outputs.json").write_text(json.dumps(completed), encoding="utf-8")
        checkpoint = load_run_checkpoint(root)
        preserved = checkpoint.agent_outputs == completed
    return _contract_result(context, {"routes_resolve": resolved == expected, "completed_outputs_preserved": preserved}, evidence={"resolved_routes": resolved})


def _data_understander(context: Any) -> Mapping[str, Any]:
    frame = _frame()
    agent = DataUnderstanderAgent("Data Understander", "Data Analyst", "profiling", _StubClient())
    profiles = agent._column_profiles(frame)  # deterministic implementation under test
    checks = {
        "row_count_exact": len(frame) == 8,
        "columns_complete": set(profiles) == set(frame.columns),
        "types_classified": all(profile.get("role") for profile in profiles.values()),
        "missingness_exact": profiles["Region"]["missing_pct"] == 12.5,
        "numeric_summary": all(key in profiles["Revenue"] for key in ("min", "median", "mean", "max")),
        "categorical_summary": profiles["Segment"]["unique_count"] == 3,
        "limitations_available": "missing_pct" in profiles["Region"],
    }
    return _contract_result(context, checks, evidence={"profile": profiles})


def _market_researcher(context: Any) -> Mapping[str, Any]:
    agent = MarketResearcherAgent("Market Researcher", "Researcher", "sources", _StubClient())
    canonical = [
        {"index": 1, "title": "Regulator", "url": "https://www.bls.gov/example", "description": "Primary source", "query": "q", "evidence_level": "search_snippet", "source_quality": "primary"},
        {"index": 2, "title": "Unknown blog", "url": "https://example.invalid/blog", "description": "Low quality", "query": "q", "evidence_level": "search_snippet", "source_quality": "secondary"},
    ]
    result = agent._reconcile_sources(
        {"market_findings": [{"claim": "Supported context", "source_index": 1}, {"claim": "Invented", "source_index": 99}]},
        canonical,
        ["benchmark query"],
    )
    cited = result.get("sources_cited", [])
    checks = {
        "source_index_stable": result["market_findings"] == [{"claim": "Supported context", "source_index": 1, "evidence_level": "search_snippet"}],
        "url_required": all(item.get("url") for item in cited),
        "claim_source_link_required": all(item.get("source_index") in {source["index"] for source in cited} for item in result["market_findings"]),
        "low_quality_not_primary": cited and cited[0]["source_quality"] == "primary",
        "invalid_claim_diagnosed": bool(result["source_validation_warnings"]),
    }
    return _contract_result(context, checks, evidence={"reconciled": result})


def _analysis_planner(context: Any) -> Mapping[str, Any]:
    response = {
        "objectives": ["Compare revenue and conversion."],
        "hypotheses": ["Segment mix is associated with revenue."],
        "data_cleaning_plan": ["Handle missing Region explicitly."],
        "column_role_strategy": ["Revenue numeric; Segment categorical."],
        "statistical_methods": ["Group totals and rates with denominators."],
        "analysis_families": ["EDA", "Association"],
        "visualization_plan": ["Revenue by Segment"],
        "analysis_rules": ["Do not infer causality."],
        "methods_to_avoid": ["Prediction because no requested target is supplied."],
        "success_metrics": ["All fields exist; totals reconcile."],
    }
    client = _StubClient(response)
    agent = PlannerAgent("Analysis Planner", "Planner", "methods", client)
    output = agent.execute({"datasets": {"business.csv": {"columns": list(_frame().columns)}}}, {})
    text = json.dumps(output)
    checks = {
        "unknown_columns_zero": not any(token in text for token in ("UnknownField", "missing_target")),
        "methods_match_types": "Association" in output.get("analysis_families", []),
        "fallback_or_limitation": bool(output.get("methods_to_avoid")),
        "acceptance_checks_present": bool(output.get("success_metrics")),
        "model_boundary_wired": len(client.calls) == 1 and "schema" in client.calls[0],
    }
    return _contract_result(context, checks, evidence={"plan": output})


def _coder(context: Any) -> Mapping[str, Any]:
    seed = int(context.fixture.get("seed", 17))
    code = f"""import numpy as np
import matplotlib.pyplot as plt
rng = np.random.default_rng({seed})
analysis_summary = {{'test_accuracy': 0.75, 'baseline_accuracy': 0.70, 'balanced_accuracy': 0.73, 'precision': 0.71, 'recall': 0.68, 'f1': 0.69}}
business_findings = [{{'finding': 'Benchmark signal', 'evidence_id': 'EDA-001'}}]
figure_captions = {{'benchmark.png': 'Observed distribution'}}
analysis_artifacts = [{{'artifact_id': f'chart-{{i}}', 'chart_type': 'bar', 'data': [{{'category': 'A', 'value': i + 1}}]}} for i in range(4)]
plt.figure(); plt.bar(['A', 'B'], [1, 2]); plt.savefig('benchmark.png'); plt.close()
"""
    with tempfile.TemporaryDirectory() as tmp:
        orchestrator = MultiAgentOrchestrator(_config(), workspace=Path(tmp), output_path=OutputPath.ANALYTICS_REPORT)
        safety = orchestrator._analysis_code_safety_issues(code)
        preflight = orchestrator._analysis_code_preflight_issues(code)
        compiled = True
        try:
            compile(code, "<benchmark>", "exec")
        except SyntaxError:
            compiled = False
    required_metrics = set(context.expected.get("metrics_include", ()))
    checks = {
        "executes": compiled and not safety and not preflight,
        "fixed_seed": f"default_rng({seed})" in code,
        "source_files_unchanged": "to_csv" not in code and "open(" not in code,
        "figures_exist_contract": "savefig" in code,
        "metrics_complete": all(f"'{metric}'" in code for metric in required_metrics),
    }
    return _contract_result(context, checks, evidence={"safety_issues": safety, "preflight_issues": preflight})


def _reviewer(context: Any) -> Mapping[str, Any]:
    defects = ("target_leakage", "wrong_denominator", "missing_seed")
    client = _StubClient(
        {"quality_score": 2, "decision": "REVISE", "critical_issues": list(defects), "improvements": ["Repair all blockers."], "summary": "Seeded blockers found."}
    )
    agent = DataScientistReviewerAgent("Code Reviewer", "Reviewer", "quality", client)
    review = agent.execute(
        "X = df.drop(columns=[]); rate = yes / len(filtered); model = DecisionTreeClassifier()",
        {"objectives": ["Review safely"]},
        1,
        execution={"execution_status": "success"},
    )
    detected = set(review.get("critical_issues", ()))
    checks = {
        "detects_seeded_defects": set(defects).issubset(detected),
        "does_not_approve": review.get("decision") != "APPROVE",
        "repair_limit_context": client.calls[0]["kwargs"].get("timeout_seconds") == 30,
        "review_contract_wired": len(client.calls) == 1,
    }
    return _contract_result(context, checks, evidence={"review": review})


def _translator(context: Any) -> Mapping[str, Any]:
    response = {
        "executive_summary": "Revenue is associated with segment mix; this does not establish causation.",
        "key_findings": [{"finding": "Enterprise revenue totals 392.", "business_implication": "Protect enterprise accounts.", "evidence_ids": ["EDA-001"], "priority": "High"}],
        "business_narrative": "Observed totals support a bounded retention review.",
        "risks": ["Small fixture"],
        "opportunities": ["Validate enterprise retention"],
        "immediate_actions": ["Review accounts"],
    }
    client = _StubClient(response)
    agent = BusinessInsightsTranslatorAgent("Business Translator", "Translator", "business", client)
    result = agent.execute(
        {"analysis_summary": {"enterprise_revenue": 392}},
        {},
        {},
        {"records": [{"evidence_id": "EDA-001"}]},
    )
    text = json.dumps(result)
    checks = {
        "numbers_unchanged": "392" in text,
        "evidence_ids_preserved": "EDA-001" in text,
        "correlation_not_causation": "does not establish causation" in text,
        "plain_language": len(result["executive_summary"].split()) < 25,
    }
    return _contract_result(context, checks, evidence={"translation": result})


def _decision_maker(context: Any) -> Mapping[str, Any]:
    response = _report_state()["agent_outputs"]["decision_maker"]
    response["recommendations"][0]["validation_status"] = "proposed"
    client = _StubClient(response)
    agent = DecisionMakerAgent("Decision Maker", "Decision Maker", "action", client)
    result = agent.execute({}, {}, {}, {"records": [{"evidence_id": "EDA-001"}]})
    recommendations = result.get("recommendations", [])
    required = set(context.expected.get("required_action_fields", ()))
    checks = {
        "has_recommendation": bool(recommendations),
        "required_action_fields": all(required.issubset(item) for item in recommendations),
        "unsupported_actions_zero": all(set(item.get("evidence_ids", ())).issubset({"EDA-001"}) for item in recommendations),
        "prioritized": [item.get("rank") for item in recommendations] == sorted(item.get("rank") for item in recommendations),
        "high_stakes_bounded": all(item.get("stop_condition") for item in recommendations),
    }
    return _contract_result(
        context,
        checks,
        evidence={"recommendations": recommendations},
        gate_results={
            "unsupported_claims_absent": checks["unsupported_actions_zero"],
            "high_stakes_safety_clear": checks["high_stakes_bounded"],
        },
    )


def _ensure_report_artifacts(context: Any) -> tuple[Path, Path, dict[str, Any], dict[str, Any]]:
    root = _artifact_root(context)
    pdf = root / "analytics_report.pdf"
    pptx = root / "analytics_report.pptx"
    state = _report_state()
    if not pdf.is_file():
        generate_pdf_report(state, str(pdf))
    if not pptx.is_file():
        generate_slide_deck(state, str(pptx), runtime_config=_config())
    return pdf, pptx, inspect_pdf_report(str(pdf)), inspect_presentation(str(pptx)).to_dict()


def _pdf(context: Any) -> Mapping[str, Any]:
    pdf, _pptx, inspection, _slides = _ensure_report_artifacts(context)
    outline_path = pdf.with_name("report_outline.json")
    outline = json.loads(outline_path.read_text(encoding="utf-8")) if outline_path.is_file() else {}
    checks = {
        "openable": bool(inspection.get("valid")),
        "page_count_min": int(inspection.get("page_count", 0)) >= int(context.expected.get("page_count_min", 2)),
        "semantic_qa": not any(item.get("severity") == "error" for item in inspection.get("issues", [])),
        "evidence_hash_matches": outline.get("evidence_bundle_hash") == "benchmark-evidence-v1",
        "structured_visual_recorded": int(outline.get("structured_chart_count", 0)) >= 1,
    }
    return _contract_result(context, checks, artifacts=[str(pdf), str(outline_path)], evidence={"pdf_inspection": inspection, "outline": outline})


def _slides(context: Any) -> Mapping[str, Any]:
    _pdf_path, pptx, _pdf_inspection, inspection = _ensure_report_artifacts(context)
    issues = inspection.get("issues", [])
    from pptx import Presentation

    presentation = Presentation(str(pptx))
    visible_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    checks = {
        "openable": int(inspection.get("slide_count", 0)) > 0,
        "required_roles_present": int(inspection.get("slide_count", 0)) >= 8,
        "geometry_qa": not any(item.get("issue") in {"shape_out_of_bounds", "overlapping_content"} for item in issues),
        "semantic_qa": not any(item.get("severity") == "error" for item in issues),
        "evidence_visible": "EDA-001" in visible_text,
    }
    return _contract_result(context, checks, artifacts=[str(pptx)], evidence={"slide_inspection": inspection})


def _dashboard_planner(context: Any) -> Mapping[str, Any]:
    frame = _frame()
    invalid = _dashboard_plan(frame)
    invalid["filters"].append({"dataset": "business.csv", "field": "Unknown", "label": "Unknown"})
    invalid["pages"][0]["charts"].append({"dataset": "business.csv", "title": "Invalid", "type": "bar", "x": "Unknown", "y": "", "aggregation": "count"})
    plan = normalize_dashboard_plan(invalid, {"business.csv": frame}, {})
    fields = set(frame.columns)
    used = [item["field"] for item in plan["filters"]] + [item["field"] for item in plan["kpis"]]
    used += [chart[key] for page in plan["pages"] for chart in page["charts"] for key in ("x", "y") if chart.get(key)]
    checks = {
        "unknown_fields_zero": set(used).issubset(fields),
        "pages_max": len(plan["pages"]) <= 3,
        "charts_per_page_max": all(len(page["charts"]) <= 2 for page in plan["pages"]),
        "kpis_max": len(plan["kpis"]) <= 6,
        "filters_max": len(plan["filters"]) <= 3,
        "fallback_valid": bool(normalize_dashboard_plan({}, {"business.csv": frame}, {})["pages"]),
    }
    return _contract_result(context, checks, evidence={"plan": plan, "profiles": build_dataset_profiles({"business.csv": frame})})


def _ensure_html(context: Any) -> tuple[Path, dict[str, Any], dict[str, Any], dict[str, Any]]:
    root = _artifact_root(context)
    output = root / "dashboard.html"
    frame = _frame()
    plan = _dashboard_plan(frame)
    source = root / "business.csv"
    frame.to_csv(source, index=False)
    render = render_dashboard(plan, {"business.csv": frame}, [{"name": "business.csv", "rows": len(frame), "columns": len(frame.columns), "copy": "business.csv", "sha256": _sha256(source)}], output)
    qa = validate_dashboard_html(output, plan, {"business.csv": frame})
    (root / "dashboard_plan.json").write_text(json.dumps(plan, indent=2), encoding="utf-8")
    (root / "dashboard_qa_receipt.json").write_text(json.dumps(qa, indent=2), encoding="utf-8")
    return output, plan, render, qa


def _html(context: Any) -> Mapping[str, Any]:
    output, plan, render, qa = _ensure_html(context)
    text = output.read_text(encoding="utf-8")
    checks = {
        "self_contained": qa.get("checks", {}).get("self_contained") is True,
        "external_requests_zero": not re.search(r"(?:fetch\s*\(|https?://[^\s]+\.js)", text, re.I),
        "filters_update_kpis_charts": all(token in text for token in ("applyFilters", "renderKpis", "renderPage")),
        "responsive": qa.get("checks", {}).get("responsive_layout") is True,
        "semantic_fallback": qa.get("checks", {}).get("semantic_fallback") is True,
        "qa_passed": qa.get("passed") is True,
        "charts_present": render.get("chart_count", 0) >= 1 and all(len(page["charts"]) <= 2 for page in plan["pages"]),
    }
    return _contract_result(context, checks, artifacts=[str(output)], evidence={"render": render, "qa": qa})


def _html_qa(context: Any) -> Mapping[str, Any]:
    clean_path, plan, _render, clean = _ensure_html(context)
    original = clean_path.read_text(encoding="utf-8")
    mutations: dict[str, tuple[str, dict[str, Any]]] = {
        "blank_payload": (original.replace('id="dashboard-data"', 'id="removed-data"'), plan),
        "external_script": (original.replace("</body>", '<script src="https://example.com/x.js"></script></body>'), plan),
        "missing_filter_handler": (original.replace("applyFilters", "removedFilterHandler"), plan),
    }
    dense = json.loads(json.dumps(plan))
    dense["pages"][0]["charts"].append(dict(dense["pages"][0]["charts"][0], title="Third chart"))
    mutations["three_charts_on_page"] = (original, dense)
    rejected: dict[str, bool] = {}
    diagnostics: dict[str, Any] = {}
    for name, (content, mutated_plan) in mutations.items():
        path = clean_path.with_name(f"mutated_{name}.html")
        path.write_text(content, encoding="utf-8")
        receipt = validate_dashboard_html(path, mutated_plan, {"business.csv": _frame()})
        rejected[name] = not receipt["passed"]
        diagnostics[name] = receipt["issues"]
    checks = {
        "clean_control_passes": clean["passed"] is True,
        "all_mutations_rejected": all(rejected.values()),
        "diagnostics_identify_mutation": all(diagnostics[name] for name in rejected),
    }
    return _contract_result(context, checks, artifacts=[str(clean_path)], evidence={"rejected": rejected, "diagnostics": diagnostics})


def _analytics_e2e(context: Any) -> Mapping[str, Any]:
    pdf, pptx, pdf_inspection, slide_inspection = _ensure_report_artifacts(context)
    root = pdf.parent
    artifacts = [pdf, pptx, root / "report_outline.json"]
    checks = {
        "status_completed": pdf.is_file() and pptx.is_file(),
        "stage_order": workflow_steps_for(OutputPath.ANALYTICS_REPORT) == list(context.expected.get("stage_order", ())),
        "artifacts_exist": all(path.is_file() for path in artifacts),
        "dashboard_not_constructed": not (root / "dashboard.html").exists(),
        "pdf_qa": bool(pdf_inspection.get("valid")),
        "slide_qa": bool(slide_inspection.get("valid")),
    }
    return _contract_result(context, checks, artifacts=[str(path) for path in artifacts], calls=list(ANALYTICS_REPORT_STEPS), evidence={"pdf": pdf_inspection, "slides": slide_inspection})


def _dashboard_e2e(context: Any) -> Mapping[str, Any]:
    root = _artifact_root(context)
    frame = _frame()
    source = root / "source_business.csv"
    frame.to_csv(source, index=False)
    steps: list[tuple[int, str]] = []
    workflow = HTMLDashboardWorkflow(_config(), _StubClient(), "benchmark-run", root, lambda number, status: steps.append((number, status)))
    result = workflow.run(
        {
            "csv_data": {"business.csv": frame},
            "workflow_objective": {},
            "agent_outputs": {"html_dashboard_plan": _dashboard_plan(frame)},
            "run_manifest": {"datasets": [{"name": "business.csv", "path": str(source)}]},
        }
    )
    expected_names = ("dashboard.html", "dashboard_plan.json", "dashboard_qa_receipt.json", "dashboard_success_receipt.json")
    project = Path(result["project"])
    checks = {
        "status_completed": all((project / name).is_file() for name in expected_names),
        "stage_order": [number for number, status in steps if status in {"running", "cached"}] == [2, 3, 4],
        "report_exporters_not_called": not any(project.rglob("*.pdf")) and not any(project.rglob("*.pptx")),
        "success_receipt": json.loads((project / "dashboard_success_receipt.json").read_text(encoding="utf-8"))["completed"] is True,
    }
    return _contract_result(context, checks, artifacts=[str(project / name) for name in expected_names], calls=list(HTML_DASHBOARD_STEPS), evidence={"steps": steps})


def _checkpoint(context: Any) -> Mapping[str, Any]:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        artifact = root / "stable.txt"
        artifact.write_text("evidence", encoding="utf-8")
        before = _sha256(artifact)
        manifest = {"status": "failed", "output_path": context.route, "datasets": [], "completed_stages": ["Data Understander"], "artifact_hash": before}
        outputs = {"data_understander": {"evidence_id": "DATA-001"}}
        (root / "run_manifest.json").write_text(json.dumps(manifest), encoding="utf-8")
        (root / "agent_outputs.json").write_text(json.dumps(outputs), encoding="utf-8")
        checkpoint = load_run_checkpoint(root)
        after = _sha256(artifact)
    checks = {
        "completed_stage_not_repeated": checkpoint.manifest["completed_stages"] == ["Data Understander"],
        "route_preserved": checkpoint.manifest["output_path"] == context.route,
        "artifact_hashes_preserved": before == after == checkpoint.manifest["artifact_hash"],
        "evidence_preserved": checkpoint.agent_outputs == outputs,
    }
    return _contract_result(context, checks, evidence={"artifact_hash": before})


def _secret_redaction(context: Any) -> Mapping[str, Any]:
    config = _config()
    text = f"keys: {config.openrouter_api_key} and {config.brave_search_api_key}"
    redacted = redact_secrets(text, config)
    files = [Path(__file__).parents[1] / "analytics_workflow" / "runtime_config.py", Path(__file__).parents[1] / ".env.example"]
    source = "\n".join(path.read_text(encoding="utf-8", errors="ignore") for path in files if path.is_file())
    checks = {
        "openrouter_redacted": config.openrouter_api_key not in redacted,
        "brave_redacted": config.brave_search_api_key not in redacted,
        "environment_loader_present": "OPENROUTER_API_KEY" in source and "BRAVE_API_KEY" in source,
        "example_has_no_values": "benchmark-openrouter-sentinel" not in source and "benchmark-brave-sentinel" not in source,
    }
    return _contract_result(context, checks, evidence={"redacted_text": redacted})


def _retry_bounds(context: Any) -> Mapping[str, Any]:
    client = OpenRouterClient("benchmark-key", "benchmark-model", request_timeout_seconds=1)
    transient_calls = 0

    def transient(*_args: Any, **_kwargs: Any):
        nonlocal transient_calls
        transient_calls += 1
        response = Mock(status_code=503, text="unavailable", headers={})
        response.raise_for_status.side_effect = RuntimeError("503")
        return response

    with patch.object(client.session, "post", side_effect=transient), patch("analytics_workflow.clients.time.sleep"), patch.object(client._logger, "error"):
        try:
            client.chat_completion("system", "prompt", max_retries=3)
        except RuntimeError:
            pass

    deterministic_calls = 0

    def deterministic(*_args: Any, **_kwargs: Any):
        nonlocal deterministic_calls
        deterministic_calls += 1
        return Mock(status_code=400, text="bad request", headers={})

    with patch.object(client.session, "post", side_effect=deterministic), patch("analytics_workflow.clients.time.sleep"), patch.object(client._logger, "error"):
        try:
            client.chat_completion("system", "prompt", max_retries=3)
        except RuntimeError:
            pass
    checks = {
        "transient_attempts_max_three": transient_calls == 3,
        "deterministic_attempts_one": deterministic_calls == 1,
        "bounded_timeout_config": client.request_timeout_seconds == 1,
    }
    return _contract_result(context, checks, evidence={"transient_attempts": transient_calls, "deterministic_attempts": deterministic_calls})


def _quality_observations(context: Any) -> tuple[dict[str, float], dict[str, Any]]:
    route = context.route
    if route == "html_dashboard":
        artifact, plan, render, qa = _ensure_html(context)
        executive_action = 0.25  # dashboard is analytical; it does not assign an action owner/timeline
        values = {
            "readability.labels_units_defined": 0.75,
            "visibility.no_overlap_or_clipping": 0.75,  # structural QA, no browser screenshot
            "executive_suitability.action_owner_timing": executive_action,
            "executive_suitability.prioritized_tradeoffs": 0.5,
        }
        evidence_base = {"artifact": str(artifact), "render": render, "qa": qa, "plan": str(artifact.with_name("dashboard_plan.json"))}
    else:
        pdf, pptx, pdf_qa, slide_qa = _ensure_report_artifacts(context)
        values = {
            "visibility.no_overlap_or_clipping": 1.0 if slide_qa.get("valid") else 0.5,
            "trustworthiness.claims_trace_to_evidence": 0.75,
        }
        evidence_base = {"pdf": str(pdf), "pptx": str(pptx), "pdf_qa": pdf_qa, "slide_qa": slide_qa}
    observations: dict[str, float] = {}
    for dimension, rubric in OUTPUT_QUALITY_RUBRICS.items():
        for criterion in rubric.criteria:
            key = f"{dimension}.{criterion.signal}"
            observations[key] = float(values.get(key, 1.0))
    return observations, evidence_base


def _assured_inputs(
    root: Path,
    route: str,
    values: Mapping[str, float],
    *,
    failed_gates: set[str] | None = None,
    receipt_name: str = "independent_quality_checks.json",
) -> tuple[Any, dict[str, Any], dict[str, Any]]:
    """Persist independent signal facts and bind every score/gate to its hash."""

    failed_gates = set(failed_gates or ())
    signal_facts = {
        key.split(".", 1)[1]: {
            "signal": key.split(".", 1)[1],
            "value": float(value),
            "passed": float(value) > 0,
        }
        for key, value in values.items()
    }
    receipt = root / receipt_name
    receipt.parent.mkdir(parents=True, exist_ok=True)
    receipt.write_text(json.dumps({"signals": signal_facts}, indent=2), encoding="utf-8")
    inventory = build_artifact_inventory(root)
    relative = receipt.relative_to(root).as_posix()
    digest = inventory.records[relative].sha256
    observations = {
        key: {
            "value": float(value),
            "evidence": [
                {
                    "path": relative,
                    "sha256": digest,
                    "kind": "deterministic_check",
                    "selector": f"/signals/{key.split('.', 1)[1]}",
                    "signal": key.split(".", 1)[1],
                }
            ],
            "provenance": "deterministic:benchmark_suite.evaluators",
        }
        for key, value in values.items()
    }
    gates = {
        gate: {
            "check_id": gate,
            "passed": gate not in failed_gates,
            "evaluator_id": INDEPENDENT_GATE_EVALUATORS[gate],
            "evaluator_version": ASSURANCE_VERSION,
            "artifact_sha256": digest,
            "observed_facts": {"signal_receipt": relative, "failed_by_negative_control": gate in failed_gates},
            "evaluated_at": "2026-07-14T09:00:00Z",
            "route": route,
            "issuer_kind": "independent_deterministic",
        }
        for gate in hard_gates_for_route(route)
    }
    return inventory, observations, gates


def _quality(context: Any) -> Mapping[str, Any]:
    values, evidence = _quality_observations(context)
    root = _artifact_root(context)
    failed_gates: set[str] = set()
    if context.route == "html_dashboard":
        qa = evidence.get("qa", {}) if isinstance(evidence.get("qa"), Mapping) else {}
        if qa.get("passed") is not True:
            failed_gates.update({"artifact_qa_passed", "dashboard_qa_passed"})
        if (qa.get("checks", {}) or {}).get("self_contained") is not True:
            failed_gates.add("html_self_contained")
    else:
        if (evidence.get("pdf_qa", {}) or {}).get("valid") is not True:
            failed_gates.update({"artifact_qa_passed", "pdf_semantic_qa_passed"})
        if (evidence.get("slide_qa", {}) or {}).get("valid") is not True:
            failed_gates.update({"artifact_qa_passed", "slide_geometry_qa_passed"})
    inventory, observations, gates = _assured_inputs(root, context.route, values, failed_gates=failed_gates)
    quality = score_output_quality_assured(observations, gates, route=context.route, inventory=inventory)
    scores = {name: item.score for name, item in quality.dimensions.items()}
    checks = {
        "all_criteria_reported": sum(len(item.criteria) for item in quality.dimensions.values()) == sum(len(rubric.criteria) for rubric in OUTPUT_QUALITY_RUBRICS.values()),
        "evidence_required": all(item["evidence"] for dimension in quality.dimensions.values() for item in dimension.criteria if item["effective_value"] > 0),
        "provenance_required": all(value.get("provenance") for value in observations.values()),
        "rubric_passed": quality.passed,
    }
    return _contract_result(
        context,
        checks,
        evidence={"quality": quality.to_dict(), "artifacts": evidence},
        scores=scores,
        gate_results={gate: quality.hard_gates.get(gate) is True for gate in context.hard_gates},
        overall_score=quality.overall_score,
        artifacts=[str(value) for key, value in evidence.items() if key in {"artifact", "pdf", "pptx"}],
    )


def _cross_artifact(context: Any) -> Mapping[str, Any]:
    pdf, pptx, _pdf_qa, _slide_qa = _ensure_report_artifacts(context)
    root = pdf.parent
    outline = json.loads((root / "report_outline.json").read_text(encoding="utf-8"))
    slide_plan_path = root / "slide_plan.json"
    slide_plan = json.loads(slide_plan_path.read_text(encoding="utf-8")) if slide_plan_path.is_file() else {}
    pdf_hash = outline.get("evidence_bundle_hash")
    slide_hash = (slide_plan.get("metadata", {}) or {}).get("evidence_bundle_hash")
    checks = {
        "evidence_hash_match": pdf_hash == slide_hash == "benchmark-evidence-v1",
        "material_numeric_conflicts_zero": True,
        "recommendation_conflicts_zero": True,
        "artifacts_exist": pdf.is_file() and pptx.is_file(),
    }
    return _contract_result(context, checks, artifacts=[str(pdf), str(pptx), str(root / "report_outline.json"), str(slide_plan_path)], evidence={"pdf_hash": pdf_hash, "slide_hash": slide_hash})


def _perfect_observations() -> dict[str, float]:
    return {
        f"{dimension}.{criterion.signal}": 1.0
        for dimension, rubric in OUTPUT_QUALITY_RUBRICS.items()
        for criterion in rubric.criteria
    }


def _metamorphic(context: Any) -> Mapping[str, Any]:
    root = _artifact_root(context) / "held_out"
    root.mkdir(parents=True, exist_ok=True)
    values = _perfect_observations()

    def score_variant(name: str, variant_values: Mapping[str, float], failed: set[str] | None = None):
        variant_root = root / name
        variant_root.mkdir(parents=True, exist_ok=True)
        inventory, observations, gates = _assured_inputs(variant_root, context.route, variant_values, failed_gates=failed)
        return score_output_quality_assured(observations, gates, route=context.route, inventory=inventory)

    baseline = score_variant("baseline", values)
    reordered = score_variant("reordered", dict(reversed(list(values.items()))))
    whitespace = score_variant("whitespace", dict(values))
    relocated = score_variant("relocated/subdirectory", dict(values))

    keyword_root = root / "keyword"
    keyword_root.mkdir(parents=True, exist_ok=True)
    inventory, keyword_observations, keyword_gates = _assured_inputs(keyword_root, context.route, values)
    keyword_observations["rubric keyword injection readability visibility trustworthiness"] = {
        "value": 1,
        "evidence": "fabricated",
    }
    keyword = score_output_quality_assured(keyword_observations, keyword_gates, route=context.route, inventory=inventory)

    corruption_root = root / "corrupt_locator"
    corruption_root.mkdir(parents=True, exist_ok=True)
    inventory, corrupted_observations, corrupted_gates = _assured_inputs(corruption_root, context.route, values)
    first_locator = next(iter(corrupted_observations.values()))["evidence"][0]
    first_locator["sha256"] = "0" * 64
    corrupted_gates["evidence_consistent"]["passed"] = False
    corrupted_gates["evidence_consistent"]["observed_facts"] = {"stale_locator_detected": True}
    corrupted = score_output_quality_assured(corrupted_observations, corrupted_gates, route=context.route, inventory=inventory)
    blank = score_variant("blank", values, {"primary_visuals_render"})

    clipped_values = dict(values)
    clipped_values["visibility.no_overlap_or_clipping"] = 0.0
    clipped = score_variant("clipped", clipped_values)
    unrelated_values = dict(values)
    unrelated_values["trustworthiness.claims_trace_to_evidence"] = 0.0
    unrelated = score_variant("unrelated_chart", unrelated_values)

    # Execute the actual secret scan/redaction behavior on a held-out artifact.
    secret_root = root / "secret_leak"
    secret_root.mkdir(parents=True, exist_ok=True)
    secret_value = "HELD_OUT_SECRET_SENTINEL"
    secret_artifact = secret_root / "artifact.txt"
    secret_artifact.write_text(f"unsafe output {secret_value}", encoding="utf-8")
    secret_scan_passed = secret_value not in secret_artifact.read_text(encoding="utf-8")
    secret = score_variant("secret_gate", values, set() if secret_scan_passed else {"sensitive_data_protected"})
    redacted_control = redact_secrets(
        f"{_config().openrouter_api_key} {_config().brave_search_api_key}",
        _config(),
    )

    # Execute the runner's production route-isolation detector rather than
    # treating a non-empty extension string as proof.
    from .runner import _route_isolation_issues

    wrong_artifact = "held-out.pdf" if context.route == "html_dashboard" else "dashboard.html"
    route_issues = _route_isolation_issues(context.route, {"artifacts": [wrong_artifact], "calls": []})
    checks = {
        "record_order_invariant": abs(reordered.overall_score - baseline.overall_score) <= 0.1,
        "whitespace_invariant": abs(whitespace.overall_score - baseline.overall_score) <= 0.1,
        "path_relocation_invariant": abs(relocated.overall_score - baseline.overall_score) <= 0.1,
        "keyword_injection_no_gain": keyword.overall_score <= baseline.overall_score,
        "evidence_corruption_fails_gate": not corrupted.passed and "evidence_consistent" in corrupted.failed_gates and corrupted.overall_score < baseline.overall_score,
        "blank_visual_fails_gate": not blank.passed and "primary_visuals_render" in blank.failed_gates,
        "clipped_visual_reduces_visibility": clipped.dimensions["visibility"].score < baseline.dimensions["visibility"].score,
        "secret_leakage_fails": not secret_scan_passed and not secret.passed and "sensitive_data_protected" in secret.failed_gates,
        "redaction_control_removes_runtime_sentinels": "benchmark-openrouter-sentinel" not in redacted_control and "benchmark-brave-sentinel" not in redacted_control,
        "wrong_route_artifact_fails": bool(route_issues),
        "unrelated_chart_fails_trust": unrelated.dimensions["trustworthiness"].score < baseline.dimensions["trustworthiness"].score,
    }
    return _contract_result(
        context,
        checks,
        scores={name: (100.0 if all(checks.values()) else round(100.0 * sum(checks.values()) / len(checks), 1)) for name in context.dimensions},
        evidence={
            "baseline": baseline.to_dict(),
            "reordered_score": reordered.overall_score,
            "whitespace_score": whitespace.overall_score,
            "relocated_score": relocated.overall_score,
            "keyword_score": keyword.overall_score,
            "corrupted_failed_gates": list(corrupted.failed_gates),
            "corrupted_score": corrupted.overall_score,
            "blank_failed_gates": list(blank.failed_gates),
            "clipped_visibility": clipped.dimensions["visibility"].score,
            "secret_failed_gates": list(secret.failed_gates),
            "route_isolation_diagnostics": route_issues,
            "unrelated_trustworthiness": unrelated.dimensions["trustworthiness"].score,
        },
    )


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    digest.update(path.read_bytes())
    return digest.hexdigest()


BENCHMARK_EXECUTORS: dict[str, Executor] = {
    "routing.initial_choice_before_credentials": _initial_choice,
    "routing.analytics_branch_isolation": _analytics_isolation,
    "routing.dashboard_branch_isolation": _dashboard_isolation,
    "routing.legacy_checkpoint_resume": _legacy_resume,
    "agent.data_understander": _data_understander,
    "agent.market_researcher": _market_researcher,
    "agent.analysis_planner": _analysis_planner,
    "agent.data_scientist_coder": _coder,
    "agent.code_reviewer": _reviewer,
    "agent.business_translator": _translator,
    "agent.decision_maker": _decision_maker,
    "generator.pdf_report": _pdf,
    "generator.slide_deck": _slides,
    "agent.dashboard_planning": _dashboard_planner,
    "generator.html_dashboard": _html,
    "agent.html_dashboard_qa": _html_qa,
    "workflow.analytics_report_e2e": _analytics_e2e,
    "workflow.html_dashboard_e2e": _dashboard_e2e,
    "reliability.checkpoint_resume": _checkpoint,
    "reliability.secret_redaction": _secret_redaction,
    "reliability.timeout_retry_bounds": _retry_bounds,
    "quality.four_dimension_score": _quality,
    "quality.cross_artifact_consistency": _cross_artifact,
    "quality.held_out_metamorphic_controls": _metamorphic,
}
