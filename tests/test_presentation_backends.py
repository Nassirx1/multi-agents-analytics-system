from __future__ import annotations

import asyncio
import tempfile
import time
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches

from analytics_workflow.presentation_backends import (
    HYBRID_SLIDE_OBJECTIVES,
    enrich_deck_executive_copy,
    PowerPointMCPBackend,
    PythonPresentationBackend,
    SAFE_AGENT9_TOOLS,
    backend_from_config,
    inspect_presentation,
)
from analytics_workflow.presentation_backends import _shape_name_from_result
from analytics_workflow.presentation_backends import _apply_hybrid_story_contract
from analytics_workflow.deck_rendering import (
    _materialize_deck_in_run_directory,
    build_consulting_deck,
)
from analytics_workflow.runtime_config import build_runtime_config, register_runtime_config
from analytics_workflow.slides.deck_spec import ContentBlock, DeckSpec, SlideSpec, VisualSpec


class PresentationBackendTests(unittest.TestCase):
    def tearDown(self) -> None:
        register_runtime_config(None)

    def test_backend_deck_is_materialized_inside_requested_run_directory(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "backend-output" / "deck.pptx"
            source.parent.mkdir()
            source.write_bytes(b"new presentation")
            requested = root / "runs" / "run-123" / "analytics_report.pptx"

            result = _materialize_deck_in_run_directory(str(source), str(requested))

            self.assertEqual(Path(result), requested.resolve())
            self.assertEqual(requested.read_bytes(), b"new presentation")

    def test_backend_specific_name_is_copied_to_canonical_run_deck(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            run_dir = Path(temp_dir) / "runs" / "run-123"
            run_dir.mkdir(parents=True)
            source = run_dir / "backend-specific-name.pptx"
            source.write_bytes(b"presentation")

            result = _materialize_deck_in_run_directory(
                str(source), str(run_dir / "analytics_report.pptx")
            )

            self.assertEqual(Path(result), (run_dir / "analytics_report.pptx").resolve())
            self.assertEqual(Path(result).read_bytes(), b"presentation")

    def test_agent9_tool_allowlist_excludes_unsafe_capabilities(self) -> None:
        self.assertIn("slide", SAFE_AGENT9_TOOLS)
        self.assertIn("chart", SAFE_AGENT9_TOOLS)
        self.assertNotIn("vba", SAFE_AGENT9_TOOLS)
        self.assertNotIn("slideshow", SAFE_AGENT9_TOOLS)
        self.assertEqual(len(HYBRID_SLIDE_OBJECTIVES), 12)
        self.assertIn("decision-tree", HYBRID_SLIDE_OBJECTIVES[6])

    def test_shape_name_parser_reads_mcp_message_instead_of_json_key(self) -> None:
        result = '{"success":true,"message":"Created shape \'Rectangle 2\' successfully."}'
        self.assertEqual(_shape_name_from_result(result, "fallback"), "Rectangle 2")

    def test_hybrid_story_contract_uses_explicit_post_eda_sequence(self) -> None:
        roles = [
            "title", "data_understanding", "market_context", "analysis", "analysis", "analysis", "analysis",
            "findings", "business_translation", "recommendations", "limitations", "summary",
        ]
        deck = DeckSpec(
            deck_title="Contract",
            slides=[
                SlideSpec(index, role, "three_finding_cards", f"Slide {index}", main_message="Evidence")
                for index, role in enumerate(roles, start=1)
            ],
        )
        _apply_hybrid_story_contract(deck)
        self.assertEqual(
            [slide.slide_role for slide in deck.slides],
            [
                "title", "data_understanding", "market_context", "analysis", "analysis", "analysis", "analysis",
                "recommendations", "recommendations", "limitations", "sources", "ending",
            ],
        )

    def test_analysis_copy_uses_business_implication_instead_of_duplicate_finding(self) -> None:
        finding = "Sales department has a 20.6% attrition rate"
        deck = DeckSpec(
            deck_title="Executive",
            slides=[SlideSpec(
                1, "analysis", "chart_left_insight_right", finding,
                main_message=finding,
                content_blocks=[ContentBlock(type="bullets", items=[finding])],
                visual=VisualSpec(type="code_figure", finding=finding, takeaway=finding),
            )],
        )
        enrich_deck_executive_copy(
            deck,
            {"business_translator": {
                "key_findings": [{
                    "finding": "Sales department has the highest attrition rate at 20.6%",
                    "business_implication": "Sales turnover threatens revenue continuity and customer relationships.",
                }],
                "immediate_actions": ["Sales leaders should run targeted stay interviews."],
            }},
        )
        items = deck.slides[0].content_blocks[0].items
        self.assertTrue(any("revenue continuity" in item for item in items))
        self.assertNotIn(finding, items)

    def test_backend_selection_prefers_mcp_in_auto_mode(self) -> None:
        config = build_runtime_config(
            "openrouter",
            "brave",
            presentation_backend="auto",
            agent_request_timeout_seconds=181,
            presentation_agent_timeout_seconds=901,
        )
        backend = backend_from_config(config, object())
        self.assertIsInstance(backend, PowerPointMCPBackend)
        self.assertEqual(backend.request_timeout_seconds, 181)
        self.assertEqual(backend.timeout_seconds, 901)

    def test_backend_selection_allows_explicit_python_fallback(self) -> None:
        config = build_runtime_config("openrouter", "brave", presentation_backend="python")
        backend = backend_from_config(config, object())
        self.assertIsInstance(backend, PythonPresentationBackend)

    def test_mcp_path_guard_blocks_writes_outside_run_directory(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "run" / "deck.pptx"
            outside = Path(temp_dir) / "outside.pptx"
            with self.assertRaisesRegex(RuntimeError, "outside the run directory"):
                PowerPointMCPBackend._validate_tool_call(
                    "file", {"action": "create", "path": str(outside)}, output, set()
                )

    def test_mcp_path_guard_blocks_untrusted_images(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "run" / "deck.pptx"
            image = Path(temp_dir) / "other" / "image.png"
            with self.assertRaisesRegex(RuntimeError, "untrusted.*image"):
                PowerPointMCPBackend._validate_tool_call(
                    "image", {"image_path": str(image)}, output, set()
                )

    def test_inspection_opens_deck_and_detects_empty_slide(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "deck.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(path)
            inspection = inspect_presentation(str(path))
            self.assertEqual(inspection.slide_count, 1)
            self.assertFalse(inspection.valid)
            self.assertIn("empty_slide", {issue["issue"] for issue in inspection.issues})

    def test_inspection_accepts_non_placeholder_visual_title(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8)).text = "Executive Summary"
            slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8), Inches(2)).text = "Evidence-backed content"
            presentation.save(path)
            inspection = inspect_presentation(str(path))
            self.assertTrue(inspection.valid, inspection.issues)

    def test_expected_deck_qa_checks_count_headline_and_visual(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.8)).text = "Wrong headline"
            presentation.save(path)
            expected = DeckSpec(
                deck_title="Expected",
                slides=[
                    SlideSpec(
                        slide_number=1,
                        slide_role="analysis",
                        template="chart_left_insight_right",
                        headline="Expected analytical headline",
                        visual=VisualSpec(type="structured_chart", chart_type="bar"),
                    ),
                    SlideSpec(
                        slide_number=2,
                        slide_role="summary",
                        template="executive_summary_closing",
                        headline="Expected closing headline",
                    ),
                ],
            )
            inspection = inspect_presentation(str(path), expected_deck=expected)
            issues = {issue["issue"] for issue in inspection.issues}
            self.assertFalse(inspection.valid)
            self.assertIn("incorrect_slide_count", issues)
            self.assertIn("expected_headline_missing", issues)
            self.assertIn("missing_expected_visual", issues)

    def test_overflow_on_recommendation_slide_is_a_qa_error(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "overflow.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
            title.text = "Recommendation"
            body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(1.0), Inches(0.3))
            body.text = "This executive recommendation is intentionally much too long for this tiny text box."
            presentation.save(path)
            expected = DeckSpec(
                deck_title="Expected",
                slides=[SlideSpec(1, "recommendations", "recommendation_cards", "Recommendation")],
            )

            inspection = inspect_presentation(str(path), expected_deck=expected)

            self.assertFalse(inspection.valid)
            self.assertTrue(
                any(
                    issue.get("issue") == "possible_text_overflow" and issue.get("severity") == "error"
                    for issue in inspection.issues
                )
            )

    def test_agent9_shared_deadline_stops_before_another_model_call(self) -> None:
        backend = PowerPointMCPBackend(object(), timeout_seconds=180)
        with self.assertRaisesRegex(TimeoutError, "total deadline"):
            asyncio.run(
                backend._run_tool_round(
                    None,
                    [],
                    [],
                    {},
                    Path("deck.pptx"),
                    set(),
                    call_budget=1,
                    deadline=time.monotonic() - 1,
                )
            )

    def test_agent9_two_slide_copy_batches_are_applied_to_deck(self) -> None:
        class Client:
            def __init__(self):
                self.calls = []

            def chat_completion_json(self, system, prompt, schema, **kwargs):
                self.calls.append(prompt)
                numbers = [int(value) for value in __import__("re").findall(r'"slide_number":\s*(\d+)', prompt)]
                return {
                    "slides": [
                        {
                            "slide_number": number,
                            "headline": f"Executive answer {number}",
                            "main_message": f"Decision implication {number}",
                        }
                        for number in numbers
                    ]
                }

        client = Client()
        backend = PowerPointMCPBackend(client)
        deck = DeckSpec(
            deck_title="Test",
            slides=[
                SlideSpec(index, "analysis", "chart_left_insight_right", f"Old {index}")
                for index in range(1, 4)
            ],
        )

        backend._polish_slide_copy_batches(deck, time.monotonic() + 30)

        self.assertEqual(len(client.calls), 2)
        self.assertEqual(deck.slides[0].headline, "Executive answer 1")
        self.assertEqual(deck.slides[2].main_message, "Decision implication 3")
        self.assertEqual(deck.metadata["agent9_copy_polish"]["slides_edited"], 3)

    def test_scaffold_precreates_every_slide_with_content_before_model_work(self) -> None:
        class Result:
            isError = False
            content = []

            def __init__(self, payload):
                self.structuredContent = payload

        class Session:
            def __init__(self):
                self.calls = []

            async def call_tool(self, name, arguments, **kwargs):
                self.calls.append((name, dict(arguments)))
                if name == "file" and arguments.get("action") == "create":
                    return Result({"session_id": "scaffold-session"})
                if name == "shape" and arguments.get("action") in {"add-shape", "add-textbox"}:
                    return Result({"shapeName": f"Shape {len(self.calls)}"})
                return Result({"success": True})

        deck = DeckSpec(
            deck_title="Scaffold",
            slides=[
                SlideSpec(1, "title", "title_cover", "Opening headline", main_message="Opening message"),
                SlideSpec(2, "summary", "executive_summary", "Summary headline", main_message="Summary message"),
            ],
        )
        backend = PowerPointMCPBackend(object())
        session = Session()
        asyncio.run(
            backend._create_scaffold(
                session,
                deck,
                Path("scaffold.pptx").resolve(),
                time.monotonic() + 30,
            )
        )
        slide_creates = [args for name, args in session.calls if name == "slide" and args.get("action") == "create"]
        textboxes = [args for name, args in session.calls if name == "shape" and args.get("action") == "add-textbox"]
        self.assertEqual(len(slide_creates), 2)
        first_shape_index = next(index for index, (name, _) in enumerate(session.calls) if name == "shape")
        self.assertTrue(all(name == "slide" for name, _ in session.calls[1:3]))
        self.assertGreaterEqual(first_shape_index, 4)
        self.assertTrue(any(args.get("text") == "Opening headline" for args in textboxes))
        self.assertTrue(any("Summary message" in args.get("text", "") for args in textboxes))
        self.assertEqual(backend._active_session_id, "scaffold-session")

    def test_malformed_tool_arguments_are_returned_to_model_for_correction(self) -> None:
        class CorrectingClient:
            def __init__(self):
                self.calls = 0

            def chat_completion_with_tools(self, messages, tools, **kwargs):
                self.calls += 1
                if self.calls == 1:
                    return {
                        "content": "",
                        "tool_calls": [
                            {
                                "id": "bad-call",
                                "function": {"name": "shape", "arguments": '{"action":"add-shape", bad}'},
                            }
                        ],
                    }
                return {"content": "corrected", "tool_calls": []}

        client = CorrectingClient()
        backend = PowerPointMCPBackend(client)
        messages = []
        calls = asyncio.run(
            backend._run_tool_round(
                object(),
                messages,
                [],
                {"shape": object()},
                Path("deck.pptx"),
                set(),
                call_budget=3,
                deadline=time.monotonic() + 30,
            )
        )
        self.assertEqual(calls, 1)
        feedback = [message for message in messages if message.get("role") == "tool"]
        self.assertIn("invalid JSON", feedback[0]["content"])

    def test_model_cannot_replace_runtime_owned_mcp_session(self) -> None:
        class Client:
            def __init__(self):
                self.calls = 0

            def chat_completion_with_tools(self, messages, tools, **kwargs):
                self.calls += 1
                if self.calls == 1:
                    return {
                        "content": "",
                        "tool_calls": [{
                            "id": "session-call",
                            "function": {
                                "name": "shape",
                                "arguments": '{"action":"list","session_id":"invented","slide_index":1}',
                            },
                        }],
                    }
                return {"content": "done", "tool_calls": []}

        class Result:
            isError = False
            structuredContent = {"success": True}
            content = []

        class Session:
            def __init__(self):
                self.arguments = None

            async def call_tool(self, name, arguments, **kwargs):
                self.arguments = dict(arguments)
                return Result()

        backend = PowerPointMCPBackend(Client())
        backend._active_session_id = "runtime-session"
        session = Session()
        asyncio.run(
            backend._run_tool_round(
                session,
                [],
                [],
                {"shape": object()},
                Path("deck.pptx"),
                set(),
                call_budget=3,
                deadline=time.monotonic() + 30,
                allowed_slide_indexes={1},
            )
        )
        self.assertEqual(session.arguments["session_id"], "runtime-session")
        self.assertEqual(backend._active_session_id, "runtime-session")

    def test_optional_enhancement_stops_after_first_mcp_error(self) -> None:
        class Client:
            def chat_completion_with_tools(self, messages, tools, **kwargs):
                return {
                    "content": "",
                    "tool_calls": [{
                        "id": "bad-mcp-call",
                        "function": {
                            "name": "shape",
                            "arguments": '{"action":"list","slide_index":1}',
                        },
                    }],
                }

        class Result:
            isError = True
            structuredContent = {"success": False, "errorMessage": "COM error"}
            content = []

        class Session:
            def __init__(self):
                self.calls = 0

            async def call_tool(self, name, arguments, **kwargs):
                self.calls += 1
                return Result()

        backend = PowerPointMCPBackend(Client())
        backend._active_session_id = "runtime-session"
        session = Session()
        calls = asyncio.run(
            backend._run_tool_round(
                session, [], [], {"shape": object()}, Path("deck.pptx"), set(),
                call_budget=16, deadline=time.monotonic() + 30,
                allowed_slide_indexes={1}, strict_budget=False,
            )
        )
        self.assertEqual(calls, 1)
        self.assertEqual(session.calls, 1)

    def test_mcp_failure_automatically_uses_python_backend(self) -> None:
        class FailingMCP(PowerPointMCPBackend):
            def render(self, deck, output_path, **kwargs):
                raise RuntimeError("MCP unavailable for test")

        register_runtime_config(build_runtime_config("openrouter", "brave", presentation_backend="auto"))
        workflow_state = {
            "agent_outputs": {"presentation_architect": {"slides": []}},
            "analysis_results": {},
            "saved_figures": [],
            "generated_reports": {},
            "run_manifest": {"warnings": []},
        }
        with tempfile.TemporaryDirectory() as temp_dir, patch(
            "analytics_workflow.presentation_backends.backend_from_config",
            return_value=FailingMCP(object()),
        ):
            result = build_consulting_deck(workflow_state, str(Path(temp_dir) / "fallback.pptx"))
            self.assertTrue(Path(result).is_file())
        self.assertEqual(workflow_state["presentation_backend_used"], "python")
        self.assertEqual(
            [entry["status"] for entry in workflow_state["presentation_backend_log"]],
            ["failed", "success"],
        )


if __name__ == "__main__":
    unittest.main()
