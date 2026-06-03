import base64
import json
import os
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from analytics_workflow.deck_rendering import normalize_slide_plan
from analytics_workflow.pipeline_runtime import generate_slide_deck
from analytics_workflow.slides.deck_spec import DeckSpec, SlideSpec, VisualSpec
from analytics_workflow.slides.pptx_renderer import PowerPointRenderer
from analytics_workflow.slides.story_builder import build_deck_spec, _dedupe_text_items
from analytics_workflow.slides.theme import THEME
from analytics_workflow.slides.templates import REQUIRED_SLIDE_ROLES, TEMPLATE_REGISTRY
from analytics_workflow.slides.text_refiner import compact_whitespace, soften_unsupported_impact_claim


PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Zl8QAAAAASUVORK5CYII="
)


def reconstructed_chart_shape_count(presentation: Presentation) -> int:
    return sum(
        1
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "name", "").startswith("reconstructed_chart")
    )


def decision_tree_shape_count(presentation: Presentation) -> int:
    return sum(
        1
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "name", "").startswith("reconstructed_chart_decision_tree")
    )


class SlideDeckTests(unittest.TestCase):
    def test_text_refiner_normalizes_unit_symbols_for_export_text(self) -> None:
        clean = compact_whitespace("Temperature rose +0.08\u00b0C and solar yield reached 6.4 kWh/m\u00b2/day \u2014 validate.")

        self.assertIn("+0.08 deg C", clean)
        self.assertIn("m2/day", clean)
        self.assertIn(" - validate", clean)

    def test_saved_figures_become_visual_slides(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            figure_paths = []
            for index in range(1, 5):
                figure_path = temp_path / f"figure_{index}.png"
                figure_path.write_bytes(PNG_BYTES)
                figure_paths.append(str(figure_path))

            output_path = temp_path / "analytics_report.pptx"
            workflow_state = {
                "saved_figures": figure_paths,
                "analysis_results": {
                    "analysis_summary": {
                        "average_return_pct": 3.1,
                        "peak_volume": "2.7M shares",
                    },
                    "figure_captions": {
                        figure_paths[0]: "Sales trend improved in the latest period.",
                        figure_paths[1]: "Volume spikes align with price volatility.",
                        figure_paths[2]: "Returns cluster around a narrow band.",
                        figure_paths[3]: "Drawdown periods are limited in duration.",
                    }
                },
                "agent_outputs": {
                    "presentation_architect": {
                        "presentation_title": "STC Analytics Review",
                        "presentation_subtitle": "Executive summary",
                        "slides": [
                            {
                                "slide_number": 1,
                                "title": "Context",
                                "main_message": "The analysis highlights performance and risk patterns.",
                                "details": ["Dataset loaded and reviewed.", "Market and business implications summarized."],
                                "visual_element": "",
                            },
                            {
                                "slide_number": 2,
                                "title": "Recommendation",
                                "main_message": "Focus on evidence-backed actions.",
                                "details": ["Preserve the strongest signal from the analysis."],
                                "visual_element": "",
                            },
                        ],
                    },
                    "decision_maker": {},
                    "business_translator": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))

            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 12)

            picture_slides = 0
            all_text = []
            for slide in presentation.slides:
                if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
                    picture_slides += 1
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)

            full_text = "\n".join(all_text)
            self.assertEqual(picture_slides, 4)
            self.assertFalse(
                any(
                    warning.get("issue") == "raw_eda_image_fallback_used"
                    for warning in workflow_state.get("slide_validation_warnings", [])
                    if isinstance(warning, dict)
                )
            )
            self.assertIn("STC Analytics Review", full_text)
            self.assertIn("The dataset is ready", full_text)
            self.assertIn("Domain context", full_text)
            self.assertIn("Average Return Pct: 3.1", full_text)
            self.assertGreaterEqual(full_text.count("Sales trend improved in the latest period"), 1)
            self.assertGreaterEqual(full_text.count("Volume spikes align with price volatility"), 1)
            self.assertIn("Limitations define", full_text)

    def test_stock_objective_text_is_sanitized_on_deck(self) -> None:
        workflow_state = {
            "user_data_description": (
                "STC stock price history dataset. Produce dataset-specific EDA, time-series visuals, "
                "volatility and volume insights, a business-readable report, and executive slides."
            ),
            "agent_outputs": {
                "data_understander": {
                    "datasets": {
                        "STC Stock Price History (5).csv": {
                            "shape": [100, 7],
                            "columns": ["Date", "Price", "Open", "High", "Low", "Vol.", "Change %"],
                        }
                    }
                },
                "presentation_architect": {},
                "decision_maker": {},
                "business_translator": {},
            },
            "analysis_results": {"business_findings": ["Price trend supports validation gates."]},
        }

        deck = build_deck_spec(workflow_state)
        text = "\n".join([deck.subtitle] + [slide.subtitle for slide in deck.slides] + [slide.headline for slide in deck.slides])

        self.assertNotIn("Produce dataset-specific", text)
        self.assertNotIn("executive slides", text)
        self.assertIn("price, volume, return", deck.subtitle)

    def test_stock_visual_queue_demotes_volume_return_scatter(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            figures = []
            for name in ("price.png", "scatter.png", "volatility.png", "seasonal.png", "volume.png"):
                path = temp_path / name
                path.write_bytes(PNG_BYTES)
                figures.append(str(path))
            workflow_state = {
                "user_data_description": "STC stock price history review.",
                "saved_figures": figures,
                "analysis_results": {
                    "figure_captions": {
                        figures[0]: "Price trend shows support and resistance levels.",
                        figures[1]: "Volume-return scatter shows correlation strength: 0.17.",
                        figures[2]: "Volatility trend shows high-risk periods.",
                        figures[3]: "Seasonal return patterns show monthly differences.",
                        figures[4]: "Volume spike trend highlights liquidity windows.",
                    }
                },
                "agent_outputs": {
                    "data_understander": {
                        "datasets": {
                            "STC Stock Price History (5).csv": {
                                "shape": [100, 7],
                                "columns": ["Date", "Price", "Open", "High", "Low", "Vol.", "Change %"],
                            }
                        }
                    },
                    "presentation_architect": {},
                    "decision_maker": {},
                    "business_translator": {},
                },
            }

            deck = build_deck_spec(workflow_state)

        analysis_visuals = [
            slide.visual.image_path for slide in deck.slides if slide.slide_role == "analysis" and slide.visual
        ]
        self.assertNotIn("scatter.png", "\n".join(analysis_visuals[:3]))
        self.assertTrue(any("volume.png" in path for path in analysis_visuals[:3]))
        self.assertTrue(any("seasonal.png" in path for path in analysis_visuals[:4]))

    def test_generation_instruction_text_does_not_trigger_stock_deck_context(self) -> None:
        workflow_state = {
            "user_data_description": (
                "Teen mental health dataset. Produce dataset-specific EDA, clear visuals, "
                "decision-tree rules for depression_label, a business-readable report, and executive slides."
            ),
            "agent_outputs": {
                "data_understander": {
                    "datasets": {
                        "Teen_Mental_Health_Dataset.csv": {
                            "shape": [1200, 13],
                            "columns": ["age", "stress_level", "sleep_hours", "depression_label"],
                        }
                    }
                },
                "presentation_architect": {},
                "decision_maker": {},
                "business_translator": {},
            },
            "analysis_results": {"business_findings": ["Stress and sleep patterns shape outreach priorities."]},
        }

        deck = build_deck_spec(workflow_state)
        text = "\n".join([deck.subtitle] + [slide.headline for slide in deck.slides])

        self.assertNotIn("price, volume, return", text)
        self.assertNotIn("portfolio", text.lower())
        self.assertIn("Teen_Mental_Health_Dataset.csv", deck.subtitle)

    def test_legacy_visual_and_data_description_render_cleanly(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            figure_path = temp_path / "price_by_brand.png"
            figure_path.write_bytes(PNG_BYTES)
            output_path = temp_path / "deck.pptx"
            workflow_state = {
                "data_description": "Used car listings with price, mileage, brand, model year, and region.",
                "saved_figures": [str(figure_path)],
                "analysis_results": {"figure_captions": {str(figure_path): "Brand and mileage explain visible price dispersion."}},
                "agent_outputs": {
                    "presentation_architect": {
                        "presentation_title": "Used Car Pricing Review",
                        "slides": [
                            {
                                "slide_number": 1,
                                "title": "Price dispersion is concentrated by brand and mileage",
                                "main_message": "The strongest pricing signal appears in brand and mileage combinations.",
                                "details": ["Premium brands retain stronger prices.", "Higher mileage compresses listed prices."],
                                "visual_element": str(figure_path),
                            }
                        ],
                    },
                    "decision_maker": {},
                    "business_translator": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            picture_slides = sum(
                1
                for slide in presentation.slides
                if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
            )
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertEqual(len(presentation.slides), 12)
            self.assertEqual(picture_slides, 1)
            self.assertIn("Used car listings", full_text)
            self.assertIn("Price dispersion is concentrated", full_text)
            self.assertNotIn("[Visualization", full_text)

    def test_slide_validation_repairs_bad_layout_and_missing_visual(self) -> None:
        workflow_state = {
            "saved_figures": [],
            "analysis_results": {},
            "agent_outputs": {
                "presentation_architect": {
                    "slides": [
                        {
                            "slide_number": 1,
                            "layout_type": "unknown_layout",
                            "title": "This title is intentionally very long because the validator should shorten it before rendering the final consulting-style executive deck",
                            "main_message": "",
                            "details": [
                                "This detail is intentionally verbose and should be shortened so the resulting slide is readable for an executive audience rather than crowded with paragraph text.",
                                "Second point",
                                "Third point",
                                "Fourth point",
                                "Fifth point",
                            ],
                            "visual_path": "missing_figure.png",
                        }
                    ]
                },
                "decision_maker": {},
                "business_translator": {},
            },
        }

        slides = normalize_slide_plan(workflow_state)
        repaired = slides[-1]

        self.assertIn(repaired["layout_type"], {"limitations_professional", "executive_summary_closing", "three_finding_cards", "comparison_matrix"})
        self.assertLessEqual(len(repaired["title"]), 95)
        self.assertTrue(repaired["main_message"])
        self.assertLessEqual(len(repaired["details"]), 4)
        self.assertEqual(repaired["visual_path"], "")
        self.assertTrue(workflow_state["slide_validation_warnings"])

    def test_default_deck_spec_has_required_12_slide_order_and_varied_templates(self) -> None:
        workflow_state = {
            "saved_figures": [],
            "analysis_results": {
                "analysis_summary": {"default_rate": "12%", "loss_rate": 3.2},
                "business_findings": [
                    "Default risk is concentrated in high-utilization customers.",
                    "Lower income bands show higher observed default rates.",
                ],
            },
            "agent_outputs": {
                "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
                "presentation_architect": {"slides": []},
            },
        }

        deck = build_deck_spec(workflow_state)

        self.assertEqual(len(deck.slides), 12)
        self.assertEqual([slide.slide_role for slide in deck.slides], REQUIRED_SLIDE_ROLES)
        analysis_templates = [slide.template for slide in deck.slides if slide.slide_role == "analysis"]
        self.assertGreaterEqual(len(set(analysis_templates)), 2)
        self.assertIn("deck_spec", workflow_state)

    def test_analysis_slide_slots_use_four_structured_visuals_before_findings_and_business_translation(self) -> None:
        artifacts = []
        for index in range(4):
            artifacts.append(
                {
                    "artifact_id": f"eda_chart_{index + 1}",
                    "artifact_type": "chart_spec",
                    "slide_candidate": True,
                    "chart_type": "bar",
                    "title": f"EDA chart {index + 1}",
                    "finding": f"EDA chart {index + 1} identifies a decision signal.",
                    "takeaway": f"Signal {index + 1} should be interpreted before recommendations.",
                    "data": [
                        {"category": "Baseline", "value": 10 + index},
                        {"category": "Focus", "value": 18 + index},
                    ],
                    "recommended_template": "unsupported_chart_layout" if index == 0 else "",
                }
            )
        workflow_state = {
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": artifacts,
                "analysis_summary": {"sample": 100, "focus_gap": 8},
                "business_findings": ["Focus segments outperform the baseline."],
            },
            "agent_outputs": {
                "presentation_architect": {"slides": []},
                "business_translator": {
                    "key_findings": [
                        {
                            "finding": "The four EDA visuals point to a focused operating signal.",
                            "business_implication": "Use that signal to set the recommendation priority.",
                            "priority": "High",
                        }
                    ]
                },
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        analysis_slides = deck.slides[3:7]
        findings_slide = deck.slides[7]
        translation_slide = deck.slides[8]

        self.assertEqual([slide.slide_number for slide in analysis_slides], [4, 5, 6, 7])
        self.assertTrue(all(slide.slide_role == "analysis" for slide in analysis_slides))
        self.assertTrue(all(slide.visual and slide.visual.type == "structured_chart" for slide in analysis_slides))
        self.assertEqual([slide.visual.artifact_id for slide in analysis_slides], [artifact["artifact_id"] for artifact in artifacts])
        self.assertTrue(all(slide.template not in {"three_finding_cards", "comparison_matrix"} for slide in analysis_slides))
        self.assertNotIn("Translate this evidence", "\n".join(bullet for slide in analysis_slides for bullet in slide.bullets))
        self.assertEqual(findings_slide.slide_number, 8)
        self.assertEqual(findings_slide.slide_role, "findings")
        self.assertIsNone(findings_slide.visual)
        self.assertIn("EDA chart 1", "\n".join(findings_slide.bullets))
        self.assertEqual(translation_slide.slide_number, 9)
        self.assertEqual(translation_slide.slide_role, "business_translation")
        self.assertIsNone(translation_slide.visual)
        self.assertIn("four EDA visuals", "\n".join(translation_slide.bullets))
        self.assertEqual(deck.slides[-1].slide_role, "summary")
        self.assertEqual(deck.slides[-1].slide_number, 12)

    def test_decision_tree_artifact_renders_as_rule_boxes_after_eda(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            figure_paths = []
            for index in range(1, 4):
                figure_path = temp_path / f"figure_{index}.png"
                figure_path.write_bytes(PNG_BYTES)
                figure_paths.append(str(figure_path))
            output_path = temp_path / "decision_tree.pptx"
            workflow_state = {
                "decision_tree_target_column": "Attrition",
                "saved_figures": figure_paths,
                "analysis_results": {
                    "figure_captions": {path: f"EDA figure {index} shows a signal." for index, path in enumerate(figure_paths, start=1)},
                    "analysis_artifacts": [
                        {
                            "artifact_id": "decision_tree_rules",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "chart_type": "decision_tree",
                            "title": "Decision tree rules for Attrition",
                            "finding": "A shallow tree explains the main attrition split with 78% accuracy.",
                            "takeaway": "The model rules provide an interpretable retention triage signal.",
                            "data": {
                                "target": "Attrition",
                                "model_type": "classification",
                                "train_accuracy": "82%",
                                "test_accuracy": "78%",
                                "baseline_accuracy": "62%",
                                "model_verified": True,
                                "rules_match_model": True,
                                "nodes": [
                                    {"id": "root", "label": "OverTime <= 0.5", "depth": 0},
                                    {"id": "left", "label": "Low risk: predict No", "depth": 1},
                                    {"id": "right", "label": "MonthlyIncome <= 4200", "depth": 1},
                                    {"id": "right_left", "label": "High risk: predict Yes", "depth": 2},
                                ],
                                "edges": [
                                    {"source": "root", "target": "left", "label": "No"},
                                    {"source": "root", "target": "right", "label": "Yes"},
                                    {"source": "right", "target": "right_left", "label": "Low income"},
                                ],
                            },
                            "recommended_template": "full_width_chart_takeaway",
                        }
                    ],
                    "analysis_summary": {"decision_tree_accuracy": "78%", "sample": 100},
                    "business_findings": ["Decision tree accuracy is 78%.", "Overtime is the first model split."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertEqual(len(presentation.slides), 12)
            self.assertGreaterEqual(decision_tree_shape_count(presentation), 4)
            self.assertIn("Decision tree rules for Attrition", full_text)
            self.assertIn("Train Accuracy: 82%", full_text)
            self.assertIn("Test Accuracy: 78%", full_text)
            tree_labels = [
                shape
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "name", "") == "reconstructed_chart_decision_tree_label"
            ]
            tree_edges = [
                shape
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "name", "") == "reconstructed_chart_decision_tree_edge"
            ]
            root_shape = next(shape for shape in tree_labels if "OverTime" in shape.text)
            left_leaf = next(shape for shape in tree_labels if "Low risk" in shape.text)
            right_split = next(shape for shape in tree_labels if "MonthlyIncome" in shape.text)
            lower_leaf = next(shape for shape in tree_labels if "High risk" in shape.text)

            self.assertLess(root_shape.top, left_leaf.top)
            self.assertLess(root_shape.top, right_split.top)
            self.assertLess(right_split.top, lower_leaf.top)
            self.assertGreater(root_shape.left, min(left_leaf.left, right_split.left))
            self.assertLess(root_shape.left, max(left_leaf.left, right_split.left))
            self.assertGreaterEqual(len(tree_edges), 9)

    def test_decision_tree_artifact_uses_shared_saved_png_when_available(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            eda_figures = []
            for index in range(1, 4):
                figure_path = temp_path / f"figure_{index}.png"
                figure_path.write_bytes(PNG_BYTES)
                eda_figures.append(str(figure_path))
            tree_path = temp_path / "decision_tree_rules.png"
            tree_path.write_bytes(PNG_BYTES)
            output_path = temp_path / "decision_tree_image.pptx"
            workflow_state = {
                "decision_tree_target_column": "Attrition",
                "saved_figures": eda_figures + [str(tree_path)],
                "analysis_results": {
                    "figure_captions": {
                        **{path: f"EDA figure {index} shows a signal." for index, path in enumerate(eda_figures, start=1)},
                        str(tree_path): "Decision tree leaf rules explain the model.",
                    },
                    "analysis_artifacts": [
                        {
                            "artifact_id": "decision_tree_rules",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "chart_type": "decision_tree",
                            "title": "Decision tree rules for Attrition",
                            "finding": "Leaf nodes show the model logic.",
                            "fallback_path": str(tree_path),
                            "data": {
                                "model_type": "classification",
                                "train_accuracy": "88%",
                                "test_accuracy": "80%",
                                "baseline_accuracy": "72%",
                                "model_verified": True,
                                "rules_match_model": True,
                                "nodes": [
                                    {"id": "root", "type": "split", "label": "OverTime <= 0.5", "depth": 0},
                                    {"id": "left", "type": "leaf", "label": "Leaf: predict No when OverTime <= 0.5", "depth": 1},
                                    {"id": "right", "type": "leaf", "label": "Leaf: predict Yes when OverTime > 0.5", "depth": 1},
                                ],
                                "edges": [
                                    {"source": "root", "target": "left", "label": "True"},
                                    {"source": "root", "target": "right", "label": "False"},
                                ],
                            },
                            "recommended_template": "full_width_chart_takeaway",
                        }
                    ],
                    "analysis_summary": {"decision_tree_accuracy": "80%", "sample": 100},
                    "business_findings": ["Decision tree accuracy is 80%.", "Leaf rules explain the model."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            tree_visual = next(
                slide.visual
                for slide in deck.slides
                if slide.visual and slide.visual.artifact_id == "decision_tree_rules"
            )
            self.assertEqual(tree_visual.type, "structured_chart")
            self.assertEqual(tree_visual.chart_type, "decision_tree")
            self.assertEqual(tree_visual.fallback_path, str(tree_path))

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            picture_count = sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            self.assertGreaterEqual(picture_count, 4)
            self.assertEqual(decision_tree_shape_count(presentation), 0)
            self.assertTrue(
                any(
                    getattr(shape, "name", "") == "shared_decision_tree_image"
                    for slide in presentation.slides
                    for shape in slide.shapes
                )
            )
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("Train Accuracy: 88%", full_text)
            self.assertIn("Test Accuracy: 80%", full_text)

    def test_decision_tree_structured_artifact_is_queued_after_eda_artifacts(self) -> None:
        artifacts = []
        for index in range(3):
            artifacts.append(
                {
                    "artifact_id": f"eda_{index + 1}",
                    "artifact_type": "chart_spec",
                    "slide_candidate": True,
                    "chart_type": "bar",
                    "title": f"EDA finding {index + 1}",
                    "finding": f"EDA finding {index + 1} explains the data before modeling.",
                    "data": [{"category": "A", "value": 10}, {"category": "B", "value": 14}],
                }
            )
        artifacts.append(
            {
                "artifact_id": "decision_tree_rules",
                "artifact_type": "chart_spec",
                "slide_candidate": True,
                "chart_type": "decision_tree",
                "title": "Decision tree rules for Attrition",
                "finding": "The decision tree follows the EDA findings.",
                "data": {
                    "accuracy": "78%",
                    "nodes": [
                        {"id": "root", "label": "OverTime <= 0.5", "depth": 0},
                        {"id": "left", "label": "Predict No", "depth": 1},
                    ],
                    "edges": [{"source": "root", "target": "left", "label": "No"}],
                },
            }
        )
        workflow_state = {
            "decision_tree_target_column": "Attrition",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": artifacts,
                "analysis_summary": {"decision_tree_accuracy": "78%", "sample": 100},
                "business_findings": ["Decision tree accuracy is 78%."],
            },
            "agent_outputs": {
                "presentation_architect": {"slides": []},
                "data_understander": {"executive_summary": "Employee dataset", "datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        analysis_visual_ids = [slide.visual.artifact_id for slide in deck.slides[3:7]]

        self.assertEqual(analysis_visual_ids, ["eda_1", "eda_2", "eda_3", "decision_tree_rules"])

    def test_data_understanding_slide_does_not_duplicate_quality_notes(self) -> None:
        workflow_state = {
            "agent_outputs": {
                "presentation_architect": {"presentation_title": "Data Quality Review", "slides": []},
                "data_understander": {
                    "executive_summary": "Employee dataset supports attrition review.",
                    "datasets": {
                        "hr.csv": {
                            "shape": [100, 6],
                            "columns": ["Attrition", "OverTime", "JobRole"],
                            "quality_summary": "No missing values detected.",
                            "cleaning_priorities": ["No missing values detected."],
                            "type_notes": [],
                        }
                    },
                },
            },
            "analysis_results": {"analysis_artifacts": [], "business_findings": []},
        }

        deck = build_deck_spec(workflow_state)
        slide = deck.slides[1]
        bullet_text = "\n".join(slide.bullets)

        self.assertNotIn("No missing values detected", bullet_text)
        self.assertIn("No missing values detected", deck.dataset_context["data_quality_notes"])

    def test_slide_polish_removes_near_duplicate_insight_text(self) -> None:
        items = _dedupe_text_items(
            [
                "Lower satisfaction scores across all dimensions correlate strongly with higher attrition rates.",
                "Lower satisfaction scores across all dimensions correlate strongly with higher attrition rates...",
                "Attrition rate by job satisfaction level",
            ]
        )

        self.assertEqual(
            items,
            [
                "Lower satisfaction scores across all dimensions correlate strongly with higher attrition rates.",
                "Attrition rate by job satisfaction level",
            ],
        )

    def test_saved_figure_captions_can_use_basename_keys(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            figure_path = Path(temp_dir) / "figure_1.png"
            figure_path.write_bytes(PNG_BYTES)
            workflow_state = {
                "saved_figures": [str(figure_path)],
                "analysis_results": {
                    "figure_captions": {
                        "figure_1.png": "Overtime employees show the highest observed attrition rate."
                    }
                },
                "agent_outputs": {
                    "presentation_architect": {"slides": []},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            first_analysis = deck.slides[3]

        self.assertIsNotNone(first_analysis.visual)
        self.assertIn("Overtime employees", first_analysis.visual.takeaway)
        self.assertIn("Overtime employees", first_analysis.headline)

    def test_analysis_slide_headline_follows_visual_not_stale_legacy_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            figure_path = Path(temp_dir) / "figure_1.png"
            figure_path.write_bytes(PNG_BYTES)
            workflow_state = {
                "saved_figures": [str(figure_path)],
                "analysis_results": {
                    "figure_captions": {
                        str(figure_path): "Travel frequency strongly correlates with attrition risk."
                    }
                },
                "agent_outputs": {
                    "presentation_architect": {
                        "slides": [
                            {
                                "slide_number": 4,
                                "slide_role": "analysis",
                                "headline": "Overtime workers face the highest attrition rate",
                                "main_message": "Overtime is the key visual.",
                                "details": ["Overtime detail"],
                            }
                        ]
                    },
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            first_analysis = deck.slides[3]

        self.assertIn("Travel frequency", first_analysis.headline)
        self.assertIn("Travel frequency", first_analysis.main_message)
        self.assertNotIn("Overtime", first_analysis.headline)

    def test_code_figure_headline_uses_caption_even_when_generic_tokens_overlap(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            figure_path = Path(temp_dir) / "figure_1.png"
            figure_path.write_bytes(PNG_BYTES)
            workflow_state = {
                "saved_figures": [str(figure_path)],
                "analysis_results": {
                    "figure_captions": {
                        str(figure_path): "Correlation matrix reveals relationships between price metrics, volume, and daily returns."
                    }
                },
                "agent_outputs": {
                    "presentation_architect": {
                        "slides": [
                            {
                                "slide_number": 4,
                                "slide_role": "analysis",
                                "headline": "Daily returns show moderate volatility with an average change",
                                "main_message": "Daily returns are important.",
                            }
                        ]
                    },
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            first_analysis = deck.slides[3]

        self.assertIn("Correlation matrix", first_analysis.headline)
        self.assertNotIn("moderate volatility", first_analysis.headline.lower())

    def test_decision_tree_headline_uses_verified_tree_splits_not_stale_finding(self) -> None:
        workflow_state = {
            "decision_tree_target_column": "Attrition",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": [
                    {
                        "artifact_id": "decision_tree_rules",
                        "artifact_type": "chart_spec",
                        "slide_candidate": True,
                        "chart_type": "decision_tree",
                        "title": "Decision Tree Rules for Attrition Prediction",
                        "finding": "Tree identifies overtime, income, promotion delay as key rules.",
                        "takeaway": "Use these rules to identify at-risk employees.",
                        "data": {
                            "target": "Attrition",
                            "model_type": "classification",
                            "test_accuracy": "76.2%",
                            "baseline_accuracy": "84.0%",
                            "performance_note": (
                                "Explanatory model only: test accuracy 76.2% trails the baseline 84.0%; "
                                "use the rules to guide investigation, not production prediction."
                            ),
                            "nodes": [
                                {
                                    "id": "root",
                                    "type": "split",
                                    "feature": "TotalWorkingYears",
                                    "label": "TotalWorkingYears <= -1.14",
                                },
                                {
                                    "id": "role",
                                    "type": "split",
                                    "feature": "JobRole_Research Scientist",
                                    "label": "JobRole_Research Scientist <= 0.5",
                                },
                                {"id": "leaf_1", "type": "leaf", "label": "Leaf: predict Yes"},
                                {"id": "leaf_2", "type": "leaf", "label": "Leaf: predict No"},
                            ],
                            "edges": [
                                {"source": "root", "target": "role", "label": "True"},
                                {"source": "role", "target": "leaf_1", "label": "True"},
                                {"source": "root", "target": "leaf_2", "label": "False"},
                            ],
                            "rules_match_model": True,
                            "model_verified": True,
                        },
                    }
                ],
                "analysis_summary": {"decision_tree_accuracy": "76.2%", "sample": 100},
                "business_findings": ["Decision tree is explanatory only."],
            },
            "agent_outputs": {
                "presentation_architect": {"slides": []},
                "data_understander": {"datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        tree_slide = next(slide for slide in deck.slides if slide.visual and slide.visual.chart_type == "decision_tree")

        self.assertIn("TotalWorkingYears", tree_slide.headline)
        self.assertIn("JobRole Research Scientist", tree_slide.headline)
        self.assertNotIn("income", tree_slide.headline.lower())
        self.assertNotIn("promotion", tree_slide.headline.lower())
        self.assertIn("Explanatory model only", tree_slide.main_message)

    def test_code_figure_headline_replaces_unsupported_ratio_with_visible_rates(self) -> None:
        workflow_state = {
            "saved_figures": ["figure_2.png"],
            "analysis_results": {
                "business_findings": ["Female teens show higher depression rates (3.5%) compared to males (1.6%)."],
                "figure_captions": {
                    "figure_2.png": "Female teens show higher depression rates (3.5%) compared to males (1.6%)."
                },
                "analysis_artifacts": [],
            },
            "agent_outputs": {
                "presentation_architect": {
                    "slides": [
                        {
                            "slide_number": 4,
                            "slide_role": "analysis",
                            "headline": "Female Teens Exhibit 2.2x Higher Depression Rates Than Males",
                            "main_message": "Female teens show higher depression rates.",
                        }
                    ]
                },
                "data_understander": {"datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        analysis_slide = next(slide for slide in deck.slides if slide.slide_role == "analysis")

        self.assertIn("3.5%", analysis_slide.headline)
        self.assertIn("1.6%", analysis_slide.headline)
        self.assertNotIn("2.2x", analysis_slide.headline)

    def test_non_analysis_slide_does_not_inherit_stale_decision_tree_visual(self) -> None:
        workflow_state = {
            "analysis_results": {
                "business_findings": ["Overtime and tenure explain the strongest retention signals."],
            },
            "agent_outputs": {
                "presentation_architect": {
                    "slides": [
                        {
                            "slide_number": 12,
                            "slide_role": "summary",
                            "headline": "Act on the clearest retention evidence",
                            "main_message": "Focus the next sprint on operational retention changes.",
                            "visual": {
                                "type": "structured_chart",
                                "chart_type": "decision_tree",
                                "artifact_id": "decision_tree_rules",
                                "title": "Decision Tree Rules for Attrition Prediction",
                                "finding": "Tree achieves 79.9% test accuracy with four interpretable rules.",
                                "data": {
                                    "target": "Attrition",
                                    "test_accuracy": 0.799,
                                    "baseline_accuracy": 0.838,
                                },
                            },
                        }
                    ]
                },
                "business_translator": {},
                "decision_maker": {
                    "conclusion": "Focus the next sprint on operational retention changes.",
                },
            },
        }

        deck = build_deck_spec(workflow_state)
        summary_slide = deck.slides[11]

        self.assertEqual(summary_slide.slide_role, "summary")
        self.assertIsNone(summary_slide.visual)
        self.assertNotIn("decision tree", summary_slide.headline.lower())

    def test_tree_findings_and_translation_ignore_stale_legacy_story_text(self) -> None:
        workflow_state = {
            "decision_tree_target_column": "Attrition",
            "user_data_description": "HR attrition dataset.",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": [
                    {
                        "artifact_id": "decision_tree_rules",
                        "artifact_type": "chart_spec",
                        "slide_candidate": True,
                        "chart_type": "decision_tree",
                        "title": "Attrition Decision Tree Rules",
                        "data": {
                            "target": "Attrition",
                            "model_type": "classification",
                            "test_accuracy": "81.9%",
                            "baseline_accuracy": "83.9%",
                            "nodes": [
                                {"id": "root", "type": "split", "feature": "TotalWorkingYears", "threshold": "-1.28"},
                                {"id": "leaf", "type": "leaf", "prediction": "Yes", "samples": 24},
                                {"id": "leaf2", "type": "leaf", "prediction": "No", "samples": 11},
                            ],
                            "edges": [
                                {"source": "root", "target": "leaf", "label": "True"},
                                {"source": "root", "target": "leaf2", "label": "False"},
                            ],
                            "rules": ["If TotalWorkingYears <= -1.28, then predict Yes."],
                        },
                    }
                ],
                "business_findings": ["Tree rules are explanatory."],
            },
            "agent_outputs": {
                "presentation_architect": {
                    "slides": [
                        {
                            "slide_number": 8,
                            "slide_role": "findings",
                            "headline": "Decision tree rules summarize model logic",
                            "main_message": "Generic model text.",
                        },
                        {
                            "slide_number": 9,
                            "slide_role": "business_translation",
                            "headline": "High-risk students need counselor review",
                            "main_message": "Generic school text.",
                        },
                    ]
                },
                "data_understander": {"datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        findings_slide = next(slide for slide in deck.slides if slide.slide_role == "findings")
        translation_slide = next(slide for slide in deck.slides if slide.slide_role == "business_translation")
        translation_text = " ".join(item for block in translation_slide.content_blocks for item in block.items)

        self.assertIn("Exploratory tree rules", findings_slide.headline)
        self.assertNotIn("summarize model logic", findings_slide.headline)
        self.assertIn("employees", translation_text)
        self.assertIn("manager and HR review", translation_text)
        self.assertNotIn("students", translation_text)
        self.assertNotIn("counselor", translation_text)

    def test_credit_tree_rules_are_business_readable_and_domain_specific(self) -> None:
        workflow_state = {
            "decision_tree_target_column": "loan_status",
            "user_data_description": "Credit risk dataset.",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": [
                    {
                        "artifact_id": "decision_tree_rules",
                        "artifact_type": "chart_spec",
                        "slide_candidate": True,
                        "chart_type": "decision_tree",
                        "title": "Loan Default Tree Rules",
                        "data": {
                            "target": "loan_status",
                            "model_type": "classification",
                            "test_accuracy": "88.2%",
                            "baseline_accuracy": "78.1%",
                            "nodes": [
                                {"id": "root", "type": "split", "feature": "person_income", "threshold": "19984"},
                                {"id": "leaf", "type": "leaf", "prediction": "Default", "samples": 21},
                                {"id": "leaf2", "type": "leaf", "prediction": "No Default", "samples": 44},
                            ],
                            "edges": [
                                {"source": "root", "target": "leaf", "label": "True", "condition": "person_income <= 19984"},
                                {"source": "root", "target": "leaf2", "label": "False", "condition": "person_income > 19984"},
                            ],
                            "rules": ["If person_income <= 19984, then predict Default."],
                        },
                    }
                ],
            },
            "agent_outputs": {
                "presentation_architect": {"slides": []},
                "data_understander": {"datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        findings_slide = next(slide for slide in deck.slides if slide.slide_role == "findings")
        translation_slide = next(slide for slide in deck.slides if slide.slide_role == "business_translation")
        findings_text = " ".join(item for block in findings_slide.content_blocks for item in block.items)
        translation_text = " ".join(item for block in translation_slide.content_blocks for item in block.items)

        self.assertIn("borrower income at or below $19,984", findings_text)
        self.assertNotIn("person_income <=", findings_text)
        self.assertIn("applications", translation_text)
        self.assertIn("underwriting or portfolio review", translation_text)

    def test_rendered_credit_tree_node_labels_are_business_readable(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "credit_tree_labels.pptx"
            deck = DeckSpec(
                deck_title="Credit Risk Review",
                slides=[
                    SlideSpec(
                        1,
                        "analysis",
                        "full_width_chart_takeaway",
                        "Decision tree splits on borrower income",
                        visual=VisualSpec(
                            type="structured_chart",
                            chart_type="decision_tree",
                            title="Loan default decision tree",
                            data={
                                "target": "loan_status",
                                "nodes": [
                                    {"id": "root", "type": "split", "feature": "person_income", "threshold": "19984"},
                                    {"id": "left", "type": "leaf", "prediction": "Default", "samples": 20},
                                    {"id": "right", "type": "leaf", "prediction": "No Default", "samples": 40},
                                ],
                                "edges": [
                                    {"source": "root", "target": "left", "label": "True"},
                                    {"source": "root", "target": "right", "label": "False"},
                                ],
                            },
                        ),
                    )
                ],
            )

            presentation_obj = Presentation()
            renderer = PowerPointRenderer(presentation_obj, deck, THEME)
            renderer.render_slide(deck.slides[0])
            presentation_obj.save(str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertIn("borrower income at or below 19984", full_text)
            self.assertNotIn("person_income", full_text)

    def test_credit_tree_callouts_use_workflow_target_when_artifact_target_missing(self) -> None:
        workflow_state = {
            "decision_tree_target_column": "loan_status",
            "user_data_description": "Credit risk dataset.",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": [
                    {
                        "artifact_id": "decision_tree_rules",
                        "artifact_type": "chart_spec",
                        "slide_candidate": True,
                        "chart_type": "decision_tree",
                        "title": "Loan Default Tree Rules",
                        "data": {
                            "model_type": "classification",
                            "nodes": [
                                {"id": "root", "type": "split", "feature": "person_income", "threshold": "19984"},
                                {"id": "leaf", "type": "leaf", "prediction": "Default", "samples": 21},
                                {"id": "leaf2", "type": "leaf", "prediction": "No Default", "samples": 44},
                            ],
                            "edges": [
                                {"source": "root", "target": "leaf", "label": "True", "condition": "person_income <= 19984"},
                                {"source": "root", "target": "leaf2", "label": "False", "condition": "person_income > 19984"},
                            ],
                            "rules": ["If person_income <= 19984, then predict Default."],
                        },
                    }
                ],
            },
            "agent_outputs": {
                "presentation_architect": {"slides": []},
                "data_understander": {"datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        findings_slide = next(slide for slide in deck.slides if slide.slide_role == "findings")
        findings_text = " ".join(item for block in findings_slide.content_blocks for item in block.items)

        self.assertIn("borrower income at or below $19,984", findings_text)
        self.assertNotIn("person_income <=", findings_text)

    def test_analysis_visual_queue_uses_structured_target_chart_before_weak_third_code_figure(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            paths = []
            for name in ("home.png", "intent.png", "age.png"):
                path = Path(temp_dir) / name
                path.write_bytes(PNG_BYTES)
                paths.append(str(path))
            workflow_state = {
                "decision_tree_target_column": "loan_status",
                "saved_figures": paths,
                "analysis_results": {
                    "figure_captions": {
                        paths[0]: "Default rate is highest by home ownership.",
                        paths[1]: "Default rate varies by loan intent.",
                        paths[2]: "Median borrower age varies by home ownership.",
                    },
                    "analysis_artifacts": [
                        {
                            "artifact_id": "default_by_grade",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "chart_type": "bar",
                            "title": "Default Rate by Loan Grade",
                            "finding": "Default rate increases with worse loan grade.",
                            "x": "loan_grade",
                            "y": "default_rate",
                            "data": [{"loan_grade": "A", "default_rate": 0.05}, {"loan_grade": "E", "default_rate": 0.28}],
                        },
                        {
                            "artifact_id": "decision_tree_rules",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "chart_type": "decision_tree",
                            "title": "Loan Default Tree Rules",
                            "data": {
                                "target": "loan_status",
                                "nodes": [
                                    {"id": "root", "type": "split", "feature": "loan_percent_income", "threshold": "0.34"},
                                    {"id": "left", "type": "leaf", "prediction": "Default"},
                                    {"id": "right", "type": "leaf", "prediction": "No Default"},
                                ],
                                "edges": [
                                    {"source": "root", "target": "left", "label": "True"},
                                    {"source": "root", "target": "right", "label": "False"},
                                ],
                            },
                        },
                    ],
                },
                "agent_outputs": {
                    "presentation_architect": {"slides": []},
                    "data_understander": {"datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            analysis_headlines = [slide.headline for slide in deck.slides if slide.slide_role == "analysis"]
            analysis_visual_titles = [
                slide.visual.title for slide in deck.slides if slide.slide_role == "analysis" and slide.visual
            ]

        self.assertTrue(any("Loan Grade" in title for title in analysis_visual_titles))
        self.assertFalse(any("Median borrower age" in headline for headline in analysis_headlines[:3]))

    def test_softens_unsupported_attrition_reduction_claims(self) -> None:
        softened = soften_unsupported_impact_claim(
            "Targeted interventions on overtime, role, and income can reduce attrition."
        )

        self.assertIn("can focus attrition validation pilots", softened)
        self.assertNotIn("can reduce attrition", softened)

    def test_business_translation_uses_distinct_template_and_content_from_findings(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "business_translation_matrix.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "business_findings": [
                        "Finding one shows a signal.",
                        "Finding two adds evidence.",
                        "Finding three narrows the decision.",
                    ],
                },
                "agent_outputs": {
                    "presentation_architect": {"slides": []},
                    "business_translator": {
                        "key_findings": [
                            "Finding four adds business meaning.",
                            "Finding five sets the priority.",
                            "Finding six defines the next check.",
                        ],
                    },
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            findings_text = "\n".join(shape.text for shape in presentation.slides[7].shapes if hasattr(shape, "text"))
            translation_text = "\n".join(shape.text for shape in presentation.slides[8].shapes if hasattr(shape, "text"))

            self.assertIn("Finding 1", findings_text)
            self.assertIn("Finding 2", findings_text)
            self.assertIn("Finding 3", findings_text)
            self.assertIn("Decision use", translation_text)
            self.assertIn("Tradeoff", translation_text)
            self.assertIn("Operating rule", translation_text)
            self.assertIn("Validation", translation_text)
            self.assertNotIn("Finding 1\n", translation_text)

    def test_required_template_registry_entries_exist(self) -> None:
        expected = {
            "title_cover",
            "data_understanding_overview",
            "market_context_bullets",
            "chart_left_insight_right",
            "chart_right_insight_left",
            "full_width_chart_takeaway",
            "metric_strip_plus_chart",
            "three_finding_cards",
            "comparison_matrix",
            "recommendation_priority",
            "limitations_professional",
            "executive_summary_closing",
            "small_multiples_with_takeaway",
            "single_bar_chart_with_insight",
            "horizontal_bar_ranking",
            "metric_cards_with_chart",
            "comparison_chart_with_interpretation",
            "distribution_with_callout",
            "segment_profile_cards",
        }

        self.assertTrue(expected.issubset(set(TEMPLATE_REGISTRY)))

    def test_native_chart_spec_renders_chart_shape(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "native_chart.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "risk_by_grade",
                            "chart_type": "bar",
                            "title": "Risk by grade",
                            "takeaway": "Grade C carries the highest default risk.",
                            "x": "grade",
                            "y": "default_rate",
                            "data": [
                                {"grade": "A", "default_rate": 0.04},
                                {"grade": "B", "default_rate": 0.08},
                                {"grade": "C", "default_rate": 0.16},
                            ],
                        }
                    ],
                    "analysis_summary": {"default_rate": "16% in Grade C", "sample": 3},
                    "business_findings": ["Grade C carries the highest default risk."],
                },
                "agent_outputs": {
                    "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                    "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))

            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 1)

    def test_slide_generation_saves_structured_plan_and_chart_artifacts(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "deck.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "risk_by_grade",
                            "chart_type": "bar",
                            "title": "Risk by grade",
                            "finding": "Grade C carries the highest default risk.",
                            "x": "grade",
                            "y": "default_rate",
                            "data": [
                                {"grade": "A", "default_rate": 0.04},
                                {"grade": "C", "default_rate": 0.16},
                            ],
                        }
                    ],
                    "analysis_summary": {"grade_c_default_rate": "16%", "sample": 2},
                    "business_findings": ["Grade C carries the highest default risk."],
                },
                "agent_outputs": {
                    "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                    "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            slide_plan_path = Path(temp_dir) / "slide_plan.json"
            chart_specs_path = Path(temp_dir) / "chart_specs.json"
            slide_plan = json.loads(slide_plan_path.read_text(encoding="utf-8"))
            chart_specs = json.loads(chart_specs_path.read_text(encoding="utf-8"))

            self.assertTrue(slide_plan_path.exists())
            self.assertTrue(chart_specs_path.exists())
            self.assertEqual(workflow_state["generated_reports"]["slide_plan"], str(slide_plan_path))
            self.assertEqual(workflow_state["generated_reports"]["chart_specs"], str(chart_specs_path))
            self.assertEqual(len(slide_plan["slides"]), 12)
            self.assertEqual(chart_specs["chart_specs"][0]["id"], "risk_by_grade")
            self.assertTrue(chart_specs["slide_visuals"])

    def test_saved_figures_take_priority_over_structured_chart_specs(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            figure_path = temp_path / "figure_1.png"
            figure_path.write_bytes(PNG_BYTES)
            output_path = temp_path / "figure_primary.pptx"
            workflow_state = {
                "saved_figures": [str(figure_path)],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "attrition_by_role",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "finding": "Sales roles show the highest attrition rate.",
                            "chart_type": "bar",
                            "title": "Attrition by role",
                            "takeaway": "Sales roles require the first retention focus.",
                            "x": "role",
                            "y": "attrition_rate",
                            "data": [
                                {"role": "Research", "attrition_rate": 9.0},
                                {"role": "Sales", "attrition_rate": 24.0},
                                {"role": "Operations", "attrition_rate": 12.0},
                            ],
                        }
                    ],
                    "analysis_summary": {"sales_attrition_rate": "24%", "sample": 3},
                    "business_findings": ["Sales roles show the highest attrition rate."],
                    "figure_captions": {str(figure_path): "Code figure should match the PDF visual."},
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee attrition dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            picture_count = sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            self.assertEqual(picture_count, 1)
            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 1)

    def test_structured_eda_precedes_dense_code_figures_when_tree_exists(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            dense_figure = temp_path / "figure_dense_segments.png"
            dense_figure.write_bytes(PNG_BYTES)
            structured = [
                {
                    "artifact_id": artifact_id,
                    "chart_type": "bar",
                    "title": title,
                    "finding": finding,
                    "x": "category",
                    "y": "value",
                    "data": [{"category": "A", "value": 1}, {"category": "B", "value": 2}],
                }
                for artifact_id, title, finding in (
                    ("attrition_by_department", "Attrition by department", "Department attrition varies."),
                    ("overtime_attrition", "Overtime attrition", "Overtime is higher."),
                    ("satisfaction_attrition", "Job satisfaction attrition", "Low satisfaction is higher."),
                )
            ]
            tree = {
                "artifact_id": "decision_tree_rules",
                "chart_type": "decision_tree",
                "title": "Decision tree rules",
                "data": {
                    "nodes": [
                        {"id": "root", "type": "split", "label": "OverTime <= 0.5"},
                        {"id": "left", "type": "leaf", "label": "Predict No"},
                        {"id": "right", "type": "leaf", "label": "Predict Yes"},
                    ],
                    "edges": [
                        {"source": "root", "target": "left", "label": "True"},
                        {"source": "root", "target": "right", "label": "False"},
                    ],
                },
            }
            workflow_state = {
                "decision_tree_target_column": "Attrition",
                "saved_figures": [str(dense_figure)],
                "analysis_results": {
                    "analysis_artifacts": structured + [tree],
                    "figure_captions": {str(dense_figure): "Multi-panel segment chart mixes several HR dimensions."},
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee attrition dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            deck = build_deck_spec(workflow_state)
            visuals = [slide.visual for slide in deck.slides if slide.slide_role == "analysis" and slide.visual]
            self.assertEqual([visual.artifact_id for visual in visuals[:3]], [item["artifact_id"] for item in structured])
            self.assertTrue(all(visual.type == "structured_chart" for visual in visuals[:3]))
            self.assertEqual(visuals[3].artifact_id, "decision_tree_rules")

    def test_duplicate_structured_eda_artifacts_do_not_repeat_analysis_slots(self) -> None:
        duplicated_grade = [
            {
                "artifact_id": artifact_id,
                "chart_type": "bar",
                "title": "Default Rate by Loan Grade",
                "finding": finding,
                "x": "category",
                "y": "value",
                "data": [{"category": "A", "value": 7.0}, {"category": "G", "value": 35.0}],
            }
            for artifact_id, finding in (
                ("risk_by_grade", "Loan grade is the strongest predictor."),
                ("grade_risk", "Grade E loans default more often than Grade A."),
            )
        ]
        distinct = [
            {
                "artifact_id": "dti_threshold",
                "chart_type": "bar",
                "title": "Default Rate by DTI Bucket",
                "finding": "Debt burden separates borrower risk.",
                "x": "category",
                "y": "value",
                "data": [{"category": "Low", "value": 8.0}, {"category": "High", "value": 24.0}],
            },
            {
                "artifact_id": "prior_default",
                "chart_type": "bar",
                "title": "Default Rate by Prior Default Status",
                "finding": "Prior defaults remain an important risk segment.",
                "x": "category",
                "y": "value",
                "data": [{"category": "N", "value": 10.0}, {"category": "Y", "value": 31.0}],
            },
        ]
        tree = {
            "artifact_id": "decision_tree_rules",
            "chart_type": "decision_tree",
            "title": "Decision tree rules",
            "data": {
                "nodes": [
                    {"id": "root", "type": "split", "label": "loan_grade <= C"},
                    {"id": "left", "type": "leaf", "label": "Predict No Default"},
                    {"id": "right", "type": "leaf", "label": "Predict Default"},
                ],
                "edges": [
                    {"source": "root", "target": "left", "label": "True"},
                    {"source": "root", "target": "right", "label": "False"},
                ],
            },
        }
        workflow_state = {
            "decision_tree_target_column": "loan_status",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": duplicated_grade + distinct + [tree],
                "business_findings": ["Loan grade and debt burden guide validation."],
            },
            "agent_outputs": {
                "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        analysis_visuals = [slide.visual for slide in deck.slides if slide.slide_role == "analysis" and slide.visual]
        analysis_titles = [visual.title for visual in analysis_visuals]

        self.assertEqual(analysis_titles.count("Default Rate by Loan Grade"), 1)
        self.assertIn("Default Rate by DTI Bucket", analysis_titles)
        self.assertIn("Default Rate by Prior Default Status", analysis_titles)

    def test_findings_slide_remains_eda_summary_when_tree_slide_exists(self) -> None:
        workflow_state = {
            "decision_tree_target_column": "loan_status",
            "saved_figures": [],
            "analysis_results": {
                "analysis_artifacts": [
                    {
                        "artifact_id": "dti_threshold",
                        "chart_type": "bar",
                        "title": "Default Rate by DTI Bucket",
                        "finding": "EDA driver: high debt-burden borrowers default more often.",
                        "data": [{"category": "Low", "value": 8.0}, {"category": "High", "value": 24.0}],
                    },
                    {
                        "artifact_id": "decision_tree_rules",
                        "chart_type": "decision_tree",
                        "title": "Decision tree rules",
                        "data": {
                            "nodes": [
                                {"id": "root", "type": "split", "label": "loan_grade <= C"},
                                {"id": "left", "type": "leaf", "label": "Predict No Default"},
                                {"id": "right", "type": "leaf", "label": "Predict Default"},
                            ],
                            "edges": [
                                {"source": "root", "target": "left", "label": "True"},
                                {"source": "root", "target": "right", "label": "False"},
                            ],
                        },
                    },
                ],
                "business_findings": ["EDA driver: debt burden separates borrower risk."],
            },
            "agent_outputs": {
                "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        deck = build_deck_spec(workflow_state)
        findings_slide = next(slide for slide in deck.slides if slide.slide_role == "findings")
        findings_text = " ".join(item for block in findings_slide.content_blocks for item in block.items)

        self.assertIn("debt-burden borrowers", findings_text)
        self.assertNotIn("Exploratory tree rules", findings_slide.headline)

    def test_category_value_chart_spec_renders_rebuilt_chart(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "category_value.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "satisfaction_attrition",
                            "chart_type": "bar",
                            "title": "Attrition rate by satisfaction level",
                            "finding": "Low satisfaction groups show consistently higher attrition risk.",
                            "takeaway": "Attrition is highest among employees with the lowest satisfaction scores.",
                            "data": [
                                {"category": "Level 1", "value": 25.4},
                                {"category": "Level 2", "value": 15.0},
                                {"category": "Level 3", "value": 13.7},
                                {"category": "Level 4", "value": 13.5},
                            ],
                            "x_label": "Satisfaction level",
                            "y_label": "Attrition rate (%)",
                            "value_format": "{:.1f}%",
                        }
                    ],
                    "analysis_summary": {"lowest_satisfaction_attrition": "25.4%", "sample": 4},
                    "business_findings": ["Low satisfaction groups show consistently higher attrition risk."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee attrition dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 4)
            self.assertIn("25.4%", full_text)
            self.assertIn("Attrition rate by satisfaction level", full_text)

    def test_series_only_bar_chart_renders_without_unavailable_placeholder(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "series_bar.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "risk_by_grade",
                            "chart_type": "bar",
                            "title": "Default rate by loan grade",
                            "finding": "Loan grade is the strongest predictor.",
                            "takeaway": "Grade G has the highest default rate.",
                            "x_label": "Loan grade",
                            "y_label": "Default rate (%)",
                            "series": [
                                {
                                    "name": "Default Rate",
                                    "data": [
                                        {"category": "A", "value": 7.0},
                                        {"category": "B", "value": 13.0},
                                        {"category": "G", "value": 35.0},
                                    ],
                                }
                            ],
                        }
                    ],
                    "business_findings": ["Loan grade is the strongest predictor."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                    "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertNotIn("Structured chart data was unavailable", full_text)
            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 1)

    def test_series_only_horizontal_bar_and_line_render_without_placeholder(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "series_charts.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "chart_specs": [
                        {
                            "id": "ranked_default_risk",
                            "chart_type": "horizontal_bar",
                            "title": "Default risk ranking",
                            "finding": "Loan grades differ materially in default risk.",
                            "x_label": "Loan grade",
                            "y_label": "Default rate (%)",
                            "series": [
                                {
                                    "name": "Default rate",
                                    "data": [
                                        {"category": "G", "value": 35.0},
                                        {"category": "B", "value": 13.0},
                                        {"category": "A", "value": 7.0},
                                    ],
                                }
                            ],
                        },
                        {
                            "id": "monthly_defaults",
                            "chart_type": "line",
                            "title": "Default rate trend",
                            "finding": "Default rates rose across the observed monthly window.",
                            "x_label": "Month",
                            "y_label": "Default rate (%)",
                            "series": [
                                {
                                    "name": "Default rate",
                                    "data": [
                                        {"category": "Jan", "value": 8.0},
                                        {"category": "Feb", "value": 10.5},
                                        {"category": "Mar", "value": 12.0},
                                    ],
                                }
                            ],
                        },
                    ],
                    "business_findings": ["Loan grade and trend evidence should guide validation."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                    "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            shape_names = "\n".join(
                getattr(shape, "name", "")
                for slide in presentation.slides
                for shape in slide.shapes
            )

            self.assertNotIn("Structured chart data was unavailable", full_text)
            self.assertIn("reconstructed_chart_horizontal_bar", shape_names)
            self.assertIn("reconstructed_chart_line", shape_names)

    def test_skipped_invalid_structured_chart_records_warning_even_with_valid_visual(self) -> None:
        workflow_state = {
            "saved_figures": [],
            "analysis_results": {
                "chart_specs": [
                    {
                        "id": "invalid_chart",
                        "chart_type": "bar",
                        "title": "Invalid chart without data",
                        "finding": "This should be skipped and warned.",
                    },
                    {
                        "id": "valid_chart",
                        "chart_type": "bar",
                        "title": "Valid attrition chart",
                        "finding": "This chart can render.",
                        "data": [{"category": "A", "value": 1}, {"category": "B", "value": 2}],
                    },
                ],
                "business_findings": ["Valid attrition chart can still be used."],
            },
            "agent_outputs": {
                "presentation_architect": {"presentation_title": "Warning Review", "slides": []},
                "data_understander": {"executive_summary": "Dataset", "datasets": {}},
                "market_researcher": {},
                "business_translator": {},
                "decision_maker": {},
            },
        }

        build_deck_spec(workflow_state)
        warnings = workflow_state.get("slide_validation_warnings", [])
        structured_warnings = [warning for warning in warnings if warning.get("issue") == "structured_chart_data_missing"]

        self.assertTrue(structured_warnings)
        self.assertIn("invalid_chart", structured_warnings[0].get("artifact_ids", []))

    def test_grouped_bar_chart_spec_renders_rebuilt_chart(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "grouped_bar.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "analysis_artifacts": [
                        {
                            "artifact_id": "attrition_by_level_and_dimension",
                            "chart_type": "grouped_bar",
                            "title": "Attrition by satisfaction level and dimension",
                            "finding": "The lowest satisfaction group is consistently highest across dimensions.",
                            "x": "level",
                            "y": "attrition_rate",
                            "group_by": "dimension",
                            "data": [
                                {"level": "1", "dimension": "Environment", "attrition_rate": 25.4},
                                {"level": "2", "dimension": "Environment", "attrition_rate": 15.0},
                                {"level": "1", "dimension": "Job", "attrition_rate": 22.8},
                                {"level": "2", "dimension": "Job", "attrition_rate": 16.4},
                            ],
                            "value_format": "{:.1f}%",
                        }
                    ],
                    "analysis_summary": {"max_attrition": "25.4%", "sample": 4},
                    "business_findings": ["The lowest satisfaction group is consistently highest across dimensions."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee attrition dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))

            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 4)

    def test_small_multiples_artifact_renders_without_raw_subplot_image(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "small_multiples.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "analysis_artifacts": [
                        {
                            "artifact_id": "satisfaction_attrition_rates",
                            "artifact_type": "chart_spec",
                            "slide_candidate": True,
                            "finding": "Low satisfaction groups show consistently higher attrition risk.",
                            "chart_type": "small_multiples_bar",
                            "title": "Lower satisfaction is consistently linked to higher attrition",
                            "x_label": "Satisfaction level",
                            "y_label": "Attrition rate (%)",
                            "series": [
                                {
                                    "name": "Environment",
                                    "data": [
                                        {"level": 1, "attrition_rate": 25.4},
                                        {"level": 2, "attrition_rate": 15.0},
                                        {"level": 3, "attrition_rate": 13.7},
                                        {"level": 4, "attrition_rate": 13.5},
                                    ],
                                },
                                {
                                    "name": "Job",
                                    "data": [
                                        {"level": 1, "attrition_rate": 22.8},
                                        {"level": 2, "attrition_rate": 16.4},
                                        {"level": 3, "attrition_rate": 16.5},
                                        {"level": 4, "attrition_rate": 11.3},
                                    ],
                                },
                            ],
                            "takeaway": "Employees with the lowest satisfaction scores show the highest attrition rates.",
                            "recommended_template": "small_multiples_with_takeaway",
                        }
                    ],
                    "analysis_summary": {"low_satisfaction_attrition": "25.4%", "sample": 2},
                    "business_findings": ["Low satisfaction groups show consistently higher attrition risk."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Attrition Review", "slides": []},
                    "data_understander": {"executive_summary": "Employee attrition dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            picture_count = sum(
                1
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertEqual(picture_count, 0)
            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 4)
            self.assertIn("Low satisfaction groups show consistently higher attrition risk", full_text)
            self.assertIn("Environment", full_text)
            self.assertIn("Job", full_text)

    def test_line_and_metric_card_visuals_render_from_structured_data(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "line_metric_cards.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {
                    "analysis_artifacts": [
                        {
                            "artifact_id": "quarterly_retention",
                            "chart_type": "line",
                            "title": "Retention improved across recent quarters",
                            "finding": "Retention improves steadily over the latest periods.",
                            "data": [
                                {"quarter": "Q1", "retention": 82},
                                {"quarter": "Q2", "retention": 85},
                                {"quarter": "Q3", "retention": 88},
                            ],
                            "x": "quarter",
                            "y": "retention",
                            "value_format": "{:.0f}%",
                        },
                        {
                            "artifact_id": "retention_metrics",
                            "chart_type": "metric_cards",
                            "title": "Retention metrics highlight the strongest signal",
                            "finding": "Three metrics summarize the retention opportunity.",
                            "data": [
                                {"label": "Retention", "value": 88},
                                {"label": "Attrition", "value": 12},
                                {"label": "At-risk segment", "value": 24},
                            ],
                            "value_format": "{:.0f}%",
                        },
                    ],
                    "analysis_summary": {"retention": "88%", "attrition": "12%"},
                    "business_findings": ["Retention improves steadily over the latest periods."],
                },
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "Retention Review", "slides": []},
                    "data_understander": {"executive_summary": "Retention dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))

            self.assertGreaterEqual(reconstructed_chart_shape_count(presentation), 6)

    def test_missing_outputs_still_generate_partial_deck(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "partial.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {},
                "agent_outputs": {"presentation_architect": {"slides": []}},
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertEqual(len(presentation.slides), 12)
            self.assertIn("Domain context", full_text)
            self.assertIn("Business translation was not captured", full_text)

    def test_recommendations_and_limitations_render_from_agent_outputs(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "recommendations_limits.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_artifact_warnings": ["Analysis produced fewer than 3 saved figures."],
                "analysis_results": {
                    "analysis_summary": {
                        "default_rate": "12%",
                        "data_quality_note": "Income values include missing observations that may affect segmentation.",
                    },
                    "business_findings": ["Default risk is concentrated in high-utilization customers."],
                },
                "agent_outputs": {
                    "data_understander": {"executive_summary": "Credit risk dataset", "datasets": {}},
                    "market_researcher": {},
                    "business_translator": {
                        "risks": ["Observed relationships are associative, not causal."],
                        "immediate_actions": ["Monitor high-utilization borrowers weekly."],
                    },
                    "decision_maker": {
                        "final_recommendation": "Tighten review for high-utilization borrowers.",
                        "recommendations": [
                            {
                                "rank": 1,
                                "action": "Tighten review for high-utilization borrowers",
                                "rationale": "This segment shows the clearest risk concentration",
                                "evidence": "Default rate reaches 12%",
                                "impact": "High",
                            }
                        ],
                        "limitations": [
                            {
                                "limitation": "Model uses historical associations only",
                                "mitigation": "Validate with a fresh holdout sample",
                                "decision_impact": "Use as a triage signal, not an automatic decline rule",
                            }
                        ],
                    },
                    "presentation_architect": {"presentation_title": "Credit Risk Review", "slides": []},
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertIn("Tighten review for high-utilization borrowers", full_text)
            self.assertIn("Model uses historical associations only", full_text)
            self.assertIn("Observed relationships are associative", full_text)

    def test_recommendation_fallback_is_domain_specific_not_generic_boilerplate(self) -> None:
        renderer = PowerPointRenderer(Presentation(), DeckSpec(deck_title="Attrition Review", slides=[]), THEME)
        spec = SlideSpec(
            slide_number=1,
            slide_role="recommendations",
            template="recommendation_priority",
            headline="Prioritized actions to reduce attrition",
            main_message="Review overtime and retention support for employees.",
            content_blocks=[],
        )

        fallback = renderer._recommendation_validation_fallback(spec, "Action: Review overtime policy")

        self.assertIn("Track attrition", fallback)
        self.assertNotIn("focused pilot and explicit success criteria", fallback)

    def test_recommendation_slide_preserves_action_owner_trigger_guardrail(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "recommendation_fields.pptx"
            workflow_state = {
                "saved_figures": [],
                "analysis_results": {"business_findings": ["Price trend supports guarded action."]},
                "agent_outputs": {
                    "presentation_architect": {"presentation_title": "STC Review", "slides": []},
                    "data_understander": {
                        "datasets": {
                            "STC Stock Price History (5).csv": {
                                "shape": [100, 7],
                                "columns": ["Date", "Price", "Vol.", "Change %"],
                            }
                        }
                    },
                    "market_researcher": {},
                    "business_translator": {},
                    "decision_maker": {
                        "recommendations": [
                            {
                                "action": "Use tactical entry only after confirmation",
                                "owner": "investment committee",
                                "trigger": "price closes above resistance with volume support",
                                "target_segment": "portfolio exposure window",
                                "validation_metric": "forward return, drawdown, and liquidity",
                                "caveat": "Risk: stop if drawdown breaches the limit.",
                            }
                        ]
                    },
                },
            }

            generate_slide_deck(workflow_state, str(output_path))
            presentation = Presentation(str(output_path))
            full_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

            self.assertIn("Action:", full_text)
            self.assertIn("Owner: investment committee", full_text)
            self.assertIn("Trigger: price closes above resistance", full_text)
            self.assertIn("Guardrail: forward return", full_text)


if __name__ == "__main__":
    unittest.main()
