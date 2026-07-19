from __future__ import annotations

import argparse
import asyncio
import json
import os
import re
from pathlib import Path
from typing import Any

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client

from analytics_workflow.presentation_backends import inspect_presentation


SLIDES = [
    ("Analytics Decision Brief", "PowerPoint MCP smoke validation"),
    ("Executive Summary", "The workflow turns evidence into clear, decision-oriented actions."),
    ("Analysis", "Validated analysis separates descriptive evidence from model-based evidence."),
    ("Charts", "A native PowerPoint chart demonstrates the preferred MCP visual path."),
    ("Findings", "Findings remain concise, quantified, and linked to the analytical evidence."),
    ("Recommendations", "Prioritize actions by expected impact, feasibility, owner, and guardrail."),
    ("Limitations", "Treat observational patterns cautiously and validate changes before scaling."),
    ("Sources", "Internal analytical workflow outputs and generated chart evidence."),
]

NAVY = "132A44"
TEAL = "168C8C"
GOLD = "D8A94A"
INK = "23384D"
PALE_BLUE = "EAF1F5"
PALE_TEAL = "E4F0EE"
WARM = "F4ECE4"


async def _call(session: ClientSession, name: str, arguments: dict[str, Any]) -> dict[str, Any]:
    for attempt in range(3):
        result = await session.call_tool(name, arguments)
        text = "\n".join(str(getattr(item, "text", item)) for item in result.content)
        try:
            payload = json.loads(text)
        except json.JSONDecodeError as exc:
            raise RuntimeError(f"PowerPoint MCP {name} returned non-JSON output: {text[:500]}") from exc
        if not result.isError and payload.get("success", False):
            return payload
        detail = str(payload.get("errorMessage") or text)
        transient = any(code in detail for code in ("0x800AC472", "RPC_E_CALL_REJECTED", "server busy"))
        if transient and attempt < 2:
            await asyncio.sleep(0.75 * (attempt + 1))
            continue
        raise RuntimeError(f"PowerPoint MCP {name} failed: {detail}")
    raise RuntimeError(f"PowerPoint MCP {name} failed after bounded retries.")


async def _build(output: Path, show: bool) -> None:
    env = dict(os.environ)
    env.setdefault("DOTNET_ROLL_FORWARD", "Major")
    params = StdioServerParameters(command="mcp-ppt", args=[], env=env, cwd=str(output.parent))
    render_dir = output.parent / f"{output.stem}_rendered"
    render_dir.mkdir(parents=True, exist_ok=True)
    async with stdio_client(params) as (read_stream, write_stream):
        async with ClientSession(read_stream, write_stream) as session:
            await session.initialize()
            tools = {tool.name for tool in (await session.list_tools()).tools}
            required = {"file", "slide", "shape", "text", "chart", "export"}
            if missing := sorted(required - tools):
                raise RuntimeError(f"Missing required PowerPoint MCP tools: {', '.join(missing)}")
            created = await _call(
                session,
                "file",
                {"action": "create", "path": str(output), "show": show},
            )
            session_id = created["session_id"]
            try:
                for slide_index, (title, body) in enumerate(SLIDES, start=1):
                    dark_slide = slide_index in {1, 8}
                    background_color = (
                        NAVY
                        if dark_slide
                        else PALE_BLUE
                        if slide_index in {2, 3, 4}
                        else PALE_TEAL
                        if slide_index in {5, 6}
                        else WARM
                    )
                    title_color = "#FFFFFF" if dark_slide else f"#{NAVY}"
                    body_color = "#DCE7EF" if dark_slide else f"#{INK}"
                    await _call(
                        session,
                        "slide",
                        {
                            "action": "create",
                            "session_id": session_id,
                            "position": slide_index,
                            "layout_name": "Blank",
                        },
                    )
                    background_result = await _call(
                        session,
                        "shape",
                        {
                            "action": "add-shape",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "left": 0,
                            "top": 0,
                            "width": 960,
                            "height": 540,
                            "auto_shape_type": 1,
                        },
                    )
                    background_shape = _shape_name(background_result, "Rectangle 1")
                    for action, extra in (
                        ("set-fill", {"color_hex": background_color}),
                        ("set-line", {"color_hex": background_color, "line_width": 0.5}),
                        ("z-order", {"z_order_cmd": 1}),
                    ):
                        await _call(
                            session,
                            "shape",
                            {
                                "action": action,
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": background_shape,
                                **extra,
                            },
                        )
                    title_result = await _call(
                        session,
                        "shape",
                        {
                            "action": "add-textbox",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "left": 54,
                            "top": 36,
                            "width": 852,
                            "height": 58,
                            "text": title,
                        },
                    )
                    title_shape = _shape_name(title_result, "TextBox 1")
                    await _call(
                        session,
                        "text",
                        {
                            "action": "format",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "shape_name": title_shape,
                            "font_name": "Aptos Display",
                            "font_size": 30,
                            "bold": True,
                            "color": title_color,
                        },
                    )
                    accent_result = await _call(
                        session,
                        "shape",
                        {
                            "action": "add-shape",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "left": 36,
                            "top": 40,
                            "width": 7,
                            "height": 48,
                            "auto_shape_type": 1,
                        },
                    )
                    accent_shape = _shape_name(accent_result, "Rectangle 2")
                    await _call(
                        session,
                        "shape",
                        {
                            "action": "set-fill",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "shape_name": accent_shape,
                            "color_hex": GOLD if dark_slide else TEAL,
                        },
                    )
                    if slide_index != 4:
                        card_result = await _call(
                            session,
                            "shape",
                            {
                                "action": "add-shape",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "left": 64,
                                "top": 145,
                                "width": 832,
                                "height": 220,
                                "auto_shape_type": 5,
                            },
                        )
                        card_shape = _shape_name(card_result, "Rounded Rectangle 1")
                        for action, extra in (
                            ("set-fill", {"color_hex": "1E3A57" if dark_slide else "FFFFFF"}),
                            (
                                "set-line",
                                {
                                    "color_hex": "31506C" if dark_slide else "D8E2E8",
                                    "line_width": 1,
                                },
                            ),
                        ):
                            await _call(
                                session,
                                "shape",
                                {
                                    "action": action,
                                    "session_id": session_id,
                                    "slide_index": slide_index,
                                    "shape_name": card_shape,
                                    **extra,
                                },
                            )
                        body_result = await _call(
                            session,
                            "shape",
                            {
                                "action": "add-textbox",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "left": 98,
                                "top": 195,
                                "width": 764,
                                "height": 125,
                                "text": body,
                            },
                        )
                        body_shape = _shape_name(body_result, "TextBox 2")
                        await _call(
                            session,
                            "text",
                            {
                                "action": "format",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": body_shape,
                                "font_name": "Aptos",
                                "font_size": 20,
                                "color": body_color,
                            },
                        )
                    else:
                        await _call(
                            session,
                            "chart",
                            {
                                "action": "create",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "chart_type": 51,
                                "left": 105,
                                "top": 120,
                                "width": 750,
                                "height": 330,
                            },
                        )
                        shapes = await _call(
                            session,
                            "shape",
                            {"action": "list", "session_id": session_id, "slide_index": slide_index},
                        )
                        chart_name = next(
                            str(shape["name"])
                            for shape in shapes.get("shapes", [])
                            if shape.get("hasChart") or shape.get("shapeType") == "Chart"
                        )
                        await _call(
                            session,
                            "chart",
                            {
                                "action": "set-data",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": chart_name,
                                "values": [
                                    ["Stage", "Confidence", "", ""],
                                    ["Data", 82, None, None],
                                    ["Analysis", 89, None, None],
                                    ["Decision", 93, None, None],
                                    ["", None, None, None],
                                ],
                            },
                        )
                        await _call(
                            session,
                            "chart",
                            {
                                "action": "set-title",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": chart_name,
                                "title": "Evidence confidence by workflow stage",
                            },
                        )
                        await _call(
                            session,
                            "chart",
                            {
                                "action": "set-legend",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": chart_name,
                                "visible": False,
                                "position": 0,
                            },
                        )
                        chart_data = await _call(
                            session,
                            "chart",
                            {
                                "action": "read-data",
                                "session_id": session_id,
                                "slide_index": slide_index,
                                "shape_name": chart_name,
                            },
                        )
                        if "Data" not in json.dumps(chart_data):
                            raise RuntimeError("PowerPoint MCP chart data did not persist after set-data.")
                    footer_result = await _call(
                        session,
                        "shape",
                        {
                            "action": "add-textbox",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "left": 858,
                            "top": 500,
                            "width": 48,
                            "height": 22,
                            "text": f"{slide_index:02d}",
                        },
                    )
                    footer_shape = _shape_name(footer_result, "TextBox 4")
                    await _call(
                        session,
                        "text",
                        {
                            "action": "format",
                            "session_id": session_id,
                            "slide_index": slide_index,
                            "shape_name": footer_shape,
                            "font_name": "Aptos",
                            "font_size": 10,
                            "bold": True,
                            "color": "#B8C8D6" if dark_slide else "#668195",
                        },
                    )
                await _call(session, "file", {"action": "save", "session_id": session_id})
                await _call(
                    session,
                    "export",
                    {
                        "action": "all-slides-to-images",
                        "session_id": session_id,
                        "destination_directory": str(render_dir),
                        "width": 1600,
                        "height": 900,
                    },
                )
            finally:
                await _call(
                    session,
                    "file",
                    {"action": "close", "session_id": session_id, "save": True},
                )


def _shape_name(payload: dict[str, Any], fallback: str) -> str:
    for key in ("shapeName", "shape_name", "name"):
        if payload.get(key):
            return str(payload[key])
    match = re.search(r"['\"]([^'\"]+)['\"]", str(payload.get("message", "")))
    return match.group(1) if match else fallback


def main() -> int:
    parser = argparse.ArgumentParser(description="Run a live eight-slide PowerPoint MCP smoke test.")
    parser.add_argument(
        "--output",
        default=".powerpoint-mcp-smoke/powerpoint_mcp_smoke.pptx",
        help="PPTX output path",
    )
    parser.add_argument("--show", action="store_true", help="Show PowerPoint while the smoke test runs")
    args = parser.parse_args()
    output = Path(args.output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    asyncio.run(_build(output, args.show))
    inspection = inspect_presentation(str(output))
    rendered = sorted(str(path) for path in (output.parent / f"{output.stem}_rendered").glob("*"))
    inspection.rendered_files.extend(rendered)
    print(json.dumps(inspection.to_dict(), indent=2))
    expected_titles = [title for title, _ in SLIDES]
    from pptx import Presentation

    actual_titles = [
        next((shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()), "")
        for slide in Presentation(str(output)).slides
    ]
    return 0 if inspection.valid and inspection.slide_count == 8 and actual_titles == expected_titles else 1


if __name__ == "__main__":
    raise SystemExit(main())
