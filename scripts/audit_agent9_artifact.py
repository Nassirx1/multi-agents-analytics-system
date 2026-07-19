from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit a verified Agent 9 run artifact.")
    parser.add_argument("run_directory", type=Path)
    parser.add_argument("--verified-deck", default="analytics_report_mcp_final_v2.pptx")
    args = parser.parse_args()
    run_dir = args.run_directory.resolve()
    main_deck = run_dir / "analytics_report.pptx"
    verified_deck = run_dir / args.verified_deck
    presentation = Presentation(str(main_deck))
    titles: list[str] = []
    picture_counts: list[int] = []
    all_text: list[str] = []
    for slide in presentation.slides:
        text = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        titles.append(text[0] if text else "")
        picture_counts.append(sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE))
        all_text.extend(text)
    manifest = json.loads((run_dir / "run_manifest.json").read_text(encoding="utf-8"))
    result = json.loads((run_dir / "agent9_mcp_rerun_result.json").read_text(encoding="utf-8"))
    sha256 = lambda path: hashlib.sha256(path.read_bytes()).hexdigest()
    raw_json_markers = any(
        marker in value for value in all_text for marker in ('"slide_role"', '"content_blocks"', '"visual"')
    )
    payload = {
        "main_matches_verified": sha256(main_deck) == sha256(verified_deck),
        "slide_count": len(presentation.slides),
        "picture_counts_slides_4_to_7": picture_counts[3:7],
        "titles": titles,
        "raw_json_markers": raw_json_markers,
        "manifest_backend": manifest.get("presentation_backend_used"),
        "rerun_backend": result.get("backend"),
        "render_count": len(result.get("rendered_files", [])),
    }
    print(json.dumps(payload, indent=2))
    accepted = (
        payload["main_matches_verified"]
        and payload["slide_count"] == 12
        and all(count >= 1 for count in payload["picture_counts_slides_4_to_7"])
        and not payload["raw_json_markers"]
        and payload["manifest_backend"] == "powerpoint_mcp"
        and payload["rerun_backend"] == "powerpoint_mcp"
        and payload["render_count"] >= 12
    )
    return 0 if accepted else 1


if __name__ == "__main__":
    raise SystemExit(main())
