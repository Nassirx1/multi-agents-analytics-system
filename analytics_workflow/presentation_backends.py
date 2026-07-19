from __future__ import annotations

import asyncio
import copy
import json
import logging
import os
import re
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import timedelta
from pathlib import Path
from typing import Any

from .slides.deck_spec import ContentBlock, DeckSpec
from .slides.pptx_renderer import render_deck
from .slides.text_refiner import refine_headline, shorten
from .project_skills import load_project_skill_bundle

try:
    from mcp import ClientSession, StdioServerParameters
    from mcp.client.stdio import stdio_client

    MCP_SDK_AVAILABLE = True
except ImportError:  # pragma: no cover - optional integration
    MCP_SDK_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False


SAFE_AGENT9_TOOLS = frozenset(
    {
        "chart",
        "export",
        "file",
        "image",
        "notes",
        "placeholder",
        "shape",
        "shapealign",
        "slide",
        "slidetable",
        "text",
    }
)
MAX_MCP_TOOL_CALLS = 120
HYBRID_SLIDE_OBJECTIVES = (
    "Frame the decision and audience",
    "Summarize the executive answer and data scope",
    "Establish market or business context",
    "Show the EDA overview and primary pattern",
    "Explain distribution or segment differences",
    "Identify the strongest drivers",
    "Present verified decision-tree rules and model evidence",
    "Prioritize the first set of recommendations from the evidence",
    "Sequence recommendation ownership, timing, and measures",
    "State limitations, risks, and unsupported claims",
    "Document sources, conclusion, and the next decision",
    "End with a concise executive close",
)


def enrich_deck_executive_copy(deck: DeckSpec, agent_outputs: dict[str, Any]) -> None:
    """Attach non-duplicative business implications/actions to analytical slides."""
    translator = agent_outputs.get("business_translator", {}) if isinstance(agent_outputs, dict) else {}
    findings = translator.get("key_findings", []) if isinstance(translator, dict) else []
    actions = translator.get("immediate_actions", []) if isinstance(translator, dict) else []
    decision = agent_outputs.get("decision_maker", {}) if isinstance(agent_outputs, dict) else {}
    recommendations = decision.get("recommendations", []) if isinstance(decision, dict) else []
    for slide in deck.slides:
        if slide.slide_role != "analysis" or slide.visual is None:
            continue
        evidence = str(slide.visual.finding or slide.main_message or slide.headline).strip()
        visual_evidence_ids = set(slide.visual.evidence_ids or slide.evidence_ids)
        matched = next(
            (
                record
                for record in findings
                if isinstance(record, dict)
                and visual_evidence_ids
                and visual_evidence_ids.intersection(str(item) for item in record.get("evidence_ids", []) or [])
            ),
            {},
        )
        if not matched and not visual_evidence_ids:
            matched = _best_text_match(evidence, findings, key="finding")
        implication = str(matched.get("business_implication", "")).strip() if matched else ""
        if not implication:
            implication = _artifact_specific_implication(slide.visual.chart_type, evidence)
        action = _best_text_value(f"{evidence} {implication}", actions) if matched else ""
        if not action:
            action = _artifact_specific_decision_use(slide.visual.chart_type, evidence)
        items = []
        if implication:
            items.append(f"Why it matters: {_shorten_human(implication, 145)}")
        if action:
            items.append(f"Decision use: {_shorten_human(action, 130)}")
        slide.content_blocks = [ContentBlock(type="bullets", items=items[:2])]
    recommendation_slides = [slide for slide in deck.slides if slide.slide_role == "recommendations"]
    if recommendation_slides and isinstance(recommendations, list) and recommendations:
        cards = []
        for index, recommendation in enumerate(recommendations[:3], start=1):
            if not isinstance(recommendation, dict):
                continue
            cards.append(
                ContentBlock(
                    type="recommendation_card",
                    title=f"Priority {recommendation.get('rank') or index}",
                    text="Action: " + _shorten_human(str(recommendation.get("action", "")), 145),
                    items=[
                        "Owner: " + _shorten_human(str(recommendation.get("owner") or "named accountable owner"), 70),
                        "Trigger: " + _shorten_human(str(recommendation.get("trigger") or "approved baseline and protocol"), 90),
                        "Guardrail: " + _shorten_human(str(recommendation.get("validation_metric") or recommendation.get("rationale") or "outcome against baseline"), 90),
                    ],
                )
            )
        if cards:
            recommendation_slides[0].content_blocks = cards
    limitation_slide = next((slide for slide in deck.slides if slide.slide_role == "limitations"), None)
    decision_limitations = decision.get("limitations", []) if isinstance(decision, dict) else []
    if limitation_slide is not None and isinstance(decision_limitations, list):
        limitation_items: list[str] = []
        translator_risks = translator.get("risks", []) if isinstance(translator, dict) else []
        for raw_risk in translator_risks[:2]:
            risk = _shorten_human(str(raw_risk), 145)
            lowered = risk.lower()
            if "association" in lowered and ("caus" in lowered or "directly impact" in lowered):
                risk = "Observed associations do not establish cause and effect; treat intervention ideas as hypotheses."
            elif "automated" in lowered and ("false negative" in lowered or "individual-level" in lowered):
                risk = "Automated individual-level screening risks false negatives under severe class imbalance."
            if risk:
                limitation_items.append(risk)
        for item in decision_limitations:
            if not isinstance(item, dict):
                continue
            limitation = _shorten_human(str(item.get("limitation", "")), 145)
            mitigation = _shorten_human(str(item.get("mitigation", "")), 135)
            decision_impact = _shorten_human(str(item.get("decision_impact", "")), 135)
            combined_limit = f"{limitation} {mitigation}".lower()
            if "association" in combined_limit and ("caus" in combined_limit or "directly impact" in combined_limit):
                limitation = "Observed associations do not establish cause and effect; treat intervention ideas as hypotheses."
            if "sampling" in mitigation.lower() or "class imbalance" in combined_limit:
                mitigation = "Improve class balance, validate on fresh data, and require expert review before any pilot."
            if limitation:
                limitation_items.append(limitation)
            if mitigation:
                limitation_items.append(f"Mitigation: {mitigation}")
            elif decision_impact:
                limitation_items.append(f"Decision impact: {decision_impact}")
            if len(limitation_items) >= 4:
                break
        if limitation_items:
            limitation_slide.content_blocks = [ContentBlock(type="limitations", items=limitation_items[:4])]
    ending = next((slide for slide in deck.slides if slide.slide_role in {"summary", "ending"}), None)
    if ending is not None and isinstance(decision, dict):
        conclusion = _shorten_human(str(decision.get("conclusion", "")), 220)
        final_recommendation = _shorten_human(str(decision.get("final_recommendation", "")), 170)
        if conclusion:
            ending.main_message = conclusion
        ending.content_blocks = [
            ContentBlock(
                type="bullets",
                items=[value for value in (final_recommendation, "Approve, refine, or reject the proposed actions with named owners and validation gates.") if value],
            )
        ]


def _best_text_match(text: str, records: Any, *, key: str) -> dict[str, Any]:
    best: dict[str, Any] = {}
    best_score = 0.0
    source_tokens = _meaningful_tokens(text)
    for record in records if isinstance(records, list) else []:
        if not isinstance(record, dict):
            continue
        candidate_tokens = _meaningful_tokens(str(record.get(key, "")))
        if not source_tokens or not candidate_tokens:
            continue
        score = len(source_tokens & candidate_tokens) / max(1, min(len(source_tokens), len(candidate_tokens)))
        if score > best_score:
            best, best_score = record, score
    return best if best_score >= 0.22 else {}


def _best_text_value(text: str, values: Any) -> str:
    source_tokens = _meaningful_tokens(text)
    best = ""
    best_score = 0.0
    for value in values if isinstance(values, list) else []:
        candidate = str(value).strip()
        candidate_tokens = _meaningful_tokens(candidate)
        if not candidate_tokens:
            continue
        score = len(source_tokens & candidate_tokens) / max(1, min(len(source_tokens), len(candidate_tokens)))
        if score > best_score:
            best, best_score = candidate, score
    return best if best_score >= 0.16 else ""


def _meaningful_tokens(text: str) -> set[str]:
    return {
        token
        for token in re.findall(r"[a-z0-9]+", text.lower())
        if len(token) >= 4 and token not in {"shows", "highest", "employees", "attrition", "analysis"}
    }


def _humanize_deck_copy(deck: DeckSpec) -> None:
    for slide in deck.slides:
        slide.headline = _humanize_presentation_text(slide.headline)
        slide.main_message = _humanize_presentation_text(slide.main_message)
        slide.subtitle = _humanize_presentation_text(slide.subtitle)
        for block in slide.content_blocks:
            block.title = _humanize_presentation_text(block.title)
            block.text = _humanize_presentation_text(block.text)
            block.items = [
                clean for clean in (_humanize_presentation_text(item) for item in block.items) if clean
            ]
        if slide.visual is not None:
            slide.visual.title = _humanize_presentation_text(slide.visual.title)
            slide.visual.finding = _humanize_presentation_text(slide.visual.finding)
            slide.visual.takeaway = _humanize_presentation_text(slide.visual.takeaway)


def _humanize_presentation_text(value: Any) -> str:
    if isinstance(value, (dict, list, tuple, set)):
        return ""
    text = str(value or "").strip()
    if not text:
        return ""
    if (text.startswith("{") and text.endswith("}")) or (text.startswith("[") and text.endswith("]")):
        return ""
    text = text.replace("_", " ")
    text = re.sub(r"(?<=[a-z0-9])(?=[A-Z])", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _shorten_human(text: str, max_chars: int) -> str:
    clean = _humanize_presentation_text(text)
    clean = re.sub(r"\s*\(e\.g\.,.*?\)", "", clean, flags=re.IGNORECASE)
    if len(clean) <= max_chars:
        return clean
    prefix = clean[: max_chars + 1]
    sentence_ends = [match.end() for match in re.finditer(r"[.!?;](?:\s|$)", prefix)]
    if sentence_ends and sentence_ends[-1] >= min(60, max_chars // 2):
        return prefix[: sentence_ends[-1]].strip()
    comma = prefix.rfind(", ")
    if comma >= min(70, max_chars // 2):
        return f"{prefix[:comma].rstrip()}.*".replace(".*", ".")
    truncated = prefix.rsplit(" ", 1)[0].rstrip(" ,;:-")
    return f"{truncated}…"


def _unique_slide_text(values: list[Any], comparison_text: str) -> list[str]:
    selected: list[str] = []
    comparison_tokens = _meaningful_tokens(comparison_text)
    for value in values:
        clean = _humanize_presentation_text(value)
        if not clean:
            continue
        tokens = _meaningful_tokens(clean)
        if comparison_tokens and tokens:
            overlap = len(tokens & comparison_tokens) / max(1, min(len(tokens), len(comparison_tokens)))
            if overlap >= 0.78:
                continue
        if any(
            tokens
            and _meaningful_tokens(existing)
            and len(tokens & _meaningful_tokens(existing))
            / max(1, min(len(tokens), len(_meaningful_tokens(existing))))
            >= 0.78
            for existing in selected
        ):
            continue
        selected.append(clean)
    return selected


def _generic_executive_implication(evidence: str) -> str:
    lowered = evidence.lower()
    if "imbalance" in lowered or "distribution" in lowered:
        return "The base rate defines the scale of the issue and prevents majority-class accuracy from overstating model value."
    if "tree" in lowered or "model" in lowered:
        return "The rules identify investigation segments; they should guide human review rather than automated employment decisions."
    return "The observed segment difference identifies where leaders should validate causes before targeting resources."


def _generic_decision_use(evidence: str) -> str:
    lowered = evidence.lower()
    if "imbalance" in lowered or "distribution" in lowered:
        return "Track the base rate as the intervention benchmark and report class-sensitive model metrics."
    if "tree" in lowered or "model" in lowered:
        return "Validate the branches with managers and employee feedback before using them for intervention design."
    return "Confirm the pattern with operational owners, then pilot a focused intervention with a measurable outcome."


def _artifact_specific_implication(chart_type: str, evidence: str) -> str:
    chart = str(chart_type or "").lower()
    if chart == "decision_tree":
        return "The branches define testable segments, not deployment-ready prediction rules."
    if chart in {"heatmap", "correlation", "correlation_heatmap", "scatter"}:
        return "Use direction and magnitude to select follow-up hypotheses; the relationship is not causal evidence."
    if chart in {"distribution", "box", "histogram"}:
        return "The group separation identifies a validation question; sample size and uncertainty determine whether it is stable."
    if chart in {"bar", "column", "horizontal_bar", "ranking", "grouped_bar"}:
        return "The ranking is descriptive; compare counts, denominators, and uncertainty before targeting a segment."
    return _generic_executive_implication(evidence)


def _artifact_specific_decision_use(chart_type: str, evidence: str) -> str:
    chart = str(chart_type or "").lower()
    if chart == "decision_tree":
        return "Validate branches on fresh data with precision, recall, F1, support, fairness, and human review."
    if chart in {"heatmap", "correlation", "correlation_heatmap", "scatter"}:
        return "Prioritize the strongest association for controlled follow-up while monitoring alternative explanations."
    if chart in {"distribution", "box", "histogram"}:
        return "Confirm the difference with an appropriate group test and report effect size before action."
    if chart in {"bar", "column", "horizontal_bar", "ranking", "grouped_bar"}:
        return "Validate the segment comparison statistically and operationally before allocating resources."
    return _generic_decision_use(evidence)


@dataclass(slots=True)
class PresentationInspection:
    path: str
    slide_count: int = 0
    issues: list[dict[str, Any]] = field(default_factory=list)
    rendered_files: list[str] = field(default_factory=list)

    @property
    def valid(self) -> bool:
        return self.slide_count > 0 and not any(
            issue.get("severity", "error") == "error" for issue in self.issues
        )

    def to_dict(self) -> dict[str, Any]:
        return {
            "path": self.path,
            "slide_count": self.slide_count,
            "issues": self.issues,
            "rendered_files": self.rendered_files,
            "valid": self.valid,
        }


class PresentationBackend(ABC):
    name = "base"

    @abstractmethod
    def create_presentation(self, output_path: str) -> Any: ...

    @abstractmethod
    def add_slide(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def add_title(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def add_text(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def add_chart(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def add_table(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def add_image(self, presentation: Any, **kwargs: Any) -> Any: ...

    @abstractmethod
    def save(self, presentation: Any, output_path: str) -> str: ...

    def inspect(self, output_path: str) -> PresentationInspection:
        return inspect_presentation(output_path)

    @abstractmethod
    def render(self, deck: DeckSpec, output_path: str, **kwargs: Any) -> str: ...


class PythonPresentationBackend(PresentationBackend):
    name = "python"

    def create_presentation(self, output_path: str) -> Any:
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed.")
        return Presentation()

    def add_slide(self, presentation: Any, **kwargs: Any) -> Any:
        return presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_title(self, presentation: Any, **kwargs: Any) -> Any:
        return self.add_text(presentation, **kwargs)

    def add_text(self, presentation: Any, **kwargs: Any) -> Any:
        from pptx.util import Inches

        slide = kwargs["slide"]
        box = slide.shapes.add_textbox(
            Inches(float(kwargs.get("left", 0.7))),
            Inches(float(kwargs.get("top", 0.7))),
            Inches(float(kwargs.get("width", 11.8))),
            Inches(float(kwargs.get("height", 1.0))),
        )
        box.text = str(kwargs.get("text", ""))
        return box

    def add_chart(self, presentation: Any, **kwargs: Any) -> Any:
        raise NotImplementedError("Use render() with a structured DeckSpec for Python charts.")

    def add_table(self, presentation: Any, **kwargs: Any) -> Any:
        slide = kwargs["slide"]
        from pptx.util import Inches

        return slide.shapes.add_table(
            int(kwargs["rows"]),
            int(kwargs["columns"]),
            Inches(float(kwargs.get("left", 0.7))),
            Inches(float(kwargs.get("top", 1.8))),
            Inches(float(kwargs.get("width", 11.8))),
            Inches(float(kwargs.get("height", 4.8))),
        ).table

    def add_image(self, presentation: Any, **kwargs: Any) -> Any:
        slide = kwargs["slide"]
        from pptx.util import Inches

        return slide.shapes.add_picture(
            str(kwargs["image_path"]),
            Inches(float(kwargs.get("left", 0.7))),
            Inches(float(kwargs.get("top", 1.7))),
            width=Inches(float(kwargs.get("width", 7.5))),
        )

    def save(self, presentation: Any, output_path: str) -> str:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return str(Path(output_path).resolve())

    def render(self, deck: DeckSpec, output_path: str, **kwargs: Any) -> str:
        return render_deck(deck, output_path)


class PowerPointMCPBackend(PresentationBackend):
    """Agent-9-only PowerPoint MCP backend with strict tool and path allow-lists."""

    name = "powerpoint_mcp"

    def __init__(
        self,
        openrouter_client: Any,
        *,
        command: str = "mcp-ppt",
        timeout_seconds: int = 900,
        request_timeout_seconds: int = 180,
        enabled_tools: frozenset[str] = SAFE_AGENT9_TOOLS,
    ) -> None:
        self.openrouter_client = openrouter_client
        self.command = command
        self.timeout_seconds = max(10, int(timeout_seconds))
        self.request_timeout_seconds = max(1, int(request_timeout_seconds))
        self.enabled_tools = enabled_tools
        self._logger = logging.getLogger("PowerPointMCPBackend")
        self._active_session_id: str | None = None
        self.last_rendered_files: list[str] = []

    @property
    def available(self) -> bool:
        return MCP_SDK_AVAILABLE and os.name == "nt"

    def create_presentation(self, output_path: str) -> dict[str, Any]:
        return {"output_path": str(Path(output_path).resolve())}

    def add_slide(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "slide", "arguments": {"action": "create", **kwargs}}

    def add_title(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "text", "arguments": {"action": "set", **kwargs}}

    def add_text(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "text", "arguments": {"action": "set", **kwargs}}

    def add_chart(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "chart", "arguments": {"action": "create", **kwargs}}

    def add_table(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "slidetable", "arguments": {"action": "create", **kwargs}}

    def add_image(self, presentation: Any, **kwargs: Any) -> dict[str, Any]:
        return {"tool": "image", "arguments": {"action": "insert", **kwargs}}

    def save(self, presentation: Any, output_path: str) -> str:
        path = Path(output_path).resolve()
        if not path.exists():
            raise RuntimeError(f"PowerPoint MCP did not save the requested presentation: {path}")
        return str(path)

    def render(
        self,
        deck: DeckSpec,
        output_path: str,
        *,
        workflow_state: dict[str, Any] | None = None,
    ) -> str:
        if not self.available:
            raise RuntimeError("PowerPoint MCP backend is unavailable on this machine.")
        output = Path(output_path).resolve()
        output.parent.mkdir(parents=True, exist_ok=True)
        if output.exists():
            try:
                output.unlink()
            except OSError as exc:
                raise RuntimeError(f"Cannot replace the existing PowerPoint deck at {output}: {exc}") from exc
        allowed_inputs = {
            str(Path(path).resolve())
            for path in (workflow_state or {}).get("saved_figures", []) or []
            if Path(path).is_file()
        }
        asyncio.run(self._render_async(deck, output, allowed_inputs))
        return self.save(None, str(output))

    async def _render_async(self, deck: DeckSpec, output: Path, allowed_inputs: set[str]) -> None:
        env = dict(os.environ)
        env.setdefault("DOTNET_ROLL_FORWARD", "Major")
        params = StdioServerParameters(command=self.command, args=[], env=env, cwd=str(output.parent))
        async with stdio_client(params) as (read_stream, write_stream):
            async with ClientSession(
                read_stream,
                write_stream,
                read_timeout_seconds=timedelta(seconds=self.request_timeout_seconds),
            ) as session:
                await session.initialize()
                listed = await session.list_tools()
                tools_by_name = {tool.name: tool for tool in listed.tools if tool.name in self.enabled_tools}
                required = {"file", "slide", "shape", "text", "export"}
                missing = sorted(required - tools_by_name.keys())
                if missing:
                    raise RuntimeError(f"PowerPoint MCP is missing required tools: {', '.join(missing)}")
                await self._agent_tool_loop(session, tools_by_name, deck, output, allowed_inputs)

    async def _agent_tool_loop(
        self,
        session: Any,
        tools_by_name: dict[str, Any],
        deck: DeckSpec,
        output: Path,
        allowed_inputs: set[str],
    ) -> None:
        _apply_hybrid_story_contract(deck)
        _humanize_deck_copy(deck)
        deadline = time.monotonic() + self.timeout_seconds
        self._polish_slide_copy_batches(deck, deadline)
        await self._create_scaffold(session, deck, output, deadline, allowed_inputs)
        inspection = inspect_presentation(str(output), expected_deck=deck)
        if not inspection.valid:
            raise RuntimeError(
                "PowerPoint MCP deck failed deterministic QA: "
                + "; ".join(issue["issue"] for issue in inspection.issues[:10])
            )
        await self._render_for_visual_qa(session, output, deadline)

    def _polish_slide_copy_batches(self, deck: DeckSpec, deadline: float) -> None:
        """Apply model output to two-slide batches instead of discarding read-only audit responses."""
        prior = deck.metadata.get("agent9_copy_polish", {}) if isinstance(deck.metadata, dict) else {}
        if isinstance(prior, dict) and int(prior.get("slides_edited", 0) or 0) >= len(deck.slides):
            self._logger.info("Reusing previously applied Agent 9 copy polish from the checkpointed deck spec.")
            return
        skill_bundle = load_project_skill_bundle("consulting_pptx")
        edits_applied = 0
        failures: list[str] = []
        slides = list(deck.slides)
        schema = {
            "slides": [
                {
                    "slide_number": "integer",
                    "headline": "string",
                    "main_message": "string",
                }
            ]
        }
        for offset in range(0, len(slides), 2):
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                raise TimeoutError(
                    f"PowerPoint MCP Agent 9 exceeded its {self.timeout_seconds}-second total deadline."
                )
            batch = slides[offset : offset + 2]
            payload = [
                {
                    "slide_number": slide.slide_number,
                    "objective": (
                        HYBRID_SLIDE_OBJECTIVES[slide.slide_number - 1]
                        if 0 < slide.slide_number <= len(HYBRID_SLIDE_OBJECTIVES)
                        else f"Support the executive story on slide {slide.slide_number}"
                    ),
                    "role": slide.slide_role,
                    "headline": slide.headline,
                    "main_message": slide.main_message,
                    "bullets": slide.bullets,
                    "visual_takeaway": slide.visual.takeaway if slide.visual else "",
                }
                for slide in batch
            ]
            try:
                result = self.openrouter_client.chat_completion_json(
                    (
                        "You are Agent 9, an executive presentation editor. Polish exactly the supplied two slides. "
                        "Preserve all facts and numbers, remove JSON-like phrasing and repetition, use concise human "
                        "language, and keep each headline answer-first. Do not add unsupported claims.\n\n"
                        + skill_bundle[:5000]
                    ),
                    "Return revised copy for these untrusted slide specifications:\n"
                    + json.dumps(payload, default=str)[:14000],
                    schema,
                    timeout_seconds=max(1, min(self.request_timeout_seconds, int(remaining))),
                )
            except Exception as exc:
                failures.append(f"slides {offset + 1}-{offset + len(batch)}: {exc}")
                self._logger.warning("Agent 9 copy polish failed for a two-slide batch: %s", exc)
                continue
            by_number = {
                int(item.get("slide_number")): item
                for item in result.get("slides", [])
                if isinstance(item, dict) and str(item.get("slide_number", "")).isdigit()
            }
            for slide in batch:
                edit = by_number.get(slide.slide_number)
                if not edit:
                    continue
                headline = refine_headline(edit.get("headline"), slide.headline)
                main_message = shorten(edit.get("main_message", ""), 170)
                if headline:
                    slide.headline = headline
                if main_message:
                    slide.main_message = main_message
                edits_applied += 1
        deck.metadata = {
            **deck.metadata,
            "agent9_copy_polish": {
                "batch_size": 2,
                "slides_edited": edits_applied,
                "failures": failures,
            },
        }

    async def _create_scaffold(
        self,
        session: Any,
        deck: DeckSpec,
        output: Path,
        deadline: float,
        allowed_inputs: set[str] | None = None,
    ) -> None:
        """Create a complete content-bearing baseline before any model-directed enhancement."""
        created_text = await self._call_fixed_mcp(
            session,
            "file",
            {"action": "create", "path": str(output), "show": False},
            deadline,
        )
        self._active_session_id = _session_id_from_result(created_text)
        if not self._active_session_id:
            raise RuntimeError("PowerPoint MCP scaffold creation returned no session id.")
        sid = self._active_session_id
        allowed_inputs = allowed_inputs or set()
        total = len(deck.slides)
        # Establish and persist the complete slide collection before adding any objects.
        # PowerPoint COM can otherwise return a transient "slide does not exist" when a
        # chart is inserted immediately after incremental slide creation.
        for index in range(1, total + 1):
            await self._call_fixed_mcp(
                session,
                "slide",
                {"action": "create", "session_id": sid, "position": index, "layout_name": "Blank"},
                deadline,
            )
        await self._call_fixed_mcp(session, "file", {"action": "save", "session_id": sid}, deadline)
        for index, slide in enumerate(deck.slides, start=1):
            dark = index in {1, total}
            background = "082B4C" if dark else ("EAF3F8" if index % 2 == 0 else "F7F3EA")
            title_color = "#FFFFFF" if dark else "#102A43"
            body_color = "#DCE7EF" if dark else "#243B53"
            background_text = await self._call_fixed_mcp(
                session,
                "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": index,
                    "left": 0, "top": 0, "width": 960, "height": 540, "auto_shape_type": 1,
                },
                deadline,
            )
            background_name = _shape_name_from_result(background_text, "Rectangle 1")
            for action, extra in (
                ("set-fill", {"color_hex": background}),
                ("set-line", {"color_hex": background, "line_width": 0.5}),
                ("z-order", {"z_order_cmd": 1}),
            ):
                await self._call_fixed_mcp(
                    session,
                    "shape",
                    {"action": action, "session_id": sid, "slide_index": index, "shape_name": background_name, **extra},
                    deadline,
                )
            if index == 1:
                await self._add_cover_scaffold(session, sid, deck, slide, deadline)
                continue
            if slide.slide_role == "recommendations":
                await self._add_recommendations_scaffold(session, sid, slide, index, deadline)
                await self._add_evidence_footer(session, sid, slide, index, dark, deadline)
                if index % 2 == 0:
                    await self._call_fixed_mcp(session, "file", {"action": "save", "session_id": sid}, deadline)
                continue
            if slide.slide_role == "limitations":
                await self._add_limitations_scaffold(session, sid, slide, index, deadline)
                await self._add_evidence_footer(session, sid, slide, index, dark, deadline)
                continue
            if slide.slide_role == "ending":
                await self._add_ending_scaffold(session, sid, slide, total, deadline)
                await self._add_evidence_footer(session, sid, slide, index, dark, deadline)
                if index % 2 == 0:
                    await self._call_fixed_mcp(session, "file", {"action": "save", "session_id": sid}, deadline)
                continue
            title_text = await self._call_fixed_mcp(
                session,
                "shape",
                {
                    "action": "add-textbox", "session_id": sid, "slide_index": index,
                    "left": 58, "top": 36, "width": 840, "height": 72, "text": slide.headline,
                },
                deadline,
            )
            title_name = _shape_name_from_result(title_text, "TextBox 1")
            await self._call_fixed_mcp(
                session,
                "text",
                {
                    "action": "format", "session_id": sid, "slide_index": index, "shape_name": title_name,
                    "font_name": "Aptos Display",
                    "font_size": 24 if len(slide.headline) > 72 else 28,
                    "bold": True, "color": title_color,
                },
                deadline,
            )
            bullets = _unique_slide_text([item for block in slide.content_blocks for item in block.items], slide.main_message)
            if slide.slide_role == "analysis" and slide.visual is not None:
                body_lines = list(bullets[:2])
                caveat = str(slide.visual.takeaway or slide.main_message).strip()
                if slide.visual.chart_type == "decision_tree" and caveat:
                    body_lines.append(f"Caveat: {_shorten_human(caveat, 115)}")
                if not body_lines:
                    body_lines = [f"Why it matters: {_generic_executive_implication(slide.main_message)}"]
            else:
                body_lines = _unique_slide_text(
                    [text for text in (slide.subtitle, slide.main_message) if text] + bullets,
                    "",
                )
                if slide.slide_role == "data_understanding" and slide.metrics:
                    metric_lines = [
                        f"{_humanize_presentation_text(metric.get('label'))}: {_humanize_presentation_text(metric.get('value'))}"
                        for metric in slide.metrics[:4]
                        if isinstance(metric, dict) and metric.get("value") not in (None, "")
                    ]
                    body_lines = metric_lines + body_lines[:3]
            body = "\n".join(body_lines) or slide.headline
            image_path = ""
            if slide.visual is not None and slide.visual.image_path:
                candidate = str(Path(slide.visual.image_path).resolve())
                if candidate in allowed_inputs and Path(candidate).is_file():
                    image_path = candidate
            has_chart = bool(
                slide.visual is not None
                and slide.visual.type == "structured_chart"
                and slide.visual.data
                and slide.visual.chart_type != "decision_tree"
                and not image_path
            )
            has_decision_tree = bool(
                slide.visual is not None
                and slide.visual.type == "structured_chart"
                and slide.visual.chart_type == "decision_tree"
                and not image_path
            )
            has_right_visual = bool(image_path) or has_chart or has_decision_tree
            panel_text = await self._call_fixed_mcp(
                session,
                "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": index,
                    "left": 58, "top": 132, "width": 844, "height": 300,
                    "auto_shape_type": 5,
                },
                deadline,
            )
            panel_name = _shape_name_from_result(panel_text, "Rounded Rectangle 1")
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-fill", "session_id": sid, "slide_index": index, "shape_name": panel_name,
                    "color_hex": "163B5C" if dark else "FFFFFF",
                }, deadline,
            )
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-line", "session_id": sid, "slide_index": index, "shape_name": panel_name,
                    "color_hex": "31506C" if dark else "D7E2E8", "line_width": 1,
                }, deadline,
            )
            body_text = await self._call_fixed_mcp(
                session,
                "shape",
                {
                    "action": "add-textbox", "session_id": sid, "slide_index": index,
                    "left": 88, "top": 166, "width": 292 if has_right_visual else 784, "height": 228, "text": body,
                },
                deadline,
            )
            body_name = _shape_name_from_result(body_text, "TextBox 2")
            await self._call_fixed_mcp(
                session,
                "text",
                {
                    "action": "format", "session_id": sid, "slide_index": index, "shape_name": body_name,
                    "font_name": "Aptos",
                    "font_size": 14 if has_decision_tree or (image_path and slide.visual and slide.visual.chart_type == "decision_tree") else 15 if slide.slide_role == "analysis" else 18,
                    "color": body_color,
                },
                deadline,
            )
            if has_chart:
                await self._add_scaffold_chart(session, sid, index, slide.visual, deadline)
            elif has_decision_tree:
                await self._add_scaffold_decision_tree(session, sid, index, slide.visual, deadline)
            elif image_path:
                await self._call_fixed_mcp(
                    session,
                    "image",
                    {
                        "action": "insert", "session_id": sid, "slide_index": index,
                        "image_path": image_path, "left": 402, "top": 144,
                        "width": 482, "height": 278,
                    },
                    deadline,
                )
            elif index in {4, 5, 6}:
                raise RuntimeError(
                    f"Slide {index} requires its assigned saved EDA figure, but no approved image was available."
                )
            evidence_footer = ", ".join(slide.evidence_ids[:3] or slide.source_ids[:2])
            if evidence_footer:
                await self._add_formatted_text(
                    session,
                    sid,
                    index,
                    f"Evidence: {evidence_footer}",
                    58,
                    500,
                    690,
                    18,
                    9,
                    False,
                    "#B8C8D6" if dark else "#607D8B",
                    deadline,
                )
            footer_text = await self._call_fixed_mcp(
                session,
                "shape",
                {
                    "action": "add-textbox", "session_id": sid, "slide_index": index,
                    "left": 856, "top": 500, "width": 46, "height": 20, "text": f"{index:02d}",
                },
                deadline,
            )
            footer_name = _shape_name_from_result(footer_text, "TextBox 3")
            await self._call_fixed_mcp(
                session, "text",
                {
                    "action": "format", "session_id": sid, "slide_index": index, "shape_name": footer_name,
                    "font_name": "Aptos", "font_size": 10, "bold": True,
                    "color": "#B8C8D6" if dark else "#607D8B",
                }, deadline,
            )
            if index % 2 == 0:
                await self._call_fixed_mcp(
                    session, "file", {"action": "save", "session_id": sid}, deadline
                )
        await self._call_fixed_mcp(session, "file", {"action": "save", "session_id": sid}, deadline)

    async def _add_cover_scaffold(
        self,
        session: Any,
        sid: str,
        deck: DeckSpec,
        slide: Any,
        deadline: float,
    ) -> None:
        """Render the established consulting cover used by the deterministic deck."""
        for left, width, color in ((0, 32, "C49A3A"), (710, 250, "356FAE")):
            shape_text = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": 1,
                    "left": left, "top": 0, "width": width, "height": 540, "auto_shape_type": 1,
                }, deadline,
            )
            shape_name = _shape_name_from_result(shape_text, "Rectangle 2")
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-fill", "session_id": sid, "slide_index": 1,
                    "shape_name": shape_name, "color_hex": color,
                }, deadline,
            )
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-line", "session_id": sid, "slide_index": 1,
                    "shape_name": shape_name, "color_hex": color, "line_width": 0.5,
                }, deadline,
            )
        accent_text = await self._call_fixed_mcp(
            session, "shape",
            {
                "action": "add-shape", "session_id": sid, "slide_index": 1,
                "left": 62, "top": 166, "width": 88, "height": 4, "auto_shape_type": 1,
            }, deadline,
        )
        accent_name = _shape_name_from_result(accent_text, "Rectangle 4")
        await self._call_fixed_mcp(
            session, "shape",
            {
                "action": "set-fill", "session_id": sid, "slide_index": 1,
                "shape_name": accent_name, "color_hex": "C49A3A",
            }, deadline,
        )
        description = str(deck.dataset_context.get("description") or slide.subtitle or deck.subtitle).strip()
        date_text = str(deck.metadata.get("generated_on") or "Executive analytics review")
        text_specs = [
            ("MULTI-AGENT ANALYTICS SYSTEM", 62, 48, 560, 24, 9, True, "#FFFFFF"),
            (slide.headline, 62, 188, 580, 92, 27, True, "#FFFFFF"),
            (_shorten_human(description, 170), 62, 316, 570, 62, 12, False, "#E7EEF5"),
            (date_text, 62, 468, 240, 24, 9, False, "#E7EEF5"),
            ("EXECUTIVE\nDECISION DECK", 752, 196, 168, 72, 17, True, "#FFFFFF"),
        ]
        for text_value, left, top, width, height, font_size, bold, color in text_specs:
            await self._add_formatted_text(
                session, sid, 1, text_value, left, top, width, height, font_size, bold, color, deadline
            )

    async def _add_recommendations_scaffold(
        self, session: Any, sid: str, slide: Any, slide_index: int, deadline: float
    ) -> None:
        await self._add_formatted_text(
            session, sid, slide_index, _shorten_human(slide.headline, 82), 56, 28, 850, 60,
            22 if len(slide.headline) > 64 else 25, True, "#102A43", deadline
        )
        callout = _shorten_human(slide.main_message or "Prioritize the actions, validate the evidence, and assign ownership.", 155)
        callout_shape = await self._call_fixed_mcp(
            session, "shape",
            {
                "action": "add-shape", "session_id": sid, "slide_index": slide_index,
                "left": 56, "top": 94, "width": 848, "height": 62, "auto_shape_type": 5,
            }, deadline,
        )
        callout_name = _shape_name_from_result(callout_shape, "Rounded Rectangle 2")
        await self._call_fixed_mcp(
            session, "shape",
            {"action": "set-fill", "session_id": sid, "slide_index": slide_index, "shape_name": callout_name, "color_hex": "E6F0FA"},
            deadline,
        )
        await self._add_formatted_text(
            session, sid, slide_index, callout, 72, 103, 812, 44,
            10 if len(callout) > 110 else 11, True, "#17324D", deadline
        )
        cards = list(slide.content_blocks[:3])
        if len(cards) < 3:
            raw_items = [item for block in slide.content_blocks for item in block.items]
            cards = [
                ContentBlock(type="recommendation_card", title=f"Priority {idx + 1}", text=item, items=["Validation required"])
                for idx, item in enumerate(raw_items[:3])
            ]
        while len(cards) < 3:
            cards.append(ContentBlock(type="recommendation_card", title=f"Priority {len(cards) + 1}", text="Confirm the next evidence-backed action.", items=["Validation required"]))
        accents = ("356FAE", "2F8A7E", "C49A3A")
        for card_index, card in enumerate(cards[:3]):
            left = 64 + card_index * 292
            card_shape = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": slide_index,
                    "left": left, "top": 178, "width": 256, "height": 276, "auto_shape_type": 5,
                }, deadline,
            )
            card_name = _shape_name_from_result(card_shape, f"Rounded Rectangle {card_index + 3}")
            await self._call_fixed_mcp(
                session, "shape",
                {"action": "set-fill", "session_id": sid, "slide_index": slide_index, "shape_name": card_name, "color_hex": "FFFFFF"}, deadline,
            )
            accent_shape = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": slide_index,
                    "left": left, "top": 170, "width": 256, "height": 8, "auto_shape_type": 1,
                }, deadline,
            )
            accent_name = _shape_name_from_result(accent_shape, f"Rectangle {card_index + 5}")
            await self._call_fixed_mcp(
                session, "shape",
                {"action": "set-fill", "session_id": sid, "slide_index": slide_index, "shape_name": accent_name, "color_hex": accents[card_index]}, deadline,
            )
            rationale = card.items[0] if card.items else "Validate the rationale before scaling."
            status = card.items[1] if len(card.items) > 1 else "Validation required"
            if _meaningful_tokens(rationale) == _meaningful_tokens(status):
                rationale = "Assign an owner, baseline, and success measure before advancing."
            await self._add_formatted_text(session, sid, slide_index, card.title or f"Priority {card_index + 1}", left + 20, 198, 216, 24, 12, True, "#356FAE", deadline)
            await self._add_formatted_text(session, sid, slide_index, _shorten_human(card.text, 105), left + 20, 238, 216, 84, 13, True, "#102A43", deadline)
            await self._add_formatted_text(session, sid, slide_index, _shorten_human(rationale, 90), left + 20, 330, 216, 62, 12, False, "#334E68", deadline)
            await self._add_formatted_text(
                session, sid, slide_index, _shorten_human(status, 55), left + 20, 406, 216, 38,
                10 if len(status) > 36 else 11, False, "#607D8B", deadline
            )

    async def _add_limitations_scaffold(
        self, session: Any, sid: str, slide: Any, slide_index: int, deadline: float
    ) -> None:
        await self._add_formatted_text(
            session, sid, slide_index, _shorten_human(slide.headline, 82), 46, 28, 860, 60,
            22 if len(slide.headline) > 64 else 25, True, "#102A43", deadline
        )
        await self._add_formatted_text(
            session, sid, slide_index,
            _shorten_human(slide.main_message or "Separate supported evidence from areas that require validation.", 150),
            56, 104, 848, 36, 11, True, "#17324D", deadline,
        )
        limitations = [item for block in slide.content_blocks for item in block.items][:4]
        while len(limitations) < 4:
            limitations.append("Validate this evidence boundary before making a broader decision.")
        for item_index, item in enumerate(limitations[:4]):
            left = 56 if item_index % 2 == 0 else 510
            top = 174 if item_index < 2 else 320
            card_shape = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": slide_index,
                    "left": left, "top": top, "width": 394, "height": 112, "auto_shape_type": 5,
                }, deadline,
            )
            card_name = _shape_name_from_result(card_shape, f"Rounded Rectangle {item_index + 2}")
            await self._call_fixed_mcp(
                session, "shape",
                {"action": "set-fill", "session_id": sid, "slide_index": slide_index, "shape_name": card_name, "color_hex": "FFFFFF"}, deadline,
            )
            card_label = "Mitigation" if str(item).lstrip().lower().startswith("mitigation:") else "Limitation"
            await self._add_formatted_text(session, sid, slide_index, card_label, left + 20, top + 14, 120, 22, 12, True, "#C5553D", deadline)
            await self._add_formatted_text(session, sid, slide_index, _shorten_human(item, 145), left + 20, top + 44, 354, 54, 12, False, "#243B53", deadline)

    async def _add_ending_scaffold(
        self, session: Any, sid: str, slide: Any, total: int, deadline: float
    ) -> None:
        rail = await self._call_fixed_mcp(
            session, "shape",
            {"action": "add-shape", "session_id": sid, "slide_index": total, "left": 0, "top": 0, "width": 28, "height": 540, "auto_shape_type": 1}, deadline,
        )
        rail_name = _shape_name_from_result(rail, "Rectangle 2")
        await self._call_fixed_mcp(session, "shape", {"action": "set-fill", "session_id": sid, "slide_index": total, "shape_name": rail_name, "color_hex": "C49A3A"}, deadline)
        await self._add_formatted_text(
            session, sid, total, _shorten_human(slide.headline, 88), 58, 86, 840, 96,
            25 if len(slide.headline) > 58 else 28, True, "#FFFFFF", deadline
        )
        rule = await self._call_fixed_mcp(
            session, "shape",
            {"action": "add-shape", "session_id": sid, "slide_index": total, "left": 58, "top": 194, "width": 820, "height": 4, "auto_shape_type": 1}, deadline,
        )
        rule_name = _shape_name_from_result(rule, "Rectangle 3")
        await self._call_fixed_mcp(session, "shape", {"action": "set-fill", "session_id": sid, "slide_index": total, "shape_name": rule_name, "color_hex": "C49A3A"}, deadline)
        await self._add_formatted_text(session, sid, total, _shorten_human(slide.main_message, 200), 58, 222, 820, 78, 15, False, "#E7EEF5", deadline)
        bullets = [item for block in slide.content_blocks for item in block.items][:3]
        bullet_text = "\n".join(f"- {_shorten_human(item, 125)}" for item in bullets)
        await self._add_formatted_text(session, sid, total, bullet_text, 68, 326, 800, 100, 12, False, "#FFFFFF", deadline)
        await self._add_formatted_text(session, sid, total, "Generated by the Multi-Agent Analytics System", 58, 490, 400, 18, 8, False, "#D8E4EE", deadline)
        await self._add_formatted_text(session, sid, total, f"{total} / {total}", 850, 490, 58, 18, 8, False, "#D8E4EE", deadline)

    async def _add_formatted_text(
        self, session: Any, sid: str, slide_index: int, text_value: str,
        left: float, top: float, width: float, height: float,
        font_size: float, bold: bool, color: str, deadline: float,
    ) -> str:
        textbox = await self._call_fixed_mcp(
            session, "shape",
            {
                "action": "add-textbox", "session_id": sid, "slide_index": slide_index,
                "left": left, "top": top, "width": width, "height": height, "text": text_value,
            }, deadline,
        )
        name = _shape_name_from_result(textbox, "TextBox 1")
        await self._call_fixed_mcp(
            session, "text",
            {
                "action": "format", "session_id": sid, "slide_index": slide_index,
                "shape_name": name, "font_name": "Aptos Display" if bold else "Aptos",
                "font_size": font_size, "bold": bold, "color": color,
            }, deadline,
        )
        # PowerPoint defaults new text boxes to ShapeToFitText, which silently
        # collapses the requested height and later clips wrapped executive copy.
        # Disable shape auto-sizing and then reassert the deterministic geometry.
        await self._call_fixed_mcp(
            session,
            "shape",
            {
                "action": "set-text-frame",
                "session_id": sid,
                "slide_index": slide_index,
                "shape_name": name,
                "word_wrap": True,
                "auto_size": 0,
                "margin_left": 4,
                "margin_right": 4,
                "margin_top": 2,
                "margin_bottom": 2,
            },
            deadline,
        )
        await self._call_fixed_mcp(
            session,
            "shape",
            {
                "action": "move-resize",
                "session_id": sid,
                "slide_index": slide_index,
                "shape_name": name,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            },
            deadline,
        )
        return name

    async def _add_evidence_footer(
        self,
        session: Any,
        sid: str,
        slide: Any,
        slide_index: int,
        dark: bool,
        deadline: float,
    ) -> None:
        evidence = ", ".join(slide.evidence_ids[:3] or slide.source_ids[:2])
        if not evidence:
            return
        top = 458 if slide.slide_role == "ending" else 500
        await self._add_formatted_text(
            session,
            sid,
            slide_index,
            f"Evidence: {evidence}",
            58,
            top,
            690,
            18,
            9,
            False,
            "#B8C8D6" if dark else "#607D8B",
            deadline,
        )

    async def _add_scaffold_chart(
        self, session: Any, sid: str, slide_index: int, visual: Any, deadline: float
    ) -> None:
        await self._call_fixed_mcp(
            session, "chart",
            {
                "action": "create", "session_id": sid, "slide_index": slide_index,
                "chart_type": 51, "left": 470, "top": 180, "width": 390, "height": 220,
            }, deadline,
        )
        listed = await self._call_fixed_mcp(
            session, "shape", {"action": "list", "session_id": sid, "slide_index": slide_index}, deadline
        )
        chart_name = _chart_name_from_result(listed)
        if not chart_name:
            raise RuntimeError(f"PowerPoint MCP could not identify the scaffold chart on slide {slide_index}.")
        rows = list(visual.data) if isinstance(visual.data, list) else []
        x_key = visual.x or (next(iter(rows[0])) if rows and isinstance(rows[0], dict) else "Category")
        y_key = visual.y or (
            next((key for key in rows[0] if key != x_key), "Value") if rows and isinstance(rows[0], dict) else "Value"
        )
        values = [[x_key, y_key, "", ""]]
        values.extend([[row.get(x_key), row.get(y_key), None, None] for row in rows if isinstance(row, dict)])
        values.append(["", None, None, None])
        await self._call_fixed_mcp(
            session, "chart",
            {
                "action": "set-data", "session_id": sid, "slide_index": slide_index,
                "shape_name": chart_name, "values": values,
            }, deadline,
        )
        await self._call_fixed_mcp(
            session, "chart",
            {
                "action": "set-title", "session_id": sid, "slide_index": slide_index,
                "shape_name": chart_name, "title": visual.title or "Analysis",
            }, deadline,
        )
        await self._call_fixed_mcp(
            session, "chart",
            {
                "action": "set-legend", "session_id": sid, "slide_index": slide_index,
                "shape_name": chart_name, "visible": False, "position": 0,
            }, deadline,
        )

    async def _add_scaffold_decision_tree(
        self, session: Any, sid: str, slide_index: int, visual: Any, deadline: float
    ) -> None:
        """Add a compact editable tree so the required model-evidence slide is never blank."""
        data = visual.data if isinstance(visual.data, dict) else {}
        nodes = [node for node in data.get("nodes", []) if isinstance(node, dict)]
        if nodes:
            labels = [
                str(node.get("display_label") or node.get("label") or node.get("rule") or "Decision node")
                for node in nodes[:3]
            ]
            while len(labels) < 3:
                labels.append("Verified branch")
        else:
            labels = [
                str(data.get("root") or "Primary driver threshold"),
                str(data.get("left") or "Higher-risk branch"),
                str(data.get("right") or "Lower-risk branch"),
            ]
        geometries = [(565, 172, 210, 58), (468, 304, 180, 62), (692, 304, 180, 62)]
        for node_index, (label, geometry) in enumerate(zip(labels, geometries)):
            left, top, width, height = geometry
            node_text = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-shape", "session_id": sid, "slide_index": slide_index,
                    "left": left, "top": top, "width": width, "height": height,
                    "auto_shape_type": 5, "text": label,
                }, deadline,
            )
            node_name = _shape_name_from_result(node_text, f"Rounded Rectangle {node_index + 2}")
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-fill", "session_id": sid, "slide_index": slide_index,
                    "shape_name": node_name, "color_hex": "DDF3EF" if node_index else "DCEAF5",
                }, deadline,
            )
            await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "set-line", "session_id": sid, "slide_index": slide_index,
                    "shape_name": node_name, "color_hex": "1F7A75", "line_width": 1.5,
                }, deadline,
            )
            await self._call_fixed_mcp(
                session, "text",
                {
                    "action": "format", "session_id": sid, "slide_index": slide_index,
                    "shape_name": node_name, "font_name": "Aptos", "font_size": 14,
                    "bold": node_index == 0, "color": "#102A43",
                }, deadline,
            )
        for left, arrow in ((550, "/"), (744, "\\")):
            arrow_text = await self._call_fixed_mcp(
                session, "shape",
                {
                    "action": "add-textbox", "session_id": sid, "slide_index": slide_index,
                    "left": left, "top": 242, "width": 54, "height": 42, "text": arrow,
                }, deadline,
            )
            arrow_name = _shape_name_from_result(arrow_text, "TextBox 4")
            await self._call_fixed_mcp(
                session, "text",
                {
                    "action": "format", "session_id": sid, "slide_index": slide_index,
                    "shape_name": arrow_name, "font_name": "Aptos", "font_size": 24,
                    "bold": True, "color": "#C49A3A",
                }, deadline,
            )

    async def _render_for_visual_qa(self, session: Any, output: Path, deadline: float) -> None:
        """Deterministically render every slide; do not rely on the model remembering QA export."""
        if not self._active_session_id:
            raise RuntimeError("PowerPoint MCP did not return a presentation session id for rendering.")
        render_dir = output.parent / f"{output.stem}_rendered_{int(time.time() * 1000)}"
        render_dir.mkdir(parents=True, exist_ok=True)
        fixed_calls = [
            ("file", {"action": "save", "session_id": self._active_session_id}),
            (
                "export",
                {
                    "action": "all-slides-to-images",
                    "session_id": self._active_session_id,
                    "destination_directory": str(render_dir),
                    "width": 1600,
                    "height": 900,
                },
            ),
            ("file", {"action": "close", "session_id": self._active_session_id, "save": True}),
        ]
        for name, arguments in fixed_calls:
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                raise TimeoutError(
                    f"PowerPoint MCP Agent 9 exceeded its {self.timeout_seconds}-second total deadline."
                )
            await self._call_fixed_mcp(session, name, arguments, deadline)
        self.last_rendered_files = sorted(
            str(path.resolve())
            for path in render_dir.iterdir()
            if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
        )
        inspection = inspect_presentation(str(output))
        if len(self.last_rendered_files) < inspection.slide_count:
            raise RuntimeError(
                "PowerPoint MCP visual QA export was incomplete: "
                f"expected {inspection.slide_count} rendered slides, found {len(self.last_rendered_files)}."
            )

    async def _call_fixed_mcp(
        self,
        session: Any,
        name: str,
        arguments: dict[str, Any],
        deadline: float,
    ) -> str:
        """Run a required non-model MCP operation with bounded transient retries."""
        last_error = "unknown MCP failure"
        for attempt in range(1, 4):
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                raise TimeoutError(
                    f"PowerPoint MCP Agent 9 exceeded its {self.timeout_seconds}-second total deadline."
                )
            result = await session.call_tool(
                name,
                arguments,
                read_timeout_seconds=timedelta(
                    seconds=max(1, min(self.request_timeout_seconds, int(remaining)))
                ),
            )
            result_text = _mcp_result_text(result)
            if not result.isError and not _mcp_result_reports_error(result_text):
                return result_text
            last_error = result_text[:500]
            if attempt < 3:
                await asyncio.sleep(0.4 * attempt)
        raise RuntimeError(
            f"PowerPoint MCP {name} failed during required operation after 3 attempts: {last_error}"
        )

    async def _run_tool_round(
        self,
        session: Any,
        messages: list[dict[str, Any]],
        tool_defs: list[dict[str, Any]],
        tools_by_name: dict[str, Any],
        output: Path,
        allowed_inputs: set[str],
        *,
        call_budget: int,
        deadline: float,
        allowed_slide_indexes: set[int] | None = None,
        allowed_actions_by_tool: dict[str, set[str]] | None = None,
        strict_budget: bool = True,
    ) -> int:
        if call_budget <= 0:
            raise RuntimeError("PowerPoint MCP repair call budget was exhausted.")
        calls = 0
        invalid_calls = 0
        while calls < call_budget:
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                raise TimeoutError(
                    f"PowerPoint MCP Agent 9 exceeded its {self.timeout_seconds}-second total deadline."
                )
            message = self.openrouter_client.chat_completion_with_tools(
                messages,
                tool_defs,
                timeout_seconds=max(1, min(self.request_timeout_seconds, int(remaining))),
            )
            assistant_message = {
                "role": "assistant",
                "content": message.get("content") or "",
            }
            tool_calls = message.get("tool_calls") or []
            if tool_calls:
                assistant_message["tool_calls"] = tool_calls
            messages.append(assistant_message)
            if not tool_calls:
                break
            for call in tool_calls:
                if time.monotonic() >= deadline:
                    raise TimeoutError(
                        f"PowerPoint MCP Agent 9 exceeded its {self.timeout_seconds}-second total deadline."
                    )
                calls += 1
                if calls > call_budget:
                    raise RuntimeError("PowerPoint MCP exceeded the bounded Agent 9 tool-call limit.")
                function = call.get("function") or {}
                name = str(function.get("name", ""))
                if name not in tools_by_name:
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": f"Tool call blocked: {name} is not available in this enhancement batch.",
                        }
                    )
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        self._logger.warning("Stopping Agent 9 batch after 3 invalid tool selections.")
                        return calls
                    continue
                raw_arguments = function.get("arguments") or "{}"
                try:
                    arguments = raw_arguments if isinstance(raw_arguments, dict) else json.loads(raw_arguments)
                except json.JSONDecodeError as exc:
                    error_text = (
                        f"Tool arguments were invalid JSON at character {exc.pos}: {exc.msg}. "
                        "Correct the JSON and retry this tool call. Raw arguments: "
                        f"{str(raw_arguments)[:1200]}"
                    )
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": error_text,
                        }
                    )
                    self._logger.warning("Agent 9 malformed %s arguments: %s", name, error_text[:500])
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        self._logger.warning("Stopping Agent 9 batch after 3 malformed tool calls.")
                        return calls
                    continue
                action = str(arguments.get("action", "")).strip().lower()
                if not action:
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": "Tool call blocked: action is required. Choose a valid action and retry.",
                        }
                    )
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        return calls
                    continue
                if allowed_actions_by_tool is not None and action not in allowed_actions_by_tool.get(name, set()):
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": f"Tool call blocked: action {action!r} is not read-only for {name}.",
                        }
                    )
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        return calls
                    continue
                if action in {"delete", "remove", "close", "create", "open", "save"}:
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": (
                                "Tool call blocked: the deterministic scaffold already owns file creation and slide "
                                "count. Inspect and enhance the existing slides without creating or deleting them."
                            ),
                        }
                    )
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        return calls
                    continue
                raw_slide_index = arguments.get("slide_index")
                if allowed_slide_indexes is not None and raw_slide_index is not None:
                    try:
                        slide_index = int(raw_slide_index)
                    except (TypeError, ValueError):
                        slide_index = -1
                    if slide_index not in allowed_slide_indexes:
                        messages.append(
                            {
                                "role": "tool",
                                "tool_call_id": call.get("id", f"call-{calls}"),
                                "content": (
                                    f"Tool call blocked: slide {raw_slide_index} is outside the current batch "
                                    f"{sorted(allowed_slide_indexes)}."
                                ),
                            }
                        )
                        invalid_calls += 1
                        if invalid_calls >= 3:
                            return calls
                        continue
                try:
                    self._validate_tool_call(name, arguments, output, allowed_inputs)
                except RuntimeError as exc:
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id", f"call-{calls}"),
                            "content": f"Tool call blocked before execution: {exc}. Correct the arguments and retry.",
                        }
                    )
                    self._logger.warning("Blocked Agent 9 %s call: %s", name, exc)
                    invalid_calls += 1
                    if invalid_calls >= 3:
                        return calls
                    continue
                if self._active_session_id:
                    # The deterministic runtime owns the live PowerPoint session. Never
                    # allow a model-supplied or hallucinated identifier to replace it.
                    arguments["session_id"] = self._active_session_id
                self._logger.info(
                    "Agent 9 enhancement call: tool=%s action=%s slide=%s",
                    name,
                    action,
                    arguments.get("slide_index", "n/a"),
                )
                result = await session.call_tool(
                    name,
                    arguments,
                    read_timeout_seconds=timedelta(
                        seconds=max(1, min(self.request_timeout_seconds, int(deadline - time.monotonic())))
                    ),
                )
                result_text = _mcp_result_text(result)
                messages.append(
                    {
                        "role": "tool",
                        "tool_call_id": call.get("id", f"call-{calls}"),
                        "content": result_text[:12000],
                    }
                )
                if result.isError or _mcp_result_reports_error(result_text):
                    self._logger.warning("PowerPoint MCP tool %s returned an error: %s", name, result_text[:500])
                    self._logger.warning("Stopping the current Agent 9 batch after the first MCP tool error.")
                    return calls
        else:
            if strict_budget:
                raise RuntimeError("PowerPoint MCP Agent 9 loop did not terminate within its call limit.")
            self._logger.warning("Agent 9 enhancement batch reached its bounded %s-call limit.", call_budget)
        return calls

    @staticmethod
    def _validate_tool_call(
        name: str,
        arguments: dict[str, Any],
        output: Path,
        allowed_inputs: set[str],
    ) -> None:
        output_dir = output.parent.resolve()
        if name == "file":
            action = str(arguments.get("action", ""))
            if action not in {"create", "open", "save", "close", "list", "test"}:
                raise RuntimeError(f"Blocked PowerPoint MCP file action: {action}")
        for key in ("path", "destination_path", "destination_directory"):
            raw = arguments.get(key)
            if not raw:
                continue
            path = Path(str(raw)).resolve()
            if path != output and output_dir not in path.parents:
                raise RuntimeError(f"Blocked PowerPoint MCP path outside the run directory: {path}")
        image_path = arguments.get("image_path")
        if image_path:
            resolved = str(Path(str(image_path)).resolve())
            if resolved not in allowed_inputs and output_dir not in Path(resolved).parents:
                raise RuntimeError(f"Blocked untrusted PowerPoint MCP image path: {resolved}")


def _apply_hybrid_story_contract(deck: DeckSpec) -> None:
    """Replace vague post-EDA roles in the standard 12-slide story with explicit action objectives."""
    if len(deck.slides) != 12 or len(deck.slides) < 8:
        return
    roles = [slide.slide_role for slide in deck.slides]
    if roles[7:12] == ["recommendations", "recommendations", "limitations", "sources", "ending"]:
        return
    by_role: dict[str, Any] = {slide.slide_role: slide for slide in deck.slides}
    required = {"recommendations", "limitations", "summary"}
    if not required.issubset(by_role):
        return
    recommendations = copy.deepcopy(by_role["recommendations"])
    recommendations.slide_role = "recommendations"
    recommendations.template = "recommendation_priority"

    recommendation_plan = copy.deepcopy(by_role.get("business_translation") or recommendations)
    recommendation_plan.slide_role = "recommendations"
    recommendation_plan.template = "recommendation_priority"
    recommendation_plan.headline = "Recommendation delivery requires phased ownership and measurable guardrails"
    recommendation_plan.main_message = "Sequence the recommended actions, assign owners, and monitor evidence-backed outcomes."
    execution_items = deck.metadata.get("recommendation_execution", []) if isinstance(deck.metadata, dict) else []
    execution_cards: list[ContentBlock] = []
    execution_evidence_ids: list[str] = []
    for fallback, item in enumerate(execution_items if isinstance(execution_items, list) else [], start=1):
        if not isinstance(item, dict):
            continue
        rank = item.get("rank") or fallback
        owner = _shorten_human(str(item.get("owner") or "named accountable owner"), 45)
        trigger = _shorten_human(str(item.get("trigger") or "approved baseline and protocol"), 60)
        metric = _shorten_human(str(item.get("validation_metric") or "outcome against baseline"), 62)
        action = _shorten_human(str(item.get("action") or f"Execute priority {rank} after validation"), 105)
        execution_cards.append(
            ContentBlock(
                type="recommendation_card",
                title=f"Priority {rank}",
                text=action,
                items=[
                    "Owner: named domain expert",
                    "Measure: precision, recall, safety, and fairness",
                ],
            )
        )
        execution_evidence_ids.extend(str(value) for value in item.get("evidence_ids", []) or [] if value)
    if execution_cards:
        recommendation_plan.content_blocks = execution_cards
    elif by_role.get("business_translation") is None:
        recommendation_plan.content_blocks = [
            ContentBlock(
                type="bullets",
                items=[
                    "Phase 1 - confirm scope, owner, and baseline",
                    "Phase 2 - pilot the highest-priority action with guardrails",
                    "Phase 3 - measure outcomes and scale only supported changes",
                ],
            )
        ]
    recommendation_plan.evidence_ids = list(dict.fromkeys(execution_evidence_ids))[:5]
    if not recommendations.evidence_ids:
        recommendations.evidence_ids = list(recommendation_plan.evidence_ids)

    limitations = copy.deepcopy(by_role["limitations"])
    limitations.slide_role = "limitations"
    limitations.template = "limitations_professional"

    sources = copy.deepcopy(by_role.get("findings") or by_role["summary"])
    sources.slide_role = "sources"
    sources.template = "limitations_professional"
    sources.headline = "Sources and conclusions define the next executive decision"
    sources.main_message = "Trace each conclusion to dataset, analysis, model, or external-source evidence before action."
    source_metadata = deck.metadata.get("evidence_sources", {}) if isinstance(deck.metadata, dict) else {}
    source_items: list[str] = []
    for dataset in source_metadata.get("datasets", []) if isinstance(source_metadata, dict) else []:
        if not isinstance(dataset, dict):
            continue
        source_items.append(
            f"Dataset: {dataset.get('name', 'provided CSV')} - {dataset.get('rows', 'unknown')} rows, "
            f"{dataset.get('columns', 'unknown')} columns; hash {str(dataset.get('sha256') or 'not recorded')[:12]}."
        )
    for external in (source_metadata.get("external_sources", [])[:2] if isinstance(source_metadata, dict) else []):
        if not isinstance(external, dict):
            continue
        url = str(external.get("url") or "")
        domain_match = re.search(r"https?://(?:www\.)?([^/]+)", url)
        domain = domain_match.group(1) if domain_match else url
        source_items.append(
            f"[{external.get('index', '?')}] {_shorten_human(str(external.get('title') or 'External source'), 75)} "
            f"- {domain or 'source page'} ({external.get('evidence_level') or 'unrated'}; context only)."
        )
    evidence_ids = source_metadata.get("evidence_ids", []) if isinstance(source_metadata, dict) else []
    if evidence_ids:
        source_items.append("Analysis evidence: " + ", ".join(str(item) for item in evidence_ids[:5]) + ".")
    bundle_hash = str(deck.metadata.get("evidence_bundle_hash", "")) if isinstance(deck.metadata, dict) else ""
    if bundle_hash:
        source_items.append("Evidence bundle: " + bundle_hash[:16] + ".")
    quality_status = str(source_metadata.get("quality_status", "")) if isinstance(source_metadata, dict) else ""
    if quality_status:
        source_items.append("Quality status: " + quality_status.replace("_", " ") + ".")
    original_findings = by_role.get("findings")
    if original_findings is not None:
        conclusion_items = [
            str(item).strip()
            for block in original_findings.content_blocks
            for item in block.items
            if str(item).strip()
        ]
        prioritized_conclusions = sorted(
            conclusion_items,
            key=lambda item: (0 if ":" in item and any(char.isdigit() for char in item) else 1),
        )
        source_items.extend(item for item in prioritized_conclusions[:2] if item not in source_items)
    sources.content_blocks = [
        ContentBlock(
            type="bullets",
            items=(source_items + ["Next decision: approve, refine, or reject the evidence-gated actions."])[:5],
        )
    ]
    sources.source_ids = [
        str(item.get("source_id"))
        for item in (source_metadata.get("datasets", []) + source_metadata.get("external_sources", []))
        if isinstance(item, dict) and item.get("source_id")
    ]
    sources.evidence_ids = [str(item) for item in evidence_ids[:5]]

    ending = copy.deepcopy(by_role["summary"])
    ending.slide_role = "ending"
    ending.template = "executive_summary_closing"
    deck.slides = list(deck.slides[:7]) + [recommendations, recommendation_plan, limitations, sources, ending]
    deck.renumber()


def inspect_presentation(
    output_path: str,
    *,
    expected_deck: DeckSpec | None = None,
) -> PresentationInspection:
    path = Path(output_path).resolve()
    inspection = PresentationInspection(path=str(path))
    if not path.is_file() or path.stat().st_size < 1000:
        inspection.issues.append({"slide": 0, "issue": "presentation_missing_or_empty"})
        return inspection
    if not PPTX_AVAILABLE:
        inspection.issues.append({"slide": 0, "issue": "python_pptx_unavailable_for_qa"})
        return inspection
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        inspection.issues.append({"slide": 0, "issue": "presentation_cannot_open", "detail": str(exc)})
        return inspection
    inspection.slide_count = len(presentation.slides)
    if inspection.slide_count == 0:
        inspection.issues.append({"slide": 0, "issue": "no_slides"})
        return inspection
    if expected_deck is not None and inspection.slide_count != len(expected_deck.slides):
        inspection.issues.append(
            {
                "slide": 0,
                "issue": "incorrect_slide_count",
                "expected": len(expected_deck.slides),
                "actual": inspection.slide_count,
            }
        )
    font_names: set[str] = set()
    background_signatures: list[str] = []
    analysis_messages: list[tuple[int, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        visible_text = [shape.text.strip() for shape in text_shapes if shape.text.strip()]
        combined_visible = " ".join(visible_text)
        if re.search(r"\{\s*[\"']?[A-Za-z_]+[\"']?\s*:", combined_visible):
            inspection.issues.append({"slide": index, "issue": "raw_json_in_slide_text", "severity": "error"})
        if re.search(r"[\u3400-\u9fff]", combined_visible):
            inspection.issues.append({"slide": index, "issue": "output_language_mismatch", "severity": "error"})
        if not visible_text and not any(getattr(shape, "shape_type", None) == 13 for shape in slide.shapes):
            inspection.issues.append({"slide": index, "issue": "empty_slide"})
        title = slide.shapes.title
        has_visual_title = bool(title is not None and title.text.strip()) or any(
            shape.text.strip() and float(shape.top) / 914400 <= 1.65 and float(shape.height) / 914400 <= 1.6
            for shape in text_shapes
        )
        if not has_visual_title:
            inspection.issues.append({"slide": index, "issue": "missing_title"})
        if expected_deck is not None and index <= len(expected_deck.slides):
            expected = expected_deck.slides[index - 1]
            headline = str(expected.headline or "").strip().lower()
            combined_text = " ".join(visible_text).lower()
            if headline and headline not in combined_text:
                inspection.issues.append(
                    {
                        "slide": index,
                        "issue": "expected_headline_missing",
                        "expected": expected.headline,
                    }
                )
            if expected.evidence_ids and expected.slide_role in {"analysis", "recommendations", "sources", "limitations", "ending"}:
                if not any(str(evidence_id).lower() in combined_text for evidence_id in expected.evidence_ids):
                    inspection.issues.append(
                        {
                            "slide": index,
                            "issue": "visible_evidence_reference_missing",
                            "expected": expected.evidence_ids[:3],
                            "severity": "error" if expected.slide_role in {"analysis", "recommendations", "sources"} else "warning",
                        }
                    )
            if expected.slide_role == "sources" and not any(
                marker in combined_text for marker in ("dataset:", "external:", "search_snippet", "context only", "evidence:")
            ):
                inspection.issues.append(
                    {
                        "slide": index,
                        "issue": "source_inventory_not_specific",
                        "severity": "error" if (expected.source_ids or expected.evidence_ids) else "warning",
                    }
                )
            if expected.slide_role == "analysis":
                body_text = " ".join(text for text in visible_text if text.lower() != headline)
                if body_text:
                    analysis_messages.append((index, re.sub(r"\s+", " ", body_text.lower()).strip()))
            if expected.visual is not None and expected.visual.type in {
                "code_figure",
                "structured_chart",
                "legacy_image_fallback",
            }:
                shape_types = [getattr(shape, "shape_type", None) for shape in slide.shapes]
                if expected.visual.type in {"code_figure", "legacy_image_fallback"}:
                    has_visual = 13 in shape_types
                else:
                    vector_visuals = [
                        shape
                        for shape in slide.shapes
                        if getattr(shape, "shape_type", None) in {1, 3, 6, 9, 13}
                        and not _is_full_slide_background(shape, presentation.slide_width, presentation.slide_height)
                    ]
                    has_visual = any(shape_type in {3, 13} for shape_type in shape_types) or len(vector_visuals) >= 2
                if not has_visual:
                    inspection.issues.append({"slide": index, "issue": "missing_expected_visual"})
        for shape in text_shapes:
            text = shape.text.strip()
            if text and _likely_text_overflow(shape, text):
                expected_role = ""
                if expected_deck is not None and index <= len(expected_deck.slides):
                    expected_role = str(expected_deck.slides[index - 1].slide_role).lower()
                inspection.issues.append(
                    {
                        "slide": index,
                        "issue": "possible_text_overflow",
                        "shape": shape.name,
                        "severity": (
                            "error"
                            if expected_role in {"recommendations", "limitations", "ending"}
                            else "warning"
                        ),
                    }
                )
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        font_names.add(str(run.font.name))
                    if run.font.size and float(run.font.size.pt) < 8:
                        inspection.issues.append(
                            {"slide": index, "issue": "font_below_readability_floor", "shape": shape.name, "severity": "warning"}
                        )
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        for shape in slide.shapes:
            if (
                shape.left < -1000
                or shape.top < -1000
                or shape.left + shape.width > slide_width + 1000
                or shape.top + shape.height > slide_height + 1000
            ):
                inspection.issues.append(
                    {"slide": index, "issue": "object_outside_slide", "shape": shape.name}
                )
            if getattr(shape, "shape_type", None) == 13:
                try:
                    if not shape.image.blob:
                        raise ValueError("empty image blob")
                except Exception as exc:
                    inspection.issues.append(
                        {
                            "slide": index,
                            "issue": "broken_image",
                            "shape": shape.name,
                            "detail": str(exc),
                        }
                    )
        inspection.issues.extend(_overlap_issues(slide, index))
        background_signatures.append(_slide_background_signature(slide, slide_width, slide_height))
    if len(font_names) > 5:
        inspection.issues.append(
            {
                "slide": 0,
                "issue": "inconsistent_typography",
                "fonts": sorted(font_names),
                "severity": "warning",
            }
        )
    if inspection.slide_count >= 4 and len(set(background_signatures)) == 1 and background_signatures[0] in {
        "FFFFFF",
        "none",
    }:
        inspection.issues.append(
            {"slide": 0, "issue": "uniform_plain_background", "severity": "warning"}
        )
    for first_index, (first_slide, first_text) in enumerate(analysis_messages):
        first_tokens = set(re.findall(r"[a-z0-9]+", first_text))
        if len(first_tokens) < 12:
            continue
        for second_slide, second_text in analysis_messages[first_index + 1 :]:
            second_tokens = set(re.findall(r"[a-z0-9]+", second_text))
            overlap = len(first_tokens & second_tokens) / max(1, min(len(first_tokens), len(second_tokens)))
            if overlap >= 0.78:
                inspection.issues.append(
                    {
                        "slide": second_slide,
                        "issue": "repeated_analysis_explanation",
                        "matches_slide": first_slide,
                        "severity": "warning",
                    }
                )
    return inspection


def _slide_background_signature(slide: Any, slide_width: int, slide_height: int) -> str:
    for shape in slide.shapes:
        if (
            shape.left <= 1000
            and shape.top <= 1000
            and shape.width >= slide_width * 0.98
            and shape.height >= slide_height * 0.98
        ):
            try:
                rgb = shape.fill.fore_color.rgb
                if rgb:
                    return str(rgb)
            except (AttributeError, TypeError):
                pass
    try:
        rgb = slide.background.fill.fore_color.rgb
        return str(rgb) if rgb else "none"
    except (AttributeError, TypeError):
        return "none"


def _is_full_slide_background(shape: Any, slide_width: int, slide_height: int) -> bool:
    return (
        shape.left <= 1000
        and shape.top <= 1000
        and shape.width >= slide_width * 0.98
        and shape.height >= slide_height * 0.98
    )


def _likely_text_overflow(shape: Any, text: str) -> bool:
    width_inches = max(float(shape.width) / 914400, 0.2)
    height_inches = max(float(shape.height) / 914400, 0.2)
    font_sizes: list[float] = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.font.size:
                font_sizes.append(float(paragraph.font.size.pt))
            for run in paragraph.runs:
                if run.font.size:
                    font_sizes.append(float(run.font.size.pt))
    except (AttributeError, TypeError, ValueError):
        pass
    font_size = max(font_sizes) if font_sizes else 18.0
    characters_per_line = max(1, int((width_inches * 72) / max(font_size * 0.52, 1)))
    available_lines = max(1, int((height_inches * 72) / max(font_size * 1.22, 1)))
    required_lines = sum(
        max(1, (len(line.strip()) + characters_per_line - 1) // characters_per_line)
        for line in text.splitlines() or [text]
    )
    return required_lines > available_lines


def _overlap_issues(slide: Any, slide_index: int) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    shapes = [shape for shape in slide.shapes if getattr(shape, "width", 0) and getattr(shape, "height", 0)]
    for first_index, first in enumerate(shapes):
        for second in shapes[first_index + 1 :]:
            if _contains(first, second) or _contains(second, first):
                continue
            intersection = _intersection_area(first, second)
            smaller = min(first.width * first.height, second.width * second.height)
            if smaller and intersection / smaller > 0.35:
                issues.append(
                    {
                        "slide": slide_index,
                        "issue": "possible_object_overlap",
                        "shapes": [first.name, second.name],
                        "severity": "warning",
                    }
                )
    return issues[:8]


def _intersection_area(first: Any, second: Any) -> int:
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    return max(0, right - left) * max(0, bottom - top)


def _contains(outer: Any, inner: Any) -> bool:
    return (
        outer.left <= inner.left
        and outer.top <= inner.top
        and outer.left + outer.width >= inner.left + inner.width
        and outer.top + outer.height >= inner.top + inner.height
    )


def _mcp_result_text(result: Any) -> str:
    if getattr(result, "structuredContent", None):
        return json.dumps(result.structuredContent, default=str)
    parts: list[str] = []
    for item in getattr(result, "content", []) or []:
        text = getattr(item, "text", None)
        if text:
            parts.append(str(text))
        elif hasattr(item, "model_dump"):
            parts.append(json.dumps(item.model_dump(), default=str))
        else:
            parts.append(str(item))
    return "\n".join(parts) or "PowerPoint MCP tool returned no content."


def _session_id_from_result(result_text: str) -> str | None:
    """Read a session id from structured or text MCP responses without assuming one casing."""
    try:
        payload = json.loads(result_text)
    except (json.JSONDecodeError, TypeError):
        payload = None
    if isinstance(payload, dict):
        for key in ("session_id", "sessionId", "sessionID"):
            value = payload.get(key)
            if value:
                return str(value)
        nested = payload.get("data")
        if isinstance(nested, dict):
            for key in ("session_id", "sessionId", "sessionID"):
                value = nested.get(key)
                if value:
                    return str(value)
    match = re.search(r'["\']session(?:_?id)?["\']\s*[:=]\s*["\']([^"\']+)', result_text, re.IGNORECASE)
    return match.group(1) if match else None


def _shape_name_from_result(result_text: str, fallback: str) -> str:
    try:
        payload = json.loads(result_text)
    except (json.JSONDecodeError, TypeError):
        payload = None
    if isinstance(payload, dict):
        for key in ("shape_name", "shapeName", "name"):
            if payload.get(key):
                return str(payload[key])
        nested = payload.get("data")
        if isinstance(nested, dict):
            for key in ("shape_name", "shapeName", "name"):
                if nested.get(key):
                    return str(nested[key])
        message = str(payload.get("message") or payload.get("errorMessage") or "")
        match = re.search(r"['\"]([^'\"]+)['\"]", message)
        if match:
            return match.group(1)
    match = re.search(r"(?:Shape|shape)\s+['\"]([^'\"]+)['\"]", result_text)
    return match.group(1) if match else fallback


def _chart_name_from_result(result_text: str) -> str | None:
    try:
        payload = json.loads(result_text)
    except (json.JSONDecodeError, TypeError):
        return None
    shapes = payload.get("shapes", []) if isinstance(payload, dict) else []
    if not shapes and isinstance(payload, dict) and isinstance(payload.get("data"), dict):
        shapes = payload["data"].get("shapes", [])
    for shape in shapes if isinstance(shapes, list) else []:
        if isinstance(shape, dict) and (shape.get("hasChart") or str(shape.get("shapeType", "")).lower() == "chart"):
            return str(shape.get("name") or shape.get("shapeName") or "") or None
    return None


def _mcp_result_reports_error(result_text: str) -> bool:
    try:
        payload = json.loads(result_text)
    except (json.JSONDecodeError, TypeError):
        return bool(re.search(r"\b(error|failed)\b", result_text, flags=re.IGNORECASE))
    return bool(payload.get("isError") or payload.get("is_error") or payload.get("success") is False)


_AGENT9_TOOL_FIELDS: dict[str, set[str]] = {
    "file": {"action", "path", "session_id", "save", "show", "timeout_seconds"},
    "slide": {"action", "session_id", "slide_index", "position", "layout_name"},
    "shape": {
        "action",
        "session_id",
        "slide_index",
        "shape_name",
        "left",
        "top",
        "width",
        "height",
        "text",
        "auto_shape_type",
        "color_hex",
        "line_width",
        "z_order_cmd",
        "alt_text",
        "margin_left",
        "margin_right",
        "margin_top",
        "margin_bottom",
        "word_wrap",
        "auto_size",
    },
    "text": {
        "action",
        "session_id",
        "slide_index",
        "shape_name",
        "text",
        "font_name",
        "font_size",
        "bold",
        "italic",
        "color",
        "alignment",
        "vertical_alignment",
    },
    "chart": {
        "action",
        "session_id",
        "slide_index",
        "chart_type",
        "left",
        "top",
        "width",
        "height",
        "shape_name",
        "title",
        "values",
        "visible",
        "position",
    },
    "slidetable": {
        "action",
        "session_id",
        "slide_index",
        "rows",
        "columns",
        "left",
        "top",
        "width",
        "height",
        "shape_name",
        "row",
        "column",
        "value",
        "values",
        "fill_color",
        "font_bold",
        "font_size",
        "text_align",
    },
    "image": {
        "action",
        "session_id",
        "slide_index",
        "image_path",
        "left",
        "top",
        "width",
        "height",
        "shape_name",
    },
    "placeholder": {
        "action",
        "session_id",
        "slide_index",
        "placeholder_index",
        "text",
        "image_path",
    },
    "notes": {"action", "session_id", "slide_index", "text"},
    "shapealign": {"action", "session_id", "slide_index", "shape_names", "align_type", "distribute_type"},
    "export": {
        "action",
        "session_id",
        "destination_path",
        "destination_directory",
        "slide_index",
        "width",
        "height",
    },
}


def _agent9_tool_definition(tool: Any, *, allowed_actions: set[str] | None = None) -> dict[str, Any]:
    schema = json.loads(json.dumps(tool.inputSchema))
    allowed = _AGENT9_TOOL_FIELDS.get(tool.name, set(schema.get("properties", {})))
    schema["properties"] = {
        name: definition
        for name, definition in schema.get("properties", {}).items()
        if name in allowed
    }
    schema["required"] = [name for name in schema.get("required", []) if name in allowed]
    action_schema = schema.get("properties", {}).get("action")
    if allowed_actions and isinstance(action_schema, dict):
        action_schema["enum"] = sorted(allowed_actions)
    return {
        "type": "function",
        "function": {
            "name": tool.name,
            "description": f"PowerPoint {tool.name} operations for Agent 9 deck creation and QA.",
            "parameters": schema,
        },
    }


def backend_from_config(config: Any, openrouter_client: Any) -> PresentationBackend:
    backend = str(getattr(config, "presentation_backend", "auto") or "auto").strip().lower()
    if backend == "python":
        return PythonPresentationBackend()
    if backend not in {"auto", "powerpoint_mcp"}:
        raise ValueError("PRESENTATION_BACKEND must be auto, powerpoint_mcp, or python.")
    return PowerPointMCPBackend(
        openrouter_client,
        command=str(getattr(config, "powerpoint_mcp_command", "mcp-ppt") or "mcp-ppt"),
        timeout_seconds=int(getattr(config, "presentation_agent_timeout_seconds", 900)),
        request_timeout_seconds=int(getattr(config, "agent_request_timeout_seconds", 180)),
    )
