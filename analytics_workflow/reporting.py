from __future__ import annotations

import html
import json
import os
import re
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Any

from .decision_tree_figure import decision_tree_imbalance_note, decision_tree_performance_note
from .slides.text_refiner import normalize_output_text, soften_unsupported_impact_claim
from .tree_diagram import build_tree_layout, clean_tree_text

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.graphics.shapes import Drawing, Line, Rect, String
    from reportlab.platypus import KeepTogether, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def _stringify(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return soften_unsupported_impact_claim(normalize_output_text(value)).strip()
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, dict):
        return "; ".join(
            f"{str(key).replace('_', ' ').title()}: {_stringify(item)}"
            for key, item in value.items()
            if _stringify(item)
        )
    if isinstance(value, (list, tuple, set)):
        return "; ".join(_stringify(item) for item in value if _stringify(item))
    return soften_unsupported_impact_claim(normalize_output_text(str(value)))


def _safe_paragraph(text: str) -> str:
    return html.escape(text).replace("\n", "<br/>")


def _add_pdf_heading(story: list[Any], text: str, style: Any) -> None:
    story.append(Paragraph(_safe_paragraph(text), style))
    story.append(Spacer(1, 0.08 * inch))


def _add_pdf_body(story: list[Any], text: str, style: Any) -> None:
    clean = soften_unsupported_impact_claim(_stringify(text))
    if clean:
        story.append(Paragraph(_safe_paragraph(clean), style))
        story.append(Spacer(1, 0.1 * inch))


def _add_pdf_bullets(story: list[Any], items: list[str], style: Any) -> None:
    for item in items:
        clean = soften_unsupported_impact_claim(_stringify(item))
        if clean:
            story.append(Paragraph(_safe_paragraph(f"- {clean}"), style))
            story.append(Spacer(1, 0.05 * inch))


def _split_report_label(value: Any, fallback: str) -> tuple[str, str]:
    clean = soften_unsupported_impact_claim(_stringify(value)).strip()
    if ":" in clean:
        label, body = clean.split(":", 1)
        if 1 <= len(label.strip()) <= 34:
            return label.strip(), body.strip()
    return fallback, clean


def _add_pdf_label_table(
    story: list[Any],
    items: list[str],
    body_style: Any,
    label_style: Any,
    *,
    fallback_label: str = "Note",
) -> None:
    rows: list[list[Any]] = []
    for index, item in enumerate(items, start=1):
        label, body = _split_report_label(item, f"{fallback_label} {index}")
        if not body:
            continue
        rows.append(
            [
                Paragraph(f"<b>{_safe_paragraph(label)}</b>", label_style),
                Paragraph(_safe_paragraph(body), body_style),
            ]
        )
    if not rows:
        return
    table = Table(rows, colWidths=[1.28 * inch, 4.92 * inch], hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (0, -1), rl_colors.HexColor("#EAF3F8")),
                ("BACKGROUND", (1, 0), (1, -1), rl_colors.HexColor("#F8FAFC")),
                ("BOX", (0, 0), (-1, -1), 0.6, rl_colors.HexColor("#C9D9E4")),
                ("INNERGRID", (0, 0), (-1, -1), 0.3, rl_colors.HexColor("#DDE6EC")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ]
        )
    )
    story.extend([table, Spacer(1, 0.14 * inch)])


def _add_pdf_numbered_items(story: list[Any], items: list[str], style: Any, *, noun: str = "Step") -> None:
    for index, item in enumerate(items, start=1):
        clean = soften_unsupported_impact_claim(_stringify(item)).strip()
        if not clean:
            continue
        label, body = _split_report_label(clean, f"{noun} {index}")
        if label.lower().startswith(("priority", "step", "question")):
            heading = label
        else:
            heading = f"{noun} {index}"
            body = clean
        story.append(
            Paragraph(
                f"<b>{_safe_paragraph(heading)}.</b> {_safe_paragraph(body)}",
                style,
            )
        )
        story.append(Spacer(1, 0.08 * inch))


def _add_pdf_caveat_table(story: list[Any], items: list[str], style: Any, label_style: Any) -> None:
    rows: list[list[Any]] = [
        [
            Paragraph('<font color="#FFFFFF"><b>Risk or assumption</b></font>', label_style),
            Paragraph('<font color="#FFFFFF"><b>Response</b></font>', label_style),
        ]
    ]
    for item in items:
        clean = soften_unsupported_impact_claim(_stringify(item)).strip()
        if not clean:
            continue
        if "Mitigation:" in clean:
            risk, response = clean.split("Mitigation:", 1)
        else:
            risk, response = clean, "Validate this assumption before scaling the decision."
        rows.append([Paragraph(_safe_paragraph(risk.strip()), style), Paragraph(_safe_paragraph(response.strip()), style)])
    if len(rows) == 1:
        return
    table = Table(rows, colWidths=[3.05 * inch, 3.15 * inch], repeatRows=1, hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#16324F")),
                ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                ("BACKGROUND", (0, 1), (-1, -1), rl_colors.HexColor("#F8FAFC")),
                ("BOX", (0, 0), (-1, -1), 0.6, rl_colors.HexColor("#C9D9E4")),
                ("INNERGRID", (0, 0), (-1, -1), 0.3, rl_colors.HexColor("#DDE6EC")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ]
        )
    )
    story.extend([table, Spacer(1, 0.14 * inch)])


def _format_objective_coverage(workflow_state: dict[str, Any]) -> list[str]:
    objective = workflow_state.get("workflow_objective", {}) or {}
    analysis_summary = workflow_state.get("analysis_results", {}).get("analysis_summary", {}) or {}
    decision = workflow_state.get("agent_outputs", {}).get("decision_maker", {}) or {}
    items: list[str] = []
    raw_description = _clean_objective_display(_stringify(objective.get("raw_description", "")), workflow_state)
    decision_question = _stringify(objective.get("decision_question", ""))
    if raw_description:
        items.append(f"Objective: {raw_description}")
    elif decision_question:
        items.append(f"Objective: {decision_question}")
    user_goal_alignment = ""
    if isinstance(analysis_summary, dict):
        user_goal_alignment = _stringify(analysis_summary.get("user_goal_alignment", ""))
    if user_goal_alignment:
        items.append(f"Analysis alignment: {user_goal_alignment}")
    target_column = _stringify(workflow_state.get("decision_tree_target_column", ""))
    if target_column:
        items.append(f"Decision tree target: {target_column}")
    final_recommendation = _stringify(decision.get("final_recommendation", ""))
    if final_recommendation:
        items.append(f"Decision answer: {final_recommendation}")
    limitations = objective.get("limitations", [])
    if limitations:
        items.append(f"Limitations: {'; '.join(_stringify(item) for item in limitations if _stringify(item))}")
    if not items:
        items.append("No user objective was provided; the workflow optimized for generally decision-useful findings.")
    return items


def _report_title(workflow_state: dict[str, Any]) -> str:
    outputs = workflow_state.get("agent_outputs", {}) or {}
    decision = outputs.get("decision_maker", {}) or {}
    title = _stringify(decision.get("title", "")).strip()
    if title:
        return title
    datasets = (workflow_state.get("evidence_bundle", {}) or {}).get("datasets", []) or []
    name = _stringify(datasets[0].get("name", "")) if datasets else ""
    return f"{Path(name).stem.replace('_', ' ')} Decision Report" if name else "Analytics Decision Report"


def _report_scope_metrics(workflow_state: dict[str, Any]) -> list[tuple[str, str]]:
    bundle = workflow_state.get("evidence_bundle", {}) or {}
    datasets = bundle.get("datasets", []) or []
    first = datasets[0] if datasets else {}
    target = _stringify(workflow_state.get("decision_tree_target_column", "")) or "Not supplied"
    quality = workflow_state.get("quality_receipt", {}) or {}
    return [
        ("Rows", f"{int(first.get('rows') or 0):,}" if first.get("rows") is not None else "Unknown"),
        ("Columns", str(first.get("columns") or "Unknown")),
        ("Target", target.replace("_", " ")),
        ("Sharing status", _stringify(quality.get("status", "not validated")).replace("_", " ").title()),
    ]


def _pdf_page_decorator(title: str, *, cover: bool = False):
    def draw(canvas: Any, doc: Any) -> None:
        width, height = A4
        if cover:
            canvas.saveState()
            canvas.setFillColor(rl_colors.HexColor("#0B2E4F"))
            canvas.rect(0, 0, width, height, stroke=0, fill=1)
            canvas.setFillColor(rl_colors.HexColor("#C39A3B"))
            canvas.rect(0, 0, 18, height, stroke=0, fill=1)
            canvas.restoreState()
            return
        canvas.saveState()
        canvas.setStrokeColor(rl_colors.HexColor("#D8E1E8"))
        canvas.line(0.7 * inch, height - 0.48 * inch, width - 0.7 * inch, height - 0.48 * inch)
        canvas.setFillColor(rl_colors.HexColor("#526677"))
        canvas.setFont("Helvetica", 8)
        canvas.drawString(0.7 * inch, height - 0.36 * inch, title[:72])
        canvas.drawRightString(width - 0.7 * inch, 0.38 * inch, f"{doc.page}")
        canvas.restoreState()
    return draw


def _format_dataset_overview(data_understander: dict[str, Any]) -> str:
    datasets = data_understander.get("datasets", {})
    parts = []
    for dataset_name, dataset_info in datasets.items() if isinstance(datasets, dict) else []:
        if not isinstance(dataset_info, dict):
            summary = _stringify(dataset_info)
            if summary:
                parts.append(f"{dataset_name}: {summary}")
            continue
        summary = dataset_info.get("quality_summary", "")
        analyses = ", ".join(dataset_info.get("recommended_analyses", [])[:3])
        segment = f"{dataset_name}: {summary}"
        if analyses:
            segment += f" Recommended analyses: {analyses}."
        parts.append(segment)
    return "\n".join(parts) or data_understander.get("executive_summary", "")


def _source_index_map(market_researcher: dict[str, Any]) -> dict[int, dict[str, Any]]:
    mapping: dict[int, dict[str, Any]] = {}
    for fallback_index, source in enumerate(market_researcher.get("sources_cited", []), start=1):
        try:
            index = int(source.get("index", fallback_index))
        except (TypeError, ValueError):
            index = fallback_index
        mapping[index] = source
    return mapping


def _source_display(index: int, source: dict[str, Any]) -> str:
    bits = [str(source.get("title", "")).strip(), str(source.get("url", "")).strip()]
    text = f"Source [{index}]: " + " - ".join(bit for bit in bits if bit)
    if source.get("evidence_level") == "search_snippet":
        text += " (search-snippet evidence; validate the source page before decision use)"
    return text


def _format_market_findings(market_researcher: dict[str, Any]) -> list[str]:
    findings: list[str] = []
    overview = market_researcher.get("industry_overview", "")
    if overview:
        findings.append(overview)

    sources = _source_index_map(market_researcher)
    market_findings = market_researcher.get("market_findings", [])
    if market_findings:
        for finding in market_findings[:5]:
            claim = finding.get("claim", "")
            try:
                source_index = int(finding.get("source_index", 1))
            except (TypeError, ValueError):
                source_index = 1
            findings.append(f"{claim} [{source_index}]")
            source = sources.get(source_index)
            if source:
                findings.append(_source_display(source_index, source))
        return findings

    for index, trend in enumerate(market_researcher.get("key_trends", [])[:4], start=1):
        findings.append(f"{trend} [{index}]")
        source = sources.get(index)
        if source:
                findings.append(_source_display(index, source))

    for index, opportunity in enumerate(market_researcher.get("opportunities", [])[:3], start=1):
        findings.append(f"{opportunity} [{index}]")
        source = sources.get(index)
        if source:
                findings.append(_source_display(index, source))

    return findings


def _market_claim_pairs(market_researcher: dict[str, Any]) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    sources = _source_index_map(market_researcher)
    if not sources:
        assumptions = market_researcher.get("key_trends", [])[:3] + market_researcher.get("opportunities", [])[:2]
        if assumptions:
            return [
                (
                    f"Portfolio assumption: {_stringify(item)}",
                    "No source/date was captured for this statement; treat it as internal context until validated.",
                )
                for item in assumptions
            ]
        return []
    market_findings = market_researcher.get("market_findings", [])
    if market_findings:
        for finding in market_findings[:5]:
            claim = finding.get("claim", "")
            try:
                source_index = int(finding.get("source_index", 1))
            except (TypeError, ValueError):
                source_index = 1
            source = sources.get(source_index, {})
            source_text = _source_display(source_index, source) if source else ""
            pairs.append((f"{claim} [{source_index}]", source_text))
        return pairs

    fallback_items = market_researcher.get("key_trends", [])[:4] + market_researcher.get("opportunities", [])[:3]
    for index, item in enumerate(fallback_items, start=1):
        source = sources.get(index, {})
        source_text = _source_display(index, source) if source else ""
        pairs.append((f"{item} [{index}]", source_text))
    return pairs


def _clean_objective_display(text: str, workflow_state: dict[str, Any]) -> str:
    clean = _stringify(text)
    lower = clean.lower()
    if not clean:
        return ""
    if any(
        phrase in lower
        for phrase in (
            "produce dataset-specific",
            "produce a business-readable",
            "executive slides",
            "do not train",
            "use only evidence visible",
        )
    ):
        if _is_stock_time_series_report_context(workflow_state):
            return "Review stock price, volume, return, volatility, and validation-gated portfolio actions from the CSV evidence."
        return "Review the dataset evidence and produce decision-ready findings, recommendations, and validation caveats."
    return clean


def _format_analysis_findings(analysis_results: dict[str, Any]) -> list[str]:
    findings: list[str] = []
    seen: set[str] = set()

    def add_finding(text: str) -> None:
        clean = _stringify(text).strip()
        if not clean:
            return
        if _looks_like_raw_diagnostic(clean):
            return
        if clean in seen:
            return
        seen.add(clean)
        findings.append(clean)

    tree_artifacts = _decision_tree_artifacts(analysis_results)
    for item in analysis_results.get("business_findings", [])[:6]:
        if tree_artifacts and "tree" in _stringify(item).lower():
            continue
        add_finding(item)

    for artifact in list(analysis_results.get("analysis_artifacts", []) or [])[:6]:
        if not isinstance(artifact, dict):
            continue
        title = _stringify(artifact.get("title", "")).strip()
        if _stringify(artifact.get("chart_type", "")).lower() == "decision_tree":
            continue
        finding = _stringify(artifact.get("finding", "") or artifact.get("takeaway", "")).strip()
        if finding:
            label = f"{title}: " if title else "Visual evidence: "
            add_finding(f"{label}{finding}")

    summary = analysis_results.get("analysis_summary", {})
    if isinstance(summary, dict):
        for key, value in list(summary.items())[:10]:
            if isinstance(value, (dict, list)):
                continue
            if any(token in str(key).lower() for token in ("model", "classification_report", "confusion", "tree_nodes", "tree_rules")):
                continue
            if _looks_like_raw_diagnostic(f"{key}: {value}"):
                continue
            if tree_artifacts and "accuracy" in str(key).lower():
                continue
            add_finding(f"{str(key).replace('_', ' ').title()}: {_stringify(value)}")

    figure_captions = analysis_results.get("figure_captions", {})
    if isinstance(figure_captions, dict):
        for figure_name, caption in list(figure_captions.items())[:6]:
            clean_caption = _stringify(caption).strip()
            if not clean_caption:
                continue
            add_finding(f"Visual evidence: {clean_caption}")

    if not findings:
        add_finding("The analysis code executed, but no structured analysis findings were captured for reporting.")

    return findings


def _looks_like_raw_diagnostic(text: str) -> bool:
    clean = _stringify(text)
    lower = clean.lower()
    raw_tokens = (
        "classification_report",
        "confusion_matrix",
        "standardscaler",
        "onehotencoder",
        "columntransformer",
        "pipeline(",
        "figure_",
        ".png",
        "dtype",
        "columns:",
        "preprocess",
        "model diagnostics",
        "tree_nodes",
        "tree_rules",
    )
    if any(token in lower for token in raw_tokens):
        return True
    if len(clean) > 220 and (clean.count("{") + clean.count("[") + clean.count(":")) >= 5:
        return True
    return False


def _decision_tree_artifacts(analysis_results: dict[str, Any]) -> list[dict[str, Any]]:
    artifacts = list(analysis_results.get("analysis_artifacts", []) or []) + list(analysis_results.get("chart_specs", []) or [])
    return [
        artifact
        for artifact in artifacts
        if isinstance(artifact, dict) and _stringify(artifact.get("chart_type", "")).lower() == "decision_tree"
    ]


def _tree_nodes_edges(artifact: dict[str, Any]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    nodes = [node for node in data.get("nodes", []) if isinstance(node, dict)]
    edges = [edge for edge in data.get("edges", []) if isinstance(edge, dict)]
    if nodes:
        return nodes[:15], edges[:20]
    rules = data.get("rules", []) or artifact.get("rules", [])
    rule_texts = [_stringify(rule) for rule in rules if _stringify(rule)]
    if not rule_texts and isinstance(artifact.get("data"), list):
        rule_texts = [_stringify(row.get("rule") or row.get("label")) for row in artifact.get("data", []) if isinstance(row, dict)]
    if not rule_texts:
        return [], []
    nodes = [{"id": "root", "label": "Decision tree rules", "depth": 0}]
    edges = []
    for index, rule in enumerate(rule_texts[:6], start=1):
        node_id = f"rule_{index}"
        nodes.append({"id": node_id, "label": rule, "depth": 1})
        edges.append({"source": "root", "target": node_id, "label": ""})
    return nodes, edges


def _decision_tree_has_valid_graph(artifact: dict[str, Any]) -> bool:
    nodes, edges = _tree_nodes_edges(artifact)
    if not nodes or not edges:
        return False
    node_ids = {_stringify(node.get("id") or node.get("node_id") or node.get("name")) for node in nodes}
    sources = {_stringify(edge.get("source") or edge.get("from") or edge.get("parent")) for edge in edges}
    targets = {_stringify(edge.get("target") or edge.get("to") or edge.get("child")) for edge in edges}
    leaves = [node for node in nodes if _stringify(node.get("id") or node.get("node_id") or node.get("name")) not in sources]
    has_split = any(
        _stringify(node.get("type")).lower() == "split"
        or _stringify(node.get("id") or node.get("node_id") or node.get("name")) in sources
        for node in nodes
    )
    return bool(node_ids and targets.issubset(node_ids) and has_split and len(leaves) >= 2)


def _decision_tree_narrative(artifact: dict[str, Any]) -> str:
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    note = _stringify(decision_tree_imbalance_note(data) or data.get("performance_note") or decision_tree_performance_note(data))
    features = _decision_tree_split_features(artifact)
    if features:
        split_text = ", ".join(features[:2])
        if len(features) > 2:
            split_text += f", and {features[2]}"
        narrative = f"Tree splits use {split_text}."
    elif _decision_tree_has_valid_graph(artifact):
        narrative = "Tree splits are rendered from the structured model artifact."
    else:
        narrative = ""
    if note:
        return f"{narrative} {note}".strip()
    return narrative


def _decision_tree_callouts(artifact: dict[str, Any], limit: int = 3) -> list[str]:
    nodes, edges = _tree_nodes_edges(artifact)
    parent_by_child = {_stringify(edge.get("target") or edge.get("to") or edge.get("child")): edge for edge in edges}
    node_by_id = {_stringify(node.get("id") or node.get("node_id") or node.get("name")): node for node in nodes}
    callouts: list[str] = []
    for node_id, node in node_by_id.items():
        label = _stringify(node.get("label") or node.get("rule") or node.get("prediction"))
        if _stringify(node.get("type")).lower() != "leaf" and "predict" not in label.lower():
            continue
        edge = parent_by_child.get(node_id, {})
        parent = node_by_id.get(_stringify(edge.get("source") or edge.get("from") or edge.get("parent")), {})
        split = _stringify(parent.get("label") or parent.get("rule"))
        condition = _stringify(edge.get("condition") or edge.get("label"))
        text = " | ".join(part for part in (split, condition, label) if part)
        if text:
            callouts.append(_business_tree_rule_text(text, artifact))
        if len(callouts) >= limit:
            break
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    if not callouts:
        callouts = [
            _business_tree_rule_text(rule, artifact)
            for rule in data.get("rules", [])
            if _stringify(rule)
        ][:limit]
    return callouts


def _business_tree_rule_text(rule: Any, artifact: dict[str, Any]) -> str:
    text = _stringify(rule)
    if not text:
        return ""
    text = re.sub(r"\bIf\s+", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\bthen\s+predict\s+", "=> predicted ", text, flags=re.IGNORECASE)
    text = text.replace("[Default]", "higher default risk").replace("[No Default]", "lower default risk")
    text = re.sub(r"\[(1|Yes|True)\]", "higher default risk", text, flags=re.IGNORECASE)
    text = re.sub(r"\[(0|No|False)\]", "lower default risk", text, flags=re.IGNORECASE)
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    target = _stringify(data.get("target") or artifact.get("target")).lower()
    text = text.replace("_", " ")
    if any(token in target for token in ("loan", "default", "credit")):
        text = _credit_rule_language(text)
    return textwrap.shorten(re.sub(r"\s+", " ", text).strip(" ."), width=190, placeholder="")


def _credit_rule_language(text: str) -> str:
    replacements = {
        "loan percent income": "debt burden",
        "loan int rate": "interest rate",
        "person income": "borrower income",
        "person age": "borrower age",
        "person home ownership": "home ownership",
        "cb person default on file": "prior default flag",
        "loan status": "loan outcome",
    }
    clean = text
    for source, replacement in replacements.items():
        clean = re.sub(source, replacement, clean, flags=re.IGNORECASE)
    clean = re.sub(r"borrower income\s*<=\s*(\d+(?:\.\d+)?)", lambda m: f"borrower income at or below ${float(m.group(1)):,.0f}", clean, flags=re.IGNORECASE)
    clean = re.sub(r"borrower income\s*>\s*(\d+(?:\.\d+)?)", lambda m: f"borrower income above ${float(m.group(1)):,.0f}", clean, flags=re.IGNORECASE)
    clean = re.sub(r"debt burden\s*<=\s*(\d+(?:\.\d+)?)", r"debt burden at or below \1", clean, flags=re.IGNORECASE)
    clean = re.sub(r"debt burden\s*>\s*(\d+(?:\.\d+)?)", r"debt burden above \1", clean, flags=re.IGNORECASE)
    clean = clean.replace("<=", "at or below").replace(">=", "at least").replace(">", "above").replace("<", "below")
    return re.sub(r"\s+", " ", clean).strip()


def _decision_tree_split_features(artifact: dict[str, Any]) -> list[str]:
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    nodes, edges = _tree_nodes_edges(artifact)
    source_ids = {_stringify(edge.get("source") or edge.get("from") or edge.get("parent")) for edge in edges}
    features: list[str] = []
    seen: set[str] = set()
    for node in nodes:
        node_id = _stringify(node.get("id") or node.get("node_id") or node.get("name"))
        is_split = _stringify(node.get("type")).lower() == "split" or node_id in source_ids
        if not is_split:
            continue
        feature = clean_tree_text(node.get("feature"))
        if not feature:
            label = clean_tree_text(node.get("label") or node.get("rule") or node.get("condition"))
            for marker in ("<=", ">=", "<", ">", "="):
                if marker in label:
                    feature = label.split(marker, 1)[0]
                    break
        feature = clean_tree_text(feature).replace("_", " ")
        if feature and feature.lower() not in seen:
            seen.add(feature.lower())
            features.append(feature)
        if len(features) >= 3:
            break
    if not features and data.get("target"):
        return [_stringify(data.get("target"))]
    return features


def _add_pdf_decision_tree_diagram(story: list[Any], artifact: dict[str, Any], styles: Any) -> None:
    if _add_pdf_decision_tree_image(story, artifact, styles):
        return
    nodes, edges = _tree_nodes_edges(artifact)
    if not nodes:
        _add_pdf_decision_tree_image(story, artifact, styles)
        return
    title = _stringify(artifact.get("title") or "Decision tree model rules")
    narrative = _decision_tree_narrative(artifact)
    story.append(Paragraph(_safe_paragraph(title), styles["ReportFigureHeading"]))
    if narrative:
        story.append(Paragraph(_safe_paragraph(narrative), styles["ReportCaption"]))
    for callout in _decision_tree_callouts(artifact):
        story.append(Paragraph(_safe_paragraph(f"Top rule: {callout}"), styles["ReportCaption"]))
    width = 6.2 * inch
    height = 3.35 * inch
    drawing = Drawing(width, height)
    layout = build_tree_layout(nodes, edges)
    layout_nodes = layout.get("nodes", [])
    layout_edges = layout.get("edges", [])
    normalized_positions = layout.get("positions", {})
    if not layout_nodes or not normalized_positions:
        return
    max_depth = int(layout.get("max_depth", 0) or 0)
    node_width = min(max(width / max(_pdf_tree_leaf_count(layout_nodes, layout_edges), 2) - 16, 78), 128)
    node_height = 36 if max_depth <= 2 else 31
    positions: dict[str, tuple[float, float, float, float]] = {}
    left_pad = 10
    right_pad = 10
    top_pad = 18
    bottom_pad = 28
    usable_width = width - left_pad - right_pad
    usable_height = height - top_pad - bottom_pad
    for node in layout_nodes:
        node_id = _stringify(node.get("_layout_id", ""))
        if node_id not in normalized_positions:
            continue
        x_norm, y_norm = normalized_positions[node_id]
        x = left_pad + (usable_width - node_width) * float(x_norm)
        row_top = top_pad + (usable_height - node_height) * float(y_norm)
        y = height - row_top - node_height
        positions[node_id] = (x, y, node_width, node_height)

    split_node_ids = {edge["source"] for edge in layout_edges}
    line_color = rl_colors.HexColor("#7A9BC3")
    blue = rl_colors.HexColor("#1F4E79")
    for edge in layout_edges:
        source = edge["source"]
        target = edge["target"]
        if source not in positions or target not in positions:
            continue
        sx, sy, sw, sh = positions[source]
        tx, ty, tw, th = positions[target]
        start_x = sx + sw / 2
        start_y = sy
        end_x = tx + tw / 2
        end_y = ty + th
        joint_y = start_y - max((start_y - end_y) * 0.42, 10)
        drawing.add(Line(start_x, start_y, start_x, joint_y, strokeColor=line_color, strokeWidth=1.0))
        drawing.add(Line(start_x, joint_y, end_x, joint_y, strokeColor=line_color, strokeWidth=1.0))
        drawing.add(Line(end_x, joint_y, end_x, end_y, strokeColor=line_color, strokeWidth=1.0))
        label = _stringify(edge.get("label")).strip()
        if label:
            drawing.add(
                String(
                    (start_x + end_x) / 2,
                    joint_y + 4,
                    label[:18],
                    fontSize=6.5,
                    fillColor=blue,
                    textAnchor="middle",
                )
            )
    for node in layout_nodes:
        node_id = _stringify(node.get("_layout_id", ""))
        if node_id not in positions:
            continue
        x, y, node_width, node_height = positions[node_id]
        is_split = node_id in split_node_ids
        fill = blue if is_split else rl_colors.HexColor("#F7F9FC")
        stroke = blue
        text_color = rl_colors.white if is_split else rl_colors.HexColor("#16324F")
        drawing.add(Rect(x, y, node_width, node_height, fillColor=fill, strokeColor=stroke, strokeWidth=1))
        label = _pdf_decision_tree_node_label(node, is_split)
        wrapped = textwrap.wrap(label, width=24)[:3] or [""]
        start_y = y + node_height - 11
        for line_index, line in enumerate(wrapped):
            drawing.add(
                String(
                    x + node_width / 2,
                    start_y - line_index * 9,
                    line,
                    fontSize=6.2,
                    fillColor=text_color,
                    textAnchor="middle",
                )
            )
    story.append(drawing)
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    metric_bits = []
    for key in _decision_tree_metric_keys():
        if data.get(key) not in (None, ""):
            metric_bits.append(f"{_decision_tree_metric_label(key, data)}: {_decision_tree_metric_value(key, data.get(key))}")
    if metric_bits:
        story.append(Paragraph(_safe_paragraph(" | ".join(metric_bits)), styles["ReportCaption"]))
    story.append(Spacer(1, 0.12 * inch))


def _decision_tree_image_path(artifact: dict[str, Any]) -> str:
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    for candidate in (
        artifact.get("image_path"),
        artifact.get("fallback_path"),
        artifact.get("visual_path"),
        data.get("image_path"),
        data.get("fallback_path"),
        data.get("figure_path"),
    ):
        path = _stringify(candidate)
        if path and os.path.exists(path):
            return path
    return ""


def _add_pdf_decision_tree_image(story: list[Any], artifact: dict[str, Any], styles: Any) -> bool:
    image_path = _decision_tree_image_path(artifact)
    if not image_path:
        return False
    from reportlab.platypus import Image as RLImage

    title = _stringify(artifact.get("title") or "Decision tree model rules")
    narrative = _decision_tree_narrative(artifact)
    story.append(Paragraph(_safe_paragraph(title), styles["ReportFigureHeading"]))
    image_width, image_height = _fit_image_size(image_path, 6.2 * inch, 3.6 * inch)
    story.append(RLImage(image_path, width=image_width, height=image_height))
    if narrative:
        story.append(Paragraph(_safe_paragraph(f"Figure note: {narrative}"), styles["ReportCaption"]))
    for callout in _decision_tree_callouts(artifact):
        story.append(Paragraph(_safe_paragraph(f"Top rule: {callout}"), styles["ReportCaption"]))
    data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
    metric_bits = []
    for key in _decision_tree_metric_keys():
        value = data.get(key) or artifact.get(key)
        if value not in (None, ""):
            metric_bits.append(f"{_decision_tree_metric_label(key, data)}: {_decision_tree_metric_value(key, value)}")
    if metric_bits:
        story.append(Paragraph(_safe_paragraph(" | ".join(metric_bits)), styles["ReportCaption"]))
    story.append(Spacer(1, 0.12 * inch))
    return True


def _decision_tree_metric_keys() -> tuple[str, ...]:
    return (
        "train_accuracy",
        "test_accuracy",
        "train_score",
        "test_score",
        "accuracy",
        "baseline_accuracy",
        "baseline_score",
        "precision",
        "recall",
        "f1",
        "f1_score",
        "support",
        "positive_class_rate",
        "train_r2",
        "test_r2",
        "r2",
        "train_mae",
        "test_mae",
        "mae",
    )


def _decision_tree_metric_label(key: str, data: dict[str, Any]) -> str:
    model_type = _stringify(data.get("model_type", "")).lower()
    if key == "train_score" and not model_type.startswith("reg"):
        return "Train Accuracy"
    if key == "test_score" and not model_type.startswith("reg"):
        return "Test Accuracy"
    if key == "baseline_score" and not model_type.startswith("reg"):
        return "Baseline Accuracy"
    return key.replace("_", " ").title()


def _decision_tree_metric_value(key: str, value: Any) -> str:
    if key in {"train_score", "test_score", "baseline_score", "train_accuracy", "test_accuracy", "baseline_accuracy", "accuracy", "precision", "recall", "f1", "f1_score", "positive_class_rate"}:
        try:
            number = float(value)
        except (TypeError, ValueError):
            return _stringify(value)
        if 0 <= number <= 1:
            return f"{number * 100:.1f}%"
    return _stringify(value)


def _pdf_tree_leaf_count(nodes: list[dict[str, Any]], edges: list[dict[str, Any]]) -> int:
    sources = {edge.get("source") for edge in edges}
    leaves = [node for node in nodes if node.get("_layout_id") not in sources]
    return max(len(leaves), 1)


def _pdf_decision_tree_node_label(node: dict[str, Any], is_split: bool) -> str:
    display = clean_tree_text(node.get("display_label") or node.get("human_label"))
    if display:
        return display
    if is_split:
        feature = clean_tree_text(node.get("feature"))
        threshold = clean_tree_text(node.get("threshold"))
        if feature and threshold:
            return f"{feature} <= {threshold}"
    else:
        prediction = clean_tree_text(node.get("prediction") or node.get("class") or node.get("value"))
        samples = clean_tree_text(node.get("samples") or node.get("n") or node.get("support"))
        if prediction:
            return f"Predict {prediction}" + (f" n={samples}" if samples else "")
    for key in ("label", "rule", "name", "condition", "prediction", "value", "class"):
        text = clean_tree_text(node.get(key))
        if text:
            return text
    return "Split node" if is_split else "Leaf"


def _eda_model_transition_text() -> str:
    return (
        "The EDA figures above show observed dataset patterns. The decision tree below is a separate "
        "explanatory model check: use it to audit possible follow-up segments, not as a production screening tool."
    )


def _figure_caption_for(figure: str, figure_captions: dict[str, Any]) -> str:
    if not isinstance(figure_captions, dict):
        return ""
    candidates = [
        figure,
        str(figure),
        os.path.basename(str(figure)),
        Path(str(figure)).name,
        Path(str(figure)).stem,
    ]
    for key in candidates:
        value = figure_captions.get(key)
        if _stringify(value).strip():
            return _stringify(value).strip()
    return ""


def _report_figures(workflow_state: dict[str, Any], limit: int = 4) -> list[str]:
    analysis_results = workflow_state.get("analysis_results", {}) or {}
    tree_artifacts = _decision_tree_artifacts(analysis_results)
    figures = [
        _stringify(figure)
        for figure in workflow_state.get("saved_figures", [])
        if _stringify(figure) and not (tree_artifacts and _is_decision_tree_figure(_stringify(figure)))
    ]
    captions = analysis_results.get("figure_captions", {}) if isinstance(analysis_results, dict) else {}
    if _is_stock_time_series_report_context(workflow_state):
        ranked = sorted(figures, key=lambda figure: _stock_figure_priority(figure, _figure_caption_for(figure, captions)))
        return ranked[:limit]
    if not _is_lending_report_context(workflow_state):
        return figures[:limit]
    ranked = sorted(figures, key=lambda figure: _credit_figure_priority(figure, _figure_caption_for(figure, captions)))
    selected = [figure for figure in ranked if _credit_figure_priority(figure, _figure_caption_for(figure, captions)) < 20]
    if len(selected) < limit:
        selected.extend(figure for figure in ranked if figure not in selected)
    return selected[:limit]


def _is_decision_tree_figure(path: str) -> bool:
    return Path(_stringify(path)).name.lower().startswith("decision_tree")


def _credit_figure_priority(figure: str, caption: str) -> int:
    text = f"{figure} {caption}".lower()
    if "decision_tree" in text:
        return 50
    preferred = (
        "loan_grade",
        "loan grade",
        "debt burden",
        "loan_percent_income",
        "percent income",
        "interest",
        "loan_int_rate",
        "income band",
        "home ownership",
        "default",
    )
    weak = ("person_age", "borrower age", " age ", "scatter")
    if any(token in text for token in preferred):
        return 0
    if any(token in text for token in weak):
        return 25
    return 10


def _stock_figure_priority(figure: str, caption: str) -> int:
    text = f"{figure} {caption}".lower()
    if any(token in text for token in ("price trend", "moving average", "support", "resistance")):
        return 0
    if any(token in text for token in ("volatility", "drawdown", "seasonal", "return distribution")):
        return 1
    if any(token in text for token in ("volume spike", "volume trend", "price-volume")):
        return 2
    if any(token in text for token in ("scatter", "volume-return", "correlation")):
        return 25
    return 10


def _format_slide_analysis_findings(analysis_results: dict[str, Any]) -> list[str]:
    findings: list[str] = []
    seen: set[str] = set()

    def add_finding(text: str) -> None:
        clean = _stringify(text).strip()
        if not clean or clean in seen:
            return
        seen.add(clean)
        findings.append(clean)

    for item in analysis_results.get("business_findings", [])[:4]:
        add_finding(item)

    summary = analysis_results.get("analysis_summary", {})
    if isinstance(summary, dict):
        for key, value in list(summary.items())[:6]:
            add_finding(f"{str(key).replace('_', ' ').title()}: {_stringify(value)}")

    if not findings:
        add_finding("The analysis code executed, but no structured analysis findings were captured for slides.")

    return findings


def _format_recommendations(recommendations: list[dict[str, Any]]) -> list[str]:
    formatted = []
    caveats = [
        "monitor approval-rate and fairness impact before scaling",
        "retain human review before adverse decisions",
        "recalibrate if the portfolio mix or macro conditions shift",
        "confirm servicing capacity before expanding outreach",
        "check drift against recent booked-loan performance",
    ]
    for index, item in enumerate(recommendations[:5]):
        if not isinstance(item, dict):
            text = _stringify(item)
            if text:
                formatted.append(_recommendation_hypothesis(text))
            continue
        action = item.get("action", "")
        rationale = item.get("rationale", "")
        evidence = item.get("evidence", "")
        target = item.get("target_segment") or item.get("segment") or _recommendation_target(index, action, evidence)
        owner = item.get("owner") or item.get("accountable_owner") or _recommendation_owner(action, evidence)
        trigger = item.get("trigger") or item.get("threshold") or _recommendation_trigger(action, evidence)
        timeline = item.get("timeline") or item.get("timeframe") or "next review cycle"
        impact = item.get("expected_impact") or item.get("impact") or ""
        risk = item.get("risk") or item.get("guardrail") or ""
        metric = item.get("validation_metric") or item.get("success_metric") or _recommendation_metric(action, evidence)
        caveat = item.get("caveat") or caveats[index % len(caveats)]
        text = f"Decision: pilot {action or 'the recommended policy change'}."
        if owner:
            text += f" Owner: {owner}."
        if trigger:
            text += f" Trigger: {trigger}."
        if target:
            text += f" Target: {target}."
        if timeline:
            text += f" Timeline: {timeline}."
        if rationale:
            text += f" Rationale: {rationale}."
        if evidence:
            text += f" Evidence: {evidence}."
        if impact:
            text += f" Expected impact: {impact}."
        text += f" Guardrail: {metric}. Risk: {risk or caveat}."
        formatted.append(text)
    return formatted


def _format_executive_recommendations(recommendations: list[Any]) -> list[str]:
    formatted: list[str] = []
    for fallback, item in enumerate(recommendations[:4], start=1):
        if not isinstance(item, dict):
            text = _stringify(item)
            if text:
                formatted.append(text)
            continue
        rank = item.get("rank") or fallback
        action = _short_report_text(item.get("action"), 260)
        owner = _short_report_text(item.get("owner"), 90)
        timeline = _short_report_text(item.get("timeline"), 80)
        metric = _short_report_text(item.get("validation_metric"), 150)
        stop = _short_report_text(item.get("stop_condition"), 150)
        evidence_ids = ", ".join(str(value) for value in item.get("evidence_ids", []) or [] if value)
        parts = [f"Priority {rank}: {action}"]
        if owner or timeline:
            parts.append(f"Owner/timing: {owner or 'named owner'}; {timeline or 'validation cycle'}")
        if evidence_ids:
            parts.append(f"Evidence: {evidence_ids}")
        if metric:
            parts.append(f"Measure: {metric}")
        if stop:
            parts.append(f"Stop condition: {stop}")
        formatted.append(". ".join(part.rstrip(".") for part in parts if part) + ".")
    return formatted


def _short_report_text(value: Any, limit: int) -> str:
    text = _stringify(value).strip()
    if len(text) <= limit:
        return text
    trimmed = text[: limit + 1].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return trimmed + "."


def _recommendation_owner(action: Any, evidence: Any) -> str:
    text = f"{_stringify(action)} {_stringify(evidence)}".lower()
    if any(token in text for token in ("stock", "price", "volatility", "drawdown", "return", "volume", "hedg")):
        return "investment lead"
    if any(token in text for token in ("loan", "default", "credit", "income", "grade", "interest", "debt")):
        return "risk owner"
    return "business owner"


def _recommendation_trigger(action: Any, evidence: Any) -> str:
    text = f"{_stringify(action)} {_stringify(evidence)}".lower()
    if any(token in text for token in ("stock", "price", "volatility", "drawdown", "return", "volume", "hedg")):
        return "price, volatility, volume, or drawdown signal crosses the review threshold"
    if any(token in text for token in ("loan", "default", "credit", "income", "grade", "interest", "debt")):
        return "segment crosses the agreed risk threshold"
    return "segment enters the agreed intervention threshold"


def _recommendation_target(index: int, action: Any, evidence: Any) -> str:
    text = f"{_stringify(action)} {_stringify(evidence)}".lower()
    if any(token in text for token in ("stock", "price", "volatility", "drawdown", "return", "volume", "hedg")):
        if "drawdown" in text or "volatility" in text:
            return "high-volatility or drawdown windows"
        if "volume" in text or "liquidity" in text:
            return "volume and liquidity spike periods"
        targets = [
            "price breakouts near support or resistance",
            "high-volatility or drawdown windows",
            "volume and liquidity spike periods",
            "forward return validation windows",
        ]
        return targets[index % len(targets)]
    if any(token in text for token in ("loan", "default", "credit", "income", "grade", "interest", "debt")):
        targets = [
            "high debt-burden applicants",
            "weaker-grade or high-interest applications",
            "prior-default or low-income review segments",
            "booked loans flagged for early collections outreach",
        ]
        return targets[index % len(targets)]
    return "the segment named by the strongest validated evidence"


def _recommendation_metric(action: Any, evidence: Any) -> str:
    text = f"{_stringify(action)} {_stringify(evidence)}".lower()
    if any(token in text for token in ("stock", "price", "volatility", "drawdown", "return", "volume", "hedg")):
        return "forward return, max drawdown, rolling volatility, volume/liquidity, and review cadence"
    if any(token in text for token in ("loan", "default", "credit", "income", "grade", "interest", "debt")):
        return "default capture, false-positive review rate, approval impact, and fairness checks"
    return "precision, recall, F1/support, adoption, and outcome lift against a holdout baseline"


def _format_executive_decision_block(workflow_state: dict[str, Any]) -> list[str]:
    outputs = workflow_state.get("agent_outputs", {}) or {}
    decision = outputs.get("decision_maker", {}) or {}
    analysis = workflow_state.get("analysis_results", {}) or {}
    recommendations = decision.get("recommendations", []) or []
    first_action = ""
    if recommendations:
        item = recommendations[0]
        first_action = _stringify(item.get("action") if isinstance(item, dict) else item)
    decision_text = _stringify(decision.get("final_recommendation") or first_action or decision.get("executive_summary"))
    artifacts = [
        artifact
        for artifact in analysis.get("analysis_artifacts", []) or []
        if isinstance(artifact, dict) and _stringify(artifact.get("chart_type")).lower() != "decision_tree"
    ]
    evidence = _stringify((artifacts[0].get("finding") or artifacts[0].get("takeaway") or artifacts[0].get("title")) if artifacts else "")
    caveat = _format_limitations(workflow_state)[0]
    target = _stringify(workflow_state.get("decision_tree_target_column", "")).replace("_", " ")
    if bool((workflow_state.get("evidence_bundle", {}) or {}).get("high_stakes")):
        return [
            f"Decision: {decision_text or 'Use the findings only to design a professionally supervised validation study.'}",
            f"Population: records in the supplied dataset related to {target or 'the stated outcome'}; broader representativeness is not established.",
            f"Evidence: {evidence or 'Use the cited EDA artifacts and exploratory model branches as hypothesis-generating evidence.'}",
            "Validation: report class-sensitive metrics, subgroup stability, fairness, safety, and expert-reviewed outcomes on fresh data.",
            "Boundary: no diagnosis, eligibility, automated adverse action, or high-impact decision without domain-expert and human review.",
        ]
    if _is_lending_report_context(workflow_state):
        return [
            f"Decision: {decision_text or 'Run bounded underwriting and portfolio-monitoring validation pilots before policy changes.'}",
            "Who affected: applicants or booked-loan segments with high debt burden, weaker grades, high interest rates, low income, or prior-default flags.",
            f"Evidence: {evidence or 'Use binned credit-risk visuals and verified tree paths as decision-support evidence.'}",
            "Pilot metric: track default capture, false-positive review rate, approval impact, and fairness checks.",
            "Caveat: explanatory screening only; not diagnostic, not automated adverse action, and not deployment-ready without governance review.",
        ]
    if _is_stock_time_series_report_context(workflow_state):
        return [
            f"Decision: {decision_text or 'Use price, volatility, drawdown, and volume triggers as validation gates before portfolio action.'}",
            "Who affected: portfolio, treasury, or investment stakeholders monitoring STC price and liquidity conditions.",
            f"Evidence: {evidence or 'Use the time-series price, volatility, return, and volume visuals as decision-support evidence.'}",
            "Pilot metric: track forward return, max drawdown, rolling volatility, volume/liquidity, and review cadence.",
            f"Caveat: {_format_limitations(workflow_state)[0]}",
        ]
    return [
        f"Decision: {decision_text or 'Prioritize a bounded validation pilot before scaling policy changes.'}",
        f"Who affected: records or segments tied to the strongest observed {target or 'outcome'} risk signals.",
        f"Evidence: {evidence or 'Use the top EDA signals and model checks as decision-support evidence.'}",
        "Pilot metric: track outcome lift, precision, recall, F1/support, and implementation quality where screening is used.",
        f"Caveat: {caveat}",
    ]


def _is_lending_report_context(workflow_state: dict[str, Any]) -> bool:
    target = _stringify(workflow_state.get("decision_tree_target_column")).lower()
    description = _stringify(workflow_state.get("user_data_description")).lower()
    outputs = workflow_state.get("agent_outputs", {}) or {}
    data_text = _stringify(outputs.get("data_understander", {})).lower()
    return any(token in f"{target} {description} {data_text}" for token in ("loan_status", "loan status", "loan_", "credit", "default", "underwriting"))


def _is_stock_time_series_report_context(workflow_state: dict[str, Any]) -> bool:
    description = _stringify(workflow_state.get("user_data_description")).lower()
    outputs = workflow_state.get("agent_outputs", {}) or {}
    data_text = _stringify(outputs.get("data_understander", {})).lower()
    fields: list[str] = []
    for df in (workflow_state.get("csv_data", {}) or {}).values():
        fields.extend(str(column) for column in getattr(df, "columns", [])[:30])
        break
    field_text = " ".join(fields).lower()
    description_text = f"{description} {data_text}"
    explicit_stock = bool(re.search(r"\b(stock|share price|ohlc|stc)\b", description_text))
    has_price_field = any(token in field_text for token in ("price", "open", "high", "low", "close"))
    has_market_field = any(token in field_text for token in ("vol.", "volume", "change %", "return"))
    return explicit_stock or (
        "date" in field_text and has_price_field and has_market_field
    )


def _format_top_risk_signals(analysis_results: dict[str, Any]) -> list[str]:
    signals: list[str] = []
    for artifact in analysis_results.get("analysis_artifacts", []) or []:
        if not isinstance(artifact, dict):
            continue
        if _stringify(artifact.get("chart_type")).lower() == "decision_tree":
            continue
        title = _stringify(artifact.get("title") or artifact.get("artifact_id"))
        evidence = _stringify(artifact.get("finding") or artifact.get("takeaway"))
        if not title and not evidence:
            continue
        implication = _risk_signal_implication(title, evidence)
        signals.append(
            f"Signal: {title or evidence}. Evidence: {evidence or 'visible in the chart.'} "
            f"Decision implication: {implication} Caveat: validate thresholds before policy changes."
        )
        if len(signals) >= 4:
            break
    return signals


def _risk_signal_implication(title: str, evidence: str) -> str:
    text = f"{title} {evidence}".lower()
    if any(token in text for token in ("loan", "default", "credit", "income", "grade", "interest")):
        return "use as an underwriting, pricing, documentation, or portfolio-monitoring pilot input."
    if any(token in text for token in ("attrition", "employee", "overtime", "jobrole")):
        return "use as a retention pilot target with manager review."
    if any(token in text for token in ("depression", "mental", "stress", "anxiety")):
        return "use as a support-screening pilot input with human review."
    return "use as a prioritized validation segment."


def _recommendation_hypothesis(text: str) -> str:
    clean = _stringify(text)
    if not clean:
        return ""
    if clean.lower().startswith("pilot hypothesis"):
        return clean
    return f"Pilot hypothesis: {clean}. Validate with explicit success criteria before scaling."


def _format_priority_findings(business_translator: dict[str, Any]) -> list[str]:
    formatted: list[str] = []
    for item in business_translator.get("key_findings", [])[:5]:
        if not isinstance(item, dict):
            text = _stringify(item)
            if text:
                formatted.append(text)
            continue
        finding = _stringify(item.get("finding", ""))
        implication = _stringify(item.get("business_implication", ""))
        priority = _stringify(item.get("priority", ""))
        text = finding
        if implication:
            text += f" Business implication: {implication}."
        if priority:
            text += f" Priority: {priority}."
        if text:
            formatted.append(text)
    return formatted


def _format_workflow_trace(workflow_state: dict[str, Any]) -> list[str]:
    manifest = workflow_state.get("run_manifest", {}) or {}
    datasets = manifest.get("datasets", []) or []
    figures = manifest.get("figures", []) or workflow_state.get("saved_figures", []) or []
    warnings = manifest.get("warnings", []) or workflow_state.get("analysis_artifact_warnings", []) or []
    agent_outputs = manifest.get("agent_outputs", []) or sorted((workflow_state.get("agent_outputs", {}) or {}).keys())
    trace = [
        f"Workflow status: {_stringify(workflow_state.get('status', 'unknown')) or 'unknown'}.",
        f"Datasets loaded: {len(datasets)}. Visual figures captured: {len(figures)}.",
        f"Agent outputs captured: {', '.join(str(item) for item in agent_outputs[:10]) or 'none recorded'}.",
    ]
    analysis_iterations = manifest.get("analysis_loop_iterations")
    if analysis_iterations:
        trace.append(f"Analysis coder-review loop iterations used: {analysis_iterations}.")
    if manifest.get("final_code_present"):
        trace.append("Final analysis code was approved for downstream reporting.")
    if warnings:
        trace.append(f"Artifact or workflow warnings recorded: {len(warnings)}.")
    return trace


def _format_data_quality_notes(data_understander: dict[str, Any]) -> list[str]:
    notes: list[str] = []
    for dataset_name, dataset_info in (data_understander.get("datasets", {}) or {}).items():
        if not isinstance(dataset_info, dict):
            text = _stringify(dataset_info)
            if text:
                notes.append(f"{dataset_name}: {text}")
            continue
        labels = {
            "cleaning_priorities": "Cleaning priority",
            "type_notes": "Field interpretation",
            "outlier_notes": "Range and outlier check",
        }
        for key in ("cleaning_priorities", "type_notes", "outlier_notes"):
            for note in dataset_info.get(key, [])[:2]:
                text = _stringify(note)
                if text:
                    notes.append(f"{labels[key]}: {text}")
    return notes or ["No material data-quality notes were recorded by the data understanding step."]


def _format_limitations(workflow_state: dict[str, Any]) -> list[str]:
    outputs = workflow_state.get("agent_outputs", {}) or {}
    decision = outputs.get("decision_maker", {}) or {}
    objective = workflow_state.get("workflow_objective", {}) or {}
    items: list[str] = []
    for limitation in decision.get("limitations", [])[:5]:
        if isinstance(limitation, dict):
            text = _stringify(limitation.get("limitation", ""))
            mitigation = _stringify(limitation.get("mitigation", ""))
            if text:
                items.append(f"{text} Mitigation: {mitigation or 'validate before scaling the decision.'}")
        else:
            text = _stringify(limitation)
            if text:
                items.append(text)
    for limitation in objective.get("limitations", [])[:3] if isinstance(objective, dict) else []:
        text = _stringify(limitation)
        if text:
            items.append(text)
    for warning in workflow_state.get("analysis_artifact_warnings", [])[:3]:
        text = _stringify(warning)
        if text:
            items.append(f"Analysis artifact warning: {text}")
    return items or ["No explicit decision limitations were recorded; validate conclusions against business context before action."]


def _fit_image_size(path: str, max_width: float, max_height: float) -> tuple[float, float]:
    if PILImage is None:
        return max_width, max_height
    try:
        with PILImage.open(path) as image:
            width, height = image.size
    except OSError:
        return max_width, max_height
    if not width or not height:
        return max_width, max_height
    ratio = min(max_width / width, max_height / height)
    return width * ratio, height * ratio


def _title_from_caption(caption: str, fallback_stem: str) -> str:
    text = (caption or "").strip()
    if not text:
        return fallback_stem.replace("_", " ").title()
    sentence = re.split(r"(?<=[.!?])\s+", text, maxsplit=1)[0]
    words = sentence.split()
    if len(words) > 9:
        sentence = " ".join(words[:9]).rstrip(",;:") + "..."
    return sentence.rstrip(".")


def _expand_slides_with_visuals(
    slides: list[dict[str, Any]],
    saved_figures: list[str],
    figure_captions: dict[str, str],
) -> list[dict[str, Any]]:
    expanded = [dict(slide) for slide in slides]
    used_visuals = {
        slide.get("visual_element", "")
        for slide in expanded
        if slide.get("visual_element") and os.path.exists(slide.get("visual_element", ""))
    }
    next_number = max([slide.get("slide_number", 0) for slide in expanded], default=0) + 1

    for figure in saved_figures:
        if not os.path.exists(figure) or figure in used_visuals:
            continue
        caption = figure_captions.get(figure, "").strip()
        title = _title_from_caption(caption, Path(figure).stem)
        expanded.append(
            {
                "slide_number": next_number,
                "title": title,
                "main_message": caption or f"Evidence from {Path(figure).stem.replace('_', ' ')}",
                "details": [],
                "visual_element": figure,
            }
        )
        next_number += 1

    return expanded


def _ensure_analysis_findings_slide(
    slides: list[dict[str, Any]],
    analysis_results: dict[str, Any],
) -> list[dict[str, Any]]:
    formatted_findings = _format_slide_analysis_findings(analysis_results)
    if not formatted_findings:
        return [dict(slide) for slide in slides]

    normalized_slides = [dict(slide) for slide in slides]
    analysis_keywords = ("analysis", "finding", "insight", "visual")

    for slide in normalized_slides:
        title = _stringify(slide.get("title", "")).lower()
        if any(keyword in title for keyword in analysis_keywords):
            existing_details = [str(item).strip() for item in (slide.get("details", []) or []) if str(item).strip()]
            merged_details: list[str] = []
            seen: set[str] = set()
            for item in existing_details + formatted_findings[1:5]:
                clean = _stringify(item).strip()
                if not clean or clean in seen:
                    continue
                seen.add(clean)
                merged_details.append(clean)
            if not _stringify(slide.get("main_message", "")).strip():
                slide["main_message"] = formatted_findings[0]
            slide["details"] = merged_details[:5]
            return _renumber_slides(normalized_slides)

    insert_at = 1 if len(normalized_slides) >= 1 else 0
    normalized_slides.insert(
        insert_at,
        {
            "slide_number": 0,
            "title": "Technical Analysis Findings",
            "main_message": formatted_findings[0],
            "details": formatted_findings[1:5],
            "visual_element": "",
        },
    )
    return _renumber_slides(normalized_slides)


def _ensure_objective_slide(slides: list[dict[str, Any]], workflow_state: dict[str, Any]) -> list[dict[str, Any]]:
    coverage = _format_objective_coverage(workflow_state)
    normalized = [dict(slide) for slide in slides]
    if any("objective" in _stringify(slide.get("title", "")).lower() for slide in normalized):
        return normalized
    normalized.insert(
        0,
        {
            "slide_number": 0,
            "title": "Objective Coverage",
            "main_message": coverage[0],
            "details": coverage[1:5],
            "visual_element": "",
        },
    )
    return _renumber_slides(normalized)


def _renumber_slides(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    renumbered: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        updated = dict(slide)
        updated["slide_number"] = index
        renumbered.append(updated)
    return renumbered


def _preferred_output_path(output_path: str) -> str:
    path = Path(output_path)
    if not path.exists():
        return str(path)
    try:
        with open(path, "ab"):
            return str(path)
    except OSError:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return str(path.with_name(f"{path.stem}_{timestamp}{path.suffix}"))


def _write_json_artifact(output_path: str, filename: str, payload: dict[str, Any]) -> str:
    artifact_path = Path(output_path).with_name(filename)
    artifact_path.write_text(json.dumps(payload, indent=2, default=str), encoding="utf-8")
    return str(artifact_path)


def _report_outline_payload(workflow_state: dict[str, Any], report_path: str) -> dict[str, Any]:
    outputs = workflow_state.get("agent_outputs", {}) or {}
    analysis_results = workflow_state.get("analysis_results", {}) or {}
    recommendations = outputs.get("decision_maker", {}).get("recommendations", []) or []
    market_claims = _market_claim_pairs(outputs.get("market_researcher", {}))
    tree_artifacts = _decision_tree_artifacts(analysis_results)
    decision_tree_count = len(tree_artifacts)
    decision_tree_notes = []
    for artifact in tree_artifacts:
        data = artifact.get("data", {}) if isinstance(artifact.get("data"), dict) else {}
        note = _stringify(data.get("performance_note") or decision_tree_performance_note(data))
        if note:
            decision_tree_notes.append(note)
    decision_tree_notes = [note for note in decision_tree_notes if note]
    sections = [
        {"title": "Executive Summary", "source_agents": ["decision_maker"], "item_count": 2},
        {"title": "Objective Coverage", "source_agents": ["workflow_objective"], "item_count": len(_format_objective_coverage(workflow_state))},
        {"title": "Dataset Overview / Data Understanding", "source_agents": ["data_understander"], "item_count": 2},
        {"title": "Market Research", "source_agents": ["market_researcher"], "item_count": len(market_claims)},
        {"title": "Analysis Plan", "source_agents": ["planner"], "item_count": len(outputs.get("planner", {}).get("objectives", []) or [])},
        {"title": "Data Analysis and Visual Findings", "source_agents": ["data_scientist_coder"], "item_count": len(_format_analysis_findings(analysis_results))},
    ]
    if decision_tree_count:
        sections.append({"title": "Decision Tree Model", "source_agents": ["data_scientist_coder"], "item_count": decision_tree_count})
    sections.extend(
        [
            {"title": "Business Translation", "source_agents": ["business_translator"], "item_count": len(_format_priority_findings(outputs.get("business_translator", {})))},
            {"title": "Decision Recommendations", "source_agents": ["decision_maker"], "item_count": len(recommendations)},
            {"title": "Limitations and Validation Steps", "source_agents": ["workflow_objective", "decision_maker", "business_translator"], "item_count": len(_format_limitations(workflow_state))},
        ]
    )
    return {
        "artifact_type": "report_outline",
        "report_path": report_path,
        "evidence_bundle_hash": (workflow_state.get("evidence_bundle", {}) or {}).get("bundle_hash", ""),
        "quality_status": (workflow_state.get("quality_receipt", {}) or {}).get("status", ""),
        "generated_on": datetime.now().isoformat(timespec="seconds"),
        "objective_coverage": _format_objective_coverage(workflow_state),
        "sections": sections,
        "analysis_visuals": [
            {
                "path": str(path),
                "caption": _figure_caption_for(str(path), analysis_results.get("figure_captions", {})),
            }
            for path in workflow_state.get("saved_figures", [])[:6]
        ],
        "recommendation_count": len(recommendations),
        "decision_tree_artifact_count": decision_tree_count,
        "decision_tree_performance_notes": decision_tree_notes[:2],
        "structured_chart_count": len(analysis_results.get("analysis_artifacts", []) or analysis_results.get("chart_specs", []) or []),
    }


def generate_pdf_report(workflow_state: dict[str, Any], output_path: str = "analytics_report.pdf") -> str:
    outputs = workflow_state.get("agent_outputs", {})
    analysis_results = workflow_state.get("analysis_results", {})
    user_description = _clean_objective_display(
        _stringify(workflow_state.get("user_data_description", "")).strip(),
        workflow_state,
    )
    resolved_output_path = _preferred_output_path(output_path)
    if REPORTLAB_AVAILABLE:
        styles = getSampleStyleSheet()
        report_title = _report_title(workflow_state)
        styles.add(ParagraphStyle(name="CoverKicker", parent=styles["Normal"], textColor=rl_colors.HexColor("#D9E7F2"), fontSize=9, leading=12, spaceAfter=18))
        styles.add(ParagraphStyle(name="CoverTitle", parent=styles["Title"], textColor=rl_colors.white, fontName="Helvetica-Bold", fontSize=26, leading=31, alignment=0, spaceAfter=18))
        styles.add(ParagraphStyle(name="CoverSub", parent=styles["BodyText"], textColor=rl_colors.HexColor("#D9E7F2"), fontSize=11, leading=16, spaceAfter=8))
        styles.add(ParagraphStyle(name="ReportHeading", parent=styles["Heading2"], textColor=rl_colors.HexColor("#16324F"), fontName="Helvetica-Bold", fontSize=16, leading=20, spaceBefore=12, spaceAfter=8))
        styles.add(ParagraphStyle(name="ReportFigureHeading", parent=styles["Heading3"], textColor=rl_colors.HexColor("#234A68"), fontName="Helvetica-Bold", fontSize=11, leading=14, spaceBefore=8, spaceAfter=5))
        styles.add(ParagraphStyle(name="ReportBody", parent=styles["BodyText"], textColor=rl_colors.HexColor("#243B53"), fontSize=10, leading=14, spaceAfter=6))
        styles.add(ParagraphStyle(name="ReportCaption", parent=styles["BodyText"], textColor=rl_colors.HexColor("#5B6570"), fontSize=8.5, leading=11, italic=True))
        styles.add(ParagraphStyle(name="EvidenceNote", parent=styles["BodyText"], backColor=rl_colors.HexColor("#EAF3F8"), borderColor=rl_colors.HexColor("#C9D9E4"), borderWidth=0.6, borderPadding=7, textColor=rl_colors.HexColor("#29485F"), fontSize=9, leading=12, spaceBefore=4, spaceAfter=8))

        story: list[Any] = [Spacer(1, 1.05 * inch)]
        story.append(Paragraph("MULTI-AGENT ANALYTICS SYSTEM", styles["CoverKicker"]))
        story.append(Paragraph(_safe_paragraph(report_title), styles["CoverTitle"]))
        story.append(Spacer(1, 0.18 * inch))
        story.append(Paragraph("Executive decision report", styles["CoverSub"]))
        if user_description:
            story.append(Paragraph(_safe_paragraph(user_description), styles["CoverSub"]))
        story.append(Spacer(1, 2.75 * inch))
        story.append(Paragraph(datetime.now().strftime("%B %d, %Y"), styles["CoverSub"]))
        story.append(PageBreak())

        _add_pdf_heading(story, "Executive Summary", styles["ReportHeading"])
        _add_pdf_label_table(
            story,
            _format_executive_decision_block(workflow_state)[:5],
            styles["ReportBody"],
            styles["ReportCaption"],
            fallback_label="Summary",
        )
        decision_output = outputs.get("decision_maker", {}) or {}
        if bool((workflow_state.get("evidence_bundle", {}) or {}).get("high_stakes")):
            summary = _stringify(decision_output.get("conclusion", ""))
        else:
            summary = _stringify(decision_output.get("executive_summary", ""))
        if summary:
            _add_pdf_body(story, summary, styles["ReportBody"])

        scope_rows = [[Paragraph(f"<b>{_safe_paragraph(label)}</b>", styles["ReportCaption"]), Paragraph(_safe_paragraph(value), styles["ReportBody"])] for label, value in _report_scope_metrics(workflow_state)]
        scope_table = Table(scope_rows, colWidths=[1.45 * inch, 4.75 * inch])
        scope_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), rl_colors.HexColor("#F3F7FA")),
            ("BOX", (0, 0), (-1, -1), 0.6, rl_colors.HexColor("#D3DFE7")),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, rl_colors.HexColor("#DDE6EC")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.extend([scope_table, Spacer(1, 0.14 * inch)])

        required_caveats = (workflow_state.get("quality_receipt", {}) or {}).get("required_caveats", []) or []
        if required_caveats:
            story.append(Paragraph(_safe_paragraph("Decision boundary: " + " ".join(required_caveats[:2])), styles["EvidenceNote"]))

        _add_pdf_heading(story, "What the data can support", styles["ReportHeading"])
        _add_pdf_body(story, outputs.get("data_understander", {}).get("executive_summary", ""), styles["ReportBody"])
        _add_pdf_label_table(
            story,
            _format_data_quality_notes(outputs.get("data_understander", {}))[:5],
            styles["ReportBody"],
            styles["ReportCaption"],
            fallback_label="Quality check",
        )

        _add_pdf_heading(story, "Key findings with visual evidence", styles["ReportHeading"])
        evidence_records = (workflow_state.get("evidence_bundle", {}) or {}).get("evidence", []) or []
        if not evidence_records:
            _add_pdf_numbered_items(story, _format_analysis_findings(analysis_results)[:6], styles["ReportBody"], noun="Finding")
        for evidence_index, item in enumerate(evidence_records, start=1):
            if not isinstance(item, dict) or item.get("kind") == "model":
                continue
            claim = _stringify(item.get("claim", ""))
            if claim:
                context_bits = []
                if item.get("sample_size"):
                    context_bits.append(f"n={item.get('sample_size')}")
                if item.get("method"):
                    context_bits.append(str(item.get("method")))
                context = f" ({'; '.join(context_bits)})" if context_bits else ""
                story.append(Paragraph(f"<b>[E{evidence_index}]</b> {_safe_paragraph(claim + context)}", styles["ReportBody"]))
        for figure in _report_figures(workflow_state):
            caption = _figure_caption_for(figure, analysis_results.get("figure_captions", {}))
            if os.path.exists(figure):
                from reportlab.platypus import Image as RLImage

                figure_title = _title_from_caption(caption, Path(figure).stem)
                image_width, image_height = _fit_image_size(figure, 6.2 * inch, 3.6 * inch)
                figure_group: list[Any] = [Paragraph(_safe_paragraph(figure_title), styles["ReportFigureHeading"]), RLImage(figure, width=image_width, height=image_height)]
                if caption:
                    figure_group.append(Paragraph(_safe_paragraph(caption), styles["ReportCaption"]))
                figure_group.append(Spacer(1, 0.12 * inch))
                story.append(KeepTogether(figure_group))
        tree_artifacts = _decision_tree_artifacts(analysis_results)
        if tree_artifacts:
            story.append(PageBreak())
            _add_pdf_heading(story, "Exploratory model check", styles["ReportHeading"])
            _add_pdf_body(story, _eda_model_transition_text(), styles["EvidenceNote"])
        for tree_artifact in tree_artifacts[:1]:
            _add_pdf_decision_tree_diagram(story, tree_artifact, styles)

        _add_pdf_heading(story, "Recommended next steps", styles["ReportHeading"])
        _add_pdf_body(story, outputs.get("decision_maker", {}).get("final_recommendation", ""), styles["ReportBody"])
        _add_pdf_numbered_items(
            story,
            _format_executive_recommendations(outputs.get("decision_maker", {}).get("recommendations", [])),
            styles["ReportBody"],
            noun="Priority",
        )

        _add_pdf_heading(story, "Further questions before scaling", styles["ReportHeading"])
        _add_pdf_numbered_items(
            story,
            [
                "Do the observed patterns hold in a fresh and appropriately sampled dataset?",
                "Which operational owner, safety guardrail, and success metric will govern the validation study?",
                "Do results remain stable across relevant subgroups and under class-sensitive evaluation?",
            ],
            styles["ReportBody"],
            noun="Question",
        )

        _add_pdf_heading(story, "Caveats and assumptions", styles["ReportHeading"])
        _add_pdf_caveat_table(
            story,
            _format_limitations(workflow_state),
            styles["ReportBody"],
            styles["ReportCaption"],
        )

        _add_pdf_heading(story, "Sources and reproducibility", styles["ReportHeading"])
        appendix_items: list[str] = []
        for source in (workflow_state.get("evidence_bundle", {}) or {}).get("datasets", []) or []:
            appendix_items.append(
                f"Dataset: {source.get('name', '')}; rows={source.get('rows', 'unknown')}; columns={source.get('columns', 'unknown')}; SHA-256={source.get('sha256') or 'not recorded'}."
                + (f" As of {source.get('as_of')}." if source.get("as_of") else "")
            )
        external_sources = (workflow_state.get("evidence_bundle", {}) or {}).get("external_sources", []) or []
        if external_sources:
            external_inventory: list[str] = []
            for source in external_sources:
                url = str(source.get("url") or "")
                domain_match = re.search(r"https?://(?:www\.)?([^/]+)", url)
                domain = domain_match.group(1) if domain_match else url
                external_inventory.append(
                    f"[{source.get('index')}] {source.get('title', '')} ({domain or 'source page'})"
                )
            appendix_items.append(
                "External context: "
                + "; ".join(external_inventory)
                + ". Search snippets are context only until the source pages are reviewed."
            )
        if not appendix_items:
            appendix_items = ["No external sources were captured; conclusions rely on the supplied dataset and saved analysis artifacts."]
        _add_pdf_label_table(
            story,
            appendix_items,
            styles["ReportBody"],
            styles["ReportCaption"],
            fallback_label="Source",
        )
        if not external_sources:
            _add_pdf_body(story, " ".join(_format_workflow_trace(workflow_state)), styles["ReportCaption"])

        doc = SimpleDocTemplate(
            resolved_output_path,
            pagesize=A4,
            rightMargin=0.72 * inch,
            leftMargin=0.72 * inch,
            topMargin=0.68 * inch,
            bottomMargin=0.62 * inch,
            title=report_title,
            author="Multi-Agent Analytics System",
        )
        _trim_trailing_spacers_and_pagebreaks(story)
        doc.build(
            story,
            onFirstPage=_pdf_page_decorator(report_title, cover=True),
            onLaterPages=_pdf_page_decorator(report_title),
        )
        pdf_qa = inspect_pdf_report(resolved_output_path)
        workflow_state["pdf_quality_receipt"] = pdf_qa
        if any(issue.get("severity") == "error" for issue in pdf_qa.get("issues", [])):
            raise RuntimeError(
                "PDF semantic QA failed: "
                + "; ".join(str(issue.get("issue")) for issue in pdf_qa["issues"][:8])
            )
        outline_path = _write_json_artifact(
            resolved_output_path,
            "report_outline.json",
            _report_outline_payload(workflow_state, resolved_output_path),
        )
        workflow_state.setdefault("generated_reports", {})["report_outline"] = outline_path
        return resolved_output_path

    resolved_pdf_path = Path(resolved_output_path)
    fallback_path = _preferred_output_path(
        str(resolved_pdf_path.with_name(f"{resolved_pdf_path.stem}_pdf_fallback.txt"))
    )
    tree_artifacts = _decision_tree_artifacts(analysis_results)
    lines = [
        "Multi-Agent Analytics Report",
        datetime.now().strftime("%B %d, %Y"),
        "",
        "Executive Summary",
        "\n".join(_format_executive_decision_block(workflow_state)),
        "",
        f"User context: {user_description}" if user_description else "",
        _stringify(outputs.get("decision_maker", {}).get("executive_summary", "")),
        "",
        "Decision Context",
        _stringify(outputs.get("decision_maker", {}).get("decision_context", "")),
        "",
        "Objective Coverage",
        "\n".join(_format_objective_coverage(workflow_state)),
        "",
        "Dataset Overview",
        _format_dataset_overview(outputs.get("data_understander", {})),
        "",
        "Data Quality Notes",
        "\n".join(_format_data_quality_notes(outputs.get("data_understander", {}))),
        "",
        "Market Research",
        "\n".join(
            [line for pair in _market_claim_pairs(outputs.get("market_researcher", {})) for line in pair if line]
        ),
        "",
        "Data Analysis and Visual Findings",
        "\n".join(_format_analysis_findings(analysis_results)),
        "",
        "Top Risk Signals",
        "\n".join(_format_top_risk_signals(analysis_results)),
    ]
    if tree_artifacts:
        lines.extend(
            [
                "",
                _eda_model_transition_text(),
                "",
                "Decision Tree Model",
                "\n".join(
                    _decision_tree_narrative(item) or _stringify(item.get("title"))
                    for item in tree_artifacts[:1]
                ),
            ]
        )
    lines.extend(
        [
            "",
            "Business Translation",
            _stringify(outputs.get("business_translator", {}).get("business_narrative", "")),
            "\n".join(_format_priority_findings(outputs.get("business_translator", {}))),
            "",
            "Recommendations",
            "\n".join(_format_recommendations(outputs.get("decision_maker", {}).get("recommendations", []))),
            _stringify(outputs.get("decision_maker", {}).get("conclusion", "")),
            "",
            "Limitations and Validation Steps",
            "\n".join(_format_limitations(workflow_state)),
            "",
            "Appendix / Workflow Trace",
            "\n".join(_format_workflow_trace(workflow_state)),
        ]
    )
    Path(fallback_path).write_text("\n".join(lines), encoding="utf-8")
    outline_path = _write_json_artifact(
        fallback_path,
        "report_outline.json",
        _report_outline_payload(workflow_state, fallback_path),
    )
    workflow_state.setdefault("generated_reports", {})["report_outline"] = outline_path
    return fallback_path


def _trim_trailing_spacers_and_pagebreaks(story: list[Any]) -> None:
    while story and story[-1].__class__.__name__ in {"Spacer", "PageBreak"}:
        story.pop()


def inspect_pdf_report(path: str) -> dict[str, Any]:
    """Run deterministic semantic/readability checks on the final PDF when PyMuPDF is available."""
    result: dict[str, Any] = {"path": str(Path(path).resolve()), "page_count": 0, "issues": [], "valid": True}
    pdf_path = Path(path)
    if not pdf_path.is_file() or pdf_path.stat().st_size < 1000:
        result["issues"].append({"issue": "pdf_missing_or_empty", "severity": "error"})
        result["valid"] = False
        return result
    try:
        import fitz  # type: ignore
    except ImportError:
        result["issues"].append({"issue": "rendered_pdf_qa_unavailable", "severity": "warning"})
        return result
    try:
        document = fitz.open(str(pdf_path))
    except Exception as exc:
        result["issues"].append({"issue": "pdf_cannot_open", "severity": "error", "detail": str(exc)})
        result["valid"] = False
        return result
    result["page_count"] = len(document)
    all_text: list[str] = []
    for index, page in enumerate(document, start=1):
        text = " ".join(page.get_text("text").split())
        all_text.append(text)
        if index > 1 and len(text) > 3200:
            result["issues"].append(
                {"page": index, "issue": "excessive_text_density", "characters": len(text), "severity": "warning"}
            )
        if index > 1 and len(text) < 45 and not page.get_images(full=True):
            result["issues"].append({"page": index, "issue": "nearly_empty_page", "severity": "warning"})
    combined = "\n".join(all_text)
    semantic_errors = {
        "raw_json_in_reader_text": bool(re.search(r"\{\s*[\"']?[A-Za-z_]+[\"']?\s*:", combined)),
        "wrong_domain_workforce_text": "observed workforce patterns" in combined.lower(),
        "unrenderable_cjk_or_language_mismatch": bool(re.search(r"[\u3400-\u9fff]", combined)),
    }
    for issue, present in semantic_errors.items():
        if present:
            result["issues"].append({"issue": issue, "severity": "error"})
    result["valid"] = not any(item.get("severity") == "error" for item in result["issues"])
    return result


def generate_slide_deck(
    workflow_state: dict[str, Any],
    output_path: str = "analytics_report.pptx",
    *,
    runtime_config: Any | None = None,
) -> str:
    from .deck_rendering import build_consulting_deck

    return build_consulting_deck(workflow_state, output_path, runtime_config=runtime_config)
