from __future__ import annotations

import math
import os
import tempfile
from typing import Any

try:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    MATPLOTLIB_AVAILABLE = True
except ImportError:  # pragma: no cover - dependency is present in normal runtime
    MATPLOTLIB_AVAILABLE = False

try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_CHARTS_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_CHARTS_AVAILABLE = False

    class _FallbackAlign:
        LEFT = "left"
        CENTER = "center"
        RIGHT = "right"

    PP_ALIGN = _FallbackAlign()

from .deck_spec import VisualSpec
from .theme import ConsultingTheme
from .text_refiner import compact_whitespace, shorten
from ..tree_diagram import build_tree_layout, clean_tree_text


SUPPORTED_CHART_TYPES = {
    "bar",
    "column",
    "grouped_bar",
    "horizontal_bar",
    "ranking",
    "small_multiples_bar",
    "line",
    "scatter",
    "distribution",
    "metric_cards",
    "comparison",
    "decision_tree",
}

_CATEGORY_KEYS = ("category", "label", "name", "level", "segment", "group", "role", "dimension", "x")
_VALUE_KEYS = ("value", "rate", "attrition_rate", "default_rate", "percent", "percentage", "score", "count", "y")
_GROUP_KEYS = ("series", "group", "segment", "dimension", "cohort")


def _number(value: Any) -> float | None:
    try:
        if value is None or value == "":
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _rows_from_visual(visual: VisualSpec) -> list[dict[str, Any]]:
    data = visual.data
    if isinstance(data, list):
        return [row for row in data if isinstance(row, dict)]
    if isinstance(data, dict):
        for key in ("rows", "data", "values"):
            rows = data.get(key)
            if isinstance(rows, list):
                return [row for row in rows if isinstance(row, dict)]
    return []


def _series_from_visual(visual: VisualSpec) -> list[dict[str, Any]]:
    if isinstance(visual.series, list):
        return [entry for entry in visual.series if isinstance(entry, dict)]
    return []


def _normalized_chart_rows(visual: VisualSpec) -> list[dict[str, Any]]:
    rows = _rows_from_visual(visual)
    if rows:
        return rows
    series_items = _series_from_visual(visual)
    return _flatten_series(series_items, visual) if series_items else []


def has_structured_chart_data(visual: VisualSpec | None) -> bool:
    if not visual:
        return False
    chart_type = compact_whitespace(visual.chart_type).lower()
    if chart_type not in SUPPORTED_CHART_TYPES:
        return False
    if chart_type == "small_multiples_bar":
        return any(isinstance(series.get("data"), list) for series in _series_from_visual(visual))
    if chart_type == "metric_cards":
        return bool(_rows_from_visual(visual) or _series_from_visual(visual))
    if chart_type == "decision_tree":
        if visual.image_path or visual.fallback_path:
            return True
        return bool(_tree_nodes_edges(visual)[0])
    rows = _normalized_chart_rows(visual)
    if not rows:
        return False
    return _infer_category_value_keys(rows, visual)[0] != "" and _infer_category_value_keys(rows, visual)[1] != ""


def can_render_native(visual: VisualSpec) -> bool:
    if not PPTX_CHARTS_AVAILABLE or not visual:
        return False
    chart_type = visual.chart_type.lower()
    rows = _normalized_chart_rows(visual)
    return chart_type in {"bar", "column", "line", "scatter"} and bool(rows)


def add_native_chart(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    if not can_render_native(visual):
        return False
    rows = _normalized_chart_rows(visual)
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return False
    chart_type = visual.chart_type.lower()
    try:
        if chart_type == "scatter":
            chart_data = XyChartData()
            series = chart_data.add_series(visual.y or "Series 1")
            for row in rows:
                x_value = _number(row.get(category_key))
                y_value = _number(row.get(value_key))
                if x_value is not None and y_value is not None:
                    series.add_data_point(x_value, y_value)
            chart_enum = XL_CHART_TYPE.XY_SCATTER
        else:
            chart_data = CategoryChartData()
            group_key = _group_key(rows, visual)
            categories = list(dict.fromkeys(compact_whitespace(row.get(category_key)) for row in rows if compact_whitespace(row.get(category_key))))
            chart_data.categories = categories
            if group_key:
                for series_name in dict.fromkeys(compact_whitespace(row.get(group_key)) for row in rows):
                    values = []
                    for category in categories:
                        matching = next(
                            (
                                row
                                for row in rows
                                if compact_whitespace(row.get(category_key)) == category
                                and compact_whitespace(row.get(group_key)) == series_name
                            ),
                            {},
                        )
                        values.append(_number(matching.get(value_key)) or 0)
                    chart_data.add_series(series_name or "Series", values)
            else:
                chart_data.add_series(visual.title or value_key, [_number(row.get(value_key)) or 0 for row in rows if compact_whitespace(row.get(category_key))])
            chart_enum = XL_CHART_TYPE.LINE if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED

        graphic_frame = slide.shapes.add_chart(chart_enum, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
        graphic_frame.name = "reconstructed_chart_native"
        chart = graphic_frame.chart
        chart.has_legend = bool(_group_key(rows, visual))
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.has_title = False
        for idx, series in enumerate(chart.series):
            color = theme.chart_palette[idx % len(theme.chart_palette)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(color)
            series.format.line.color.rgb = _rgb(color)
        return True
    except Exception:
        return False


def render_matplotlib_fallback(visual: VisualSpec, theme: ConsultingTheme) -> str:
    if not MATPLOTLIB_AVAILABLE:
        return ""
    rows = _normalized_chart_rows(visual)
    if not rows:
        return ""
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return ""
    chart_type = visual.chart_type.lower()
    x_values = [compact_whitespace(row.get(category_key)) for row in rows]
    y_values = [_number(row.get(value_key)) or 0 for row in rows]
    if not x_values or not y_values:
        return ""

    fig, ax = plt.subplots(figsize=(8.5, 4.2), dpi=160)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    colors = ["#2c65a4", "#35807f", "#bb9241", "#49845c", "#a8534c"]
    if chart_type == "line":
        ax.plot(x_values, y_values, color=colors[0], linewidth=2.4, marker="o", markersize=4)
    elif chart_type == "scatter":
        ax.scatter(x_values, y_values, color=colors[0], s=36)
    elif chart_type == "horizontal_bar":
        ax.barh(x_values, y_values, color=colors[0])
    else:
        ax.bar(x_values, y_values, color=colors[0], width=0.62)
    ax.set_title(visual.title or "", loc="left", fontsize=12, color="#1f2730", pad=12)
    ax.set_xlabel(visual.x_label or "", fontsize=9, color="#5e6974")
    ax.set_ylabel(visual.y_label or "", fontsize=9, color="#5e6974")
    ax.tick_params(axis="x", labelrotation=30, labelsize=8, colors="#5e6974")
    ax.tick_params(axis="y", labelsize=8, colors="#5e6974")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#d1d7dc")
    ax.spines["bottom"].set_color("#d1d7dc")
    ax.grid(axis="y", color="#e7eaed", linewidth=0.8)
    fig.tight_layout()
    handle = tempfile.NamedTemporaryFile(prefix="figure_slide_fallback_", suffix=".png", delete=False)
    fallback_path = handle.name
    handle.close()
    fig.savefig(fallback_path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return fallback_path


def add_structured_chart(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    chart_type = compact_whitespace(visual.chart_type).lower()
    if not chart_type:
        return False
    if chart_type == "small_multiples_bar":
        return _add_small_multiples_bar(slide, visual, left, top, width, height, theme)
    if chart_type in {"horizontal_bar", "ranking"}:
        return _add_horizontal_bar(slide, visual, left, top, width, height, theme)
    if chart_type in {"bar", "column", "distribution"}:
        return _add_grouped_or_vertical_bar(slide, visual, left, top, width, height, theme)
    if chart_type == "grouped_bar":
        return _add_grouped_bar(slide, visual, left, top, width, height, theme)
    if chart_type in {"line", "scatter"}:
        return _add_line_chart(slide, visual, left, top, width, height, theme)
    if chart_type == "metric_cards":
        return _add_metric_cards(slide, visual, left, top, width, height, theme)
    if chart_type == "comparison":
        return _add_horizontal_bar(slide, visual, left, top, width, height, theme)
    if chart_type == "decision_tree":
        return _add_decision_tree(slide, visual, left, top, width, height, theme)
    return False


def _tree_nodes_edges(visual: VisualSpec) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    data = visual.data if isinstance(visual.data, dict) else {}
    raw_nodes = data.get("nodes", []) if isinstance(data, dict) else []
    raw_edges = data.get("edges", []) if isinstance(data, dict) else []
    nodes = [node for node in raw_nodes if isinstance(node, dict)]
    edges = [edge for edge in raw_edges if isinstance(edge, dict)]
    if nodes:
        return _simplify_tree_for_slide(nodes, edges)
    rules = data.get("rules", []) if isinstance(data, dict) else []
    rule_texts = [compact_whitespace(rule) for rule in rules if compact_whitespace(rule)]
    if not rule_texts:
        rows = _rows_from_visual(visual)
        rule_texts = [compact_whitespace(row.get("rule") or row.get("label")) for row in rows if compact_whitespace(row.get("rule") or row.get("label"))]
    if not rule_texts:
        return [], []
    nodes = [{"id": "root", "label": "Decision tree rules", "depth": 0}]
    edges = []
    for index, rule in enumerate(rule_texts[:6], start=1):
        node_id = f"rule_{index}"
        nodes.append({"id": node_id, "label": rule, "depth": 1})
        edges.append({"source": "root", "target": node_id, "label": ""})
    return nodes, edges


def _simplify_tree_for_slide(nodes: list[dict[str, Any]], edges: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    if len(nodes) <= 7:
        return nodes, edges[:10]
    node_by_id = {
        compact_whitespace(node.get("id") or node.get("node_id") or node.get("name")): node
        for node in nodes
        if compact_whitespace(node.get("id") or node.get("node_id") or node.get("name"))
    }
    root_id = ""
    targets = {compact_whitespace(edge.get("target") or edge.get("to") or edge.get("child")) for edge in edges}
    for node_id in node_by_id:
        if node_id not in targets:
            root_id = node_id
            break
    root_id = root_id or next(iter(node_by_id), "")
    selected = {root_id}
    for edge in edges:
        source = compact_whitespace(edge.get("source") or edge.get("from") or edge.get("parent"))
        target = compact_whitespace(edge.get("target") or edge.get("to") or edge.get("child"))
        if source == root_id and target:
            selected.add(target)
    for edge in edges:
        if len(selected) >= 7:
            break
        source = compact_whitespace(edge.get("source") or edge.get("from") or edge.get("parent"))
        target = compact_whitespace(edge.get("target") or edge.get("to") or edge.get("child"))
        if source in selected and target:
            selected.add(target)
    simplified_nodes = [node for node in nodes if compact_whitespace(node.get("id") or node.get("node_id") or node.get("name")) in selected]
    simplified_ids = {compact_whitespace(node.get("id") or node.get("node_id") or node.get("name")) for node in simplified_nodes}
    simplified_edges = [
        edge
        for edge in edges
        if compact_whitespace(edge.get("source") or edge.get("from") or edge.get("parent")) in simplified_ids
        and compact_whitespace(edge.get("target") or edge.get("to") or edge.get("child")) in simplified_ids
    ][:10]
    return simplified_nodes, simplified_edges


def _add_decision_tree(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    nodes, edges = _tree_nodes_edges(visual)
    if not nodes:
        return False
    _chart_title(slide, visual, left, top, width, theme)
    layout = build_tree_layout(nodes, edges)
    layout_nodes = layout.get("nodes", [])
    layout_edges = layout.get("edges", [])
    normalized_positions = layout.get("positions", {})
    if not layout_nodes or not normalized_positions:
        return False

    chart_top = top + 0.48
    chart_h = max(height - 0.8, 1.4)
    chart_left = left + 0.1
    chart_w = width - 0.2
    max_depth = int(layout.get("max_depth", 0) or 0)
    node_w = min(max(chart_w / max(_tree_leaf_count(layout_nodes, layout_edges), 2) - 0.1, 1.65), 2.45)
    node_h = 0.68 if max_depth <= 2 else 0.58
    positions: dict[str, tuple[float, float, float, float]] = {}
    for node in layout_nodes:
        node_id = str(node.get("_layout_id", ""))
        if node_id not in normalized_positions:
            continue
        x_norm, y_norm = normalized_positions[node_id]
        node_left = chart_left + (chart_w - node_w) * float(x_norm)
        node_top = chart_top + (chart_h - node_h - 0.24) * float(y_norm)
        positions[node_id] = (node_left, node_top, node_w, node_h)

    split_node_ids = {edge["source"] for edge in layout_edges}
    for edge in layout_edges:
        source = edge["source"]
        target = edge["target"]
        if source not in positions or target not in positions:
            continue
        sx, sy, sw, sh = positions[source]
        tx, ty, tw, _ = positions[target]
        start_x = sx + sw / 2
        start_y = sy + sh
        end_x = tx + tw / 2
        end_y = ty
        joint_y = start_y + max((end_y - start_y) * 0.42, 0.12)
        _decision_tree_line(slide, start_x, start_y, start_x, joint_y, theme)
        _decision_tree_line(slide, start_x, joint_y, end_x, joint_y, theme)
        _decision_tree_line(slide, end_x, joint_y, end_x, end_y, theme)
        label = compact_whitespace(edge.get("label"))
        if label:
            label_left = min(start_x, end_x) + abs(end_x - start_x) / 2 - 0.45
            _text(slide, shorten(label, 18), label_left, joint_y - 0.2, 0.9, 0.16, theme, size=6, color=theme.blue, bold=True, align=PP_ALIGN.CENTER)

    for index, node in enumerate(layout_nodes):
        node_id = str(node.get("_layout_id", ""))
        if node_id not in positions:
            continue
        node_left, node_top, node_w, node_h = positions[node_id]
        is_split = node_id in split_node_ids
        fill = theme.blue if is_split else theme.white
        line = theme.blue if not is_split else theme.blue
        shape = _rect(slide, node_left, node_top, node_w, node_h, fill, line=line)
        shape.name = "reconstructed_chart_decision_tree_node"
        label = _decision_tree_node_label(node, is_split, visual)
        text_color = theme.white if is_split else theme.navy
        label_box = _text(
            slide,
            shorten(label, 78),
            node_left + 0.08,
            node_top + 0.08,
            node_w - 0.16,
            node_h - 0.13,
            theme,
            size=8 if max_depth <= 2 else 7,
            color=text_color,
            bold=is_split,
            align=PP_ALIGN.CENTER,
        )
        label_box.name = "reconstructed_chart_decision_tree_label"
    metrics = []
    data = visual.data if isinstance(visual.data, dict) else {}
    for key in (
        "train_accuracy",
        "test_accuracy",
        "train_score",
        "test_score",
        "accuracy",
        "baseline_accuracy",
        "baseline_score",
        "precision",
        "recall",
        "f1",
        "f1_score",
        "support",
        "positive_class_rate",
        "train_r2",
        "test_r2",
        "r2",
        "train_mae",
        "test_mae",
        "mae",
    ):
        value = data.get(key)
        if value not in (None, ""):
            metrics.append(f"{_decision_tree_metric_label(key, data)}: {_decision_tree_metric_value(key, value)}")
    if metrics:
        _text(slide, " | ".join(metrics[:4]), left + 0.05, top + height - 0.2, width - 0.1, 0.17, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
    return True


def _tree_leaf_count(nodes: list[dict[str, Any]], edges: list[dict[str, Any]]) -> int:
    sources = {edge.get("source") for edge in edges}
    leaves = [node for node in nodes if node.get("_layout_id") not in sources]
    return max(len(leaves), 1)


def _decision_tree_node_label(node: dict[str, Any], is_split: bool, visual: VisualSpec | None = None) -> str:
    display = clean_tree_text(node.get("display_label") or node.get("human_label"))
    if display:
        return _clean_tree_label_for_context(display, visual)
    if is_split:
        feature = _clean_tree_label_for_context(clean_tree_text(node.get("feature")), visual)
        threshold = clean_tree_text(node.get("threshold"))
        if feature and threshold:
            return _clean_tree_label_for_context(f"{feature} <= {threshold}", visual)
    else:
        prediction = clean_tree_text(node.get("prediction") or node.get("class") or node.get("value"))
        samples = clean_tree_text(node.get("samples") or node.get("n") or node.get("support"))
        if prediction:
            return _clean_tree_label_for_context(f"Predict {prediction}" + (f"\nn={samples}" if samples else ""), visual)
    for key in ("label", "rule", "name", "condition", "prediction", "value", "class"):
        text = clean_tree_text(node.get(key))
        if text:
            return _clean_tree_label_for_context(text, visual)
    return "Split node" if is_split else "Leaf"


def _clean_tree_label_for_context(text: str, visual: VisualSpec | None) -> str:
    clean = compact_whitespace(text).replace("_", " ")
    data = visual.data if visual and isinstance(visual.data, dict) else {}
    target = compact_whitespace(data.get("target") or data.get("target_column")).lower()
    context = " ".join(
        compact_whitespace(value).lower()
        for value in (target, visual.title if visual else "", visual.finding if visual else "", visual.takeaway if visual else "")
        if compact_whitespace(value)
    )
    if any(token in context for token in ("loan", "default", "credit")):
        replacements = {
            "loan percent income": "debt burden",
            "loan int rate": "interest rate",
            "person income": "borrower income",
            "person age": "borrower age",
            "person home ownership": "home ownership",
            "cb person default on file": "prior default flag",
            "loan grade": "loan grade",
        }
        for source, replacement in replacements.items():
            clean = clean.replace(source, replacement)
        clean = clean.replace("<=", "at or below").replace(">=", "at least").replace(">", "above").replace("<", "below")
    return compact_whitespace(clean)


def _decision_tree_line(slide: Any, x1: float, y1: float, x2: float, y2: float, theme: ConsultingTheme) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = "reconstructed_chart_decision_tree_edge"
    line.line.color.rgb = RGBColor(*theme.blue)
    line.line.width = Pt(1.1)


def _decision_tree_metric_label(key: str, data: dict[str, Any]) -> str:
    model_type = str(data.get("model_type", "")).lower()
    if key == "train_score" and not model_type.startswith("reg"):
        return "Train Accuracy"
    if key == "test_score" and not model_type.startswith("reg"):
        return "Test Accuracy"
    if key == "baseline_score" and not model_type.startswith("reg"):
        return "Baseline Accuracy"
    return key.replace("_", " ").title()


def _decision_tree_metric_value(key: str, value: Any) -> str:
    if key in {"train_score", "test_score", "baseline_score", "train_accuracy", "test_accuracy", "baseline_accuracy", "accuracy", "precision", "recall", "f1", "f1_score", "positive_class_rate"}:
        try:
            number = float(value)
        except (TypeError, ValueError):
            return str(value)
        if 0 <= number <= 1:
            return f"{number * 100:.1f}%"
    return str(value)


def _add_grouped_or_vertical_bar(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    if not rows:
        return _add_grouped_bar(slide, visual, left, top, width, height, theme)
    if _group_key(rows, visual):
        return _add_grouped_bar(slide, visual, left, top, width, height, theme)
    return _add_vertical_bar(slide, visual, left, top, width, height, theme)


def _add_vertical_bar(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return False
    values = [(compact_whitespace(row.get(category_key)), _number(row.get(value_key))) for row in rows]
    values = [(label, value) for label, value in values if label and value is not None][:9]
    if not values:
        return False
    _chart_title(slide, visual, left, top, width, theme)
    chart_left, chart_top, chart_w, chart_h = _plot_area(left, top, width, height, title=True, bottom_label=True)
    max_value = _nice_max([value for _, value in values])
    _draw_y_grid(slide, chart_left, chart_top, chart_w, chart_h, max_value, theme)
    baseline = chart_top + chart_h
    bar_gap = 0.08
    bar_w = max((chart_w - bar_gap * (len(values) - 1)) / len(values), 0.12)
    for index, (label, value) in enumerate(values):
        bar_h = chart_h * value / max_value if max_value else 0
        bar_left = chart_left + index * (bar_w + bar_gap)
        shape = _rect(slide, bar_left, baseline - bar_h, bar_w, max(bar_h, 0.035), theme.chart_palette[index % len(theme.chart_palette)])
        shape.name = "reconstructed_chart_bar"
        _text(slide, _format_value(value, visual), bar_left - 0.04, max(baseline - bar_h - 0.24, chart_top - 0.04), bar_w + 0.08, 0.17, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
        _text(slide, shorten(label, 18), bar_left - 0.06, baseline + 0.06, bar_w + 0.12, 0.38, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
    _axis_labels(slide, visual, left, top, width, height, theme)
    return True


def _add_grouped_bar(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    category_key, value_key = _infer_category_value_keys(rows, visual)
    group_key = _group_key(rows, visual)
    if not rows or not category_key or not value_key or not group_key:
        return _add_vertical_bar(slide, visual, left, top, width, height, theme)
    categories = list(dict.fromkeys(compact_whitespace(row.get(category_key)) for row in rows if compact_whitespace(row.get(category_key))))[:6]
    groups = list(dict.fromkeys(compact_whitespace(row.get(group_key)) for row in rows if compact_whitespace(row.get(group_key))))[:5]
    if not categories or not groups:
        return False
    _chart_title(slide, visual, left, top, width, theme)
    chart_left, chart_top, chart_w, chart_h = _plot_area(left, top, width, height, title=True, bottom_label=True, legend=True)
    values = [_number(row.get(value_key)) or 0 for row in rows]
    max_value = _nice_max(values)
    _draw_y_grid(slide, chart_left, chart_top, chart_w, chart_h, max_value, theme)
    baseline = chart_top + chart_h
    cluster_gap = 0.15
    cluster_w = max((chart_w - cluster_gap * (len(categories) - 1)) / len(categories), 0.28)
    bar_gap = 0.025
    bar_w = max((cluster_w - bar_gap * (len(groups) - 1)) / len(groups), 0.045)
    for category_index, category in enumerate(categories):
        cluster_left = chart_left + category_index * (cluster_w + cluster_gap)
        for group_index, group in enumerate(groups):
            row = next(
                (
                    item
                    for item in rows
                    if compact_whitespace(item.get(category_key)) == category
                    and compact_whitespace(item.get(group_key)) == group
                ),
                {},
            )
            value = _number(row.get(value_key)) or 0
            bar_h = chart_h * value / max_value if max_value else 0
            bar_left = cluster_left + group_index * (bar_w + bar_gap)
            shape = _rect(slide, bar_left, baseline - bar_h, bar_w, max(bar_h, 0.03), theme.chart_palette[group_index % len(theme.chart_palette)])
            shape.name = "reconstructed_chart_grouped_bar"
        _text(slide, shorten(category, 18), cluster_left - 0.02, baseline + 0.06, cluster_w + 0.04, 0.32, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
    _legend(slide, groups, chart_left, chart_top + chart_h + 0.48, chart_w, theme)
    _axis_labels(slide, visual, left, top, width, height, theme)
    return True


def _add_small_multiples_bar(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    series_items = _series_from_visual(visual)
    if not series_items:
        rows = _rows_from_visual(visual)
        category_key, value_key = _infer_category_value_keys(rows, visual)
        group_key = _group_key(rows, visual)
        if not category_key or not value_key or not group_key:
            return False
        series_items = [
            {"name": group, "data": [row for row in rows if compact_whitespace(row.get(group_key)) == group]}
            for group in dict.fromkeys(compact_whitespace(row.get(group_key)) for row in rows if compact_whitespace(row.get(group_key)))
        ]
    series_items = series_items[:4]
    all_values: list[float] = []
    for series in series_items:
        for point in series.get("data", []) or []:
            if isinstance(point, dict):
                value = _first_numeric(point, preferred=visual.y or "value", exclude={visual.x or "level", "level", "category", "label"})
                if value is not None:
                    all_values.append(value)
    max_value = _nice_max(all_values)
    if max_value <= 0:
        return False

    _chart_title(slide, visual, left, top, width, theme)
    top += 0.34
    height -= 0.34
    cols = 2 if len(series_items) > 1 else 1
    rows_count = 2 if len(series_items) > 2 else 1
    gap_x = 0.22
    gap_y = 0.28
    panel_w = (width - gap_x * (cols - 1)) / cols
    panel_h = (height - gap_y * (rows_count - 1)) / rows_count

    for index, series in enumerate(series_items):
        col = index % cols
        row = index // cols
        panel_left = left + col * (panel_w + gap_x)
        panel_top = top + row * (panel_h + gap_y)
        panel = _rect(slide, panel_left, panel_top, panel_w, panel_h, theme.white, line=theme.line)
        panel.name = "reconstructed_chart_small_multiple_panel"
        _text(slide, compact_whitespace(series.get("name")) or f"Series {index + 1}", panel_left + 0.1, panel_top + 0.08, panel_w - 0.2, 0.22, theme, size=8, bold=True, color=theme.navy)
        points = [point for point in series.get("data", []) or [] if isinstance(point, dict)][:6]
        if not points:
            continue
        chart_left = panel_left + 0.22
        chart_top = panel_top + 0.48
        chart_w = panel_w - 0.42
        chart_h = panel_h - 0.82
        baseline = chart_top + chart_h
        _rect(slide, chart_left, baseline, chart_w, 0.01, theme.line).name = "reconstructed_chart_axis"
        bar_gap = 0.04
        bar_w = max((chart_w - bar_gap * (len(points) - 1)) / max(len(points), 1), 0.08)
        for point_index, point in enumerate(points):
            value = _first_numeric(point, preferred=visual.y or "value", exclude={visual.x or "level", "level", "category", "label"}) or 0
            label = compact_whitespace(point.get(visual.x or "level") or point.get("level") or point.get("category") or point.get("label") or point_index + 1)
            bar_h = chart_h * (value / max_value)
            bar_left = chart_left + point_index * (bar_w + bar_gap)
            bar_top = baseline - bar_h
            shape = _rect(slide, bar_left, bar_top, bar_w, max(bar_h, 0.02), theme.chart_palette[index % len(theme.chart_palette)])
            shape.name = "reconstructed_chart_small_multiple_bar"
            _text(slide, _format_value(value, visual), bar_left - 0.04, max(bar_top - 0.18, chart_top - 0.05), bar_w + 0.08, 0.14, theme, size=6, color=theme.muted, align=PP_ALIGN.CENTER)
            _text(slide, shorten(label, 10), bar_left - 0.03, baseline + 0.04, bar_w + 0.06, 0.14, theme, size=6, color=theme.muted, align=PP_ALIGN.CENTER)
    return True


def _add_horizontal_bar(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return False
    ranked = []
    for row in rows:
        value = _number(row.get(value_key))
        label = compact_whitespace(row.get(category_key))
        if value is not None and label:
            ranked.append((label, value))
    ranked = sorted(ranked, key=lambda item: item[1], reverse=True)[:8]
    if not ranked:
        return False
    _chart_title(slide, visual, left, top, width, theme)
    chart_top = top + 0.48
    chart_h = height - 0.62
    max_value = _nice_max([value for _, value in ranked])
    label_w = min(width * 0.34, 2.7)
    bar_left = left + label_w + 0.18
    bar_w = width - label_w - 0.76
    row_h = min(chart_h / max(len(ranked), 1), 0.46)
    for index, (label, value) in enumerate(ranked):
        row_top = chart_top + index * row_h
        _text(slide, shorten(label, 28), left, row_top + 0.03, label_w, 0.22, theme, size=8, color=theme.ink)
        _rect(slide, bar_left, row_top + 0.075, bar_w, 0.14, (235, 239, 242)).name = "reconstructed_chart_bar_background"
        shape = _rect(slide, bar_left, row_top + 0.075, max(bar_w * value / max_value, 0.04), 0.14, theme.chart_palette[index % len(theme.chart_palette)])
        shape.name = "reconstructed_chart_horizontal_bar"
        _text(slide, _format_value(value, visual), bar_left + bar_w + 0.08, row_top + 0.03, 0.6, 0.2, theme, size=8, color=theme.muted)
    return True


def _add_line_chart(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return False
    values = [(compact_whitespace(row.get(category_key)), _number(row.get(value_key))) for row in rows]
    values = [(label, value) for label, value in values if label and value is not None][:10]
    if len(values) < 2:
        return False
    _chart_title(slide, visual, left, top, width, theme)
    chart_left, chart_top, chart_w, chart_h = _plot_area(left, top, width, height, title=True, bottom_label=True)
    max_value = _nice_max([value for _, value in values])
    min_value = min(value for _, value in values)
    value_range = max(max_value - min_value, 1)
    _draw_y_grid(slide, chart_left, chart_top, chart_w, chart_h, max_value, theme)
    points: list[tuple[float, float, str, float]] = []
    for index, (label, value) in enumerate(values):
        x = chart_left + (chart_w * index / max(len(values) - 1, 1))
        y = chart_top + chart_h - ((value - min_value) / value_range * chart_h)
        points.append((x, y, label, value))
    for index in range(len(points) - 1):
        x1, y1, _, _ = points[index]
        x2, y2, _, _ = points[index + 1]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.name = "reconstructed_chart_line"
        line.line.color.rgb = RGBColor(*theme.blue)
        line.line.width = Pt(1.8)
    for x, y, label, value in points:
        marker = _oval(slide, x - 0.045, y - 0.045, 0.09, 0.09, theme.blue)
        marker.name = "reconstructed_chart_line_marker"
        _text(slide, _format_value(value, visual), x - 0.22, y - 0.28, 0.44, 0.16, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
        _text(slide, shorten(label, 14), x - 0.28, chart_top + chart_h + 0.06, 0.56, 0.3, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
    _axis_labels(slide, visual, left, top, width, height, theme)
    return True


def _add_metric_cards(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> bool:
    rows = _normalized_chart_rows(visual)
    if not rows:
        return False
    category_key, value_key = _infer_category_value_keys(rows, visual)
    if not category_key or not value_key:
        return False
    values = [(compact_whitespace(row.get(category_key)), _number(row.get(value_key))) for row in rows]
    values = [(label, value) for label, value in values if label and value is not None][:6]
    if not values:
        return False
    cols = 3 if len(values) > 2 else len(values)
    rows_count = math.ceil(len(values) / cols)
    gap = 0.18
    card_w = (width - gap * (cols - 1)) / cols
    card_h = min((height - gap * (rows_count - 1)) / rows_count, 1.25)
    for index, (label, value) in enumerate(values):
        col = index % cols
        row = index // cols
        card_left = left + col * (card_w + gap)
        card_top = top + row * (card_h + gap)
        card = _rect(slide, card_left, card_top, card_w, card_h, theme.white, line=theme.line)
        card.name = "reconstructed_chart_metric_card"
        _rect(slide, card_left, card_top, 0.06, card_h, theme.chart_palette[index % len(theme.chart_palette)])
        _text(slide, _format_value(value, visual), card_left + 0.16, card_top + 0.2, card_w - 0.28, 0.35, theme, size=16, color=theme.blue, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, shorten(label, 34), card_left + 0.16, card_top + 0.62, card_w - 0.28, 0.32, theme, size=8, color=theme.muted, align=PP_ALIGN.CENTER)
    return True


def _infer_category_value_keys(rows: list[dict[str, Any]], visual: VisualSpec) -> tuple[str, str]:
    if not rows:
        return "", ""
    keys = list(dict.fromkeys(key for row in rows for key in row.keys()))
    category_key = _first_existing_key(keys, visual.x, _CATEGORY_KEYS)
    value_key = _first_existing_key(keys, visual.y, _VALUE_KEYS, numeric_rows=rows, exclude={category_key})
    if not category_key:
        category_key = next((key for key in keys if key != value_key and any(_number(row.get(key)) is None for row in rows)), "")
    if not value_key:
        value_key = next((key for key in keys if key != category_key and any(_number(row.get(key)) is not None for row in rows)), "")
    return category_key, value_key


def _first_existing_key(
    keys: list[str],
    preferred: str,
    fallbacks: tuple[str, ...],
    *,
    numeric_rows: list[dict[str, Any]] | None = None,
    exclude: set[str] | None = None,
) -> str:
    exclude = exclude or set()
    if preferred and preferred in keys and preferred not in exclude:
        return preferred
    lowered = {str(key).lower(): key for key in keys}
    for candidate in fallbacks:
        key = lowered.get(candidate.lower())
        if key and key not in exclude:
            if numeric_rows is None or any(_number(row.get(key)) is not None for row in numeric_rows):
                return key
    return ""


def _group_key(rows: list[dict[str, Any]], visual: VisualSpec) -> str:
    if not rows:
        return ""
    keys = list(dict.fromkeys(key for row in rows for key in row.keys()))
    preferred = visual.group_by or (visual.series if isinstance(visual.series, str) else "")
    category_key, value_key = _infer_category_value_keys(rows, visual)
    key = _first_existing_key(keys, preferred, _GROUP_KEYS, exclude={category_key, value_key})
    if key and len({compact_whitespace(row.get(key)) for row in rows}) > 1:
        return key
    return ""


def _flatten_series(series_items: list[dict[str, Any]], visual: VisualSpec) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for series in series_items:
        series_name = compact_whitespace(series.get("name")) or "Series"
        for point in series.get("data", []) or []:
            if not isinstance(point, dict):
                continue
            row = dict(point)
            row.setdefault("series", series_name)
            rows.append(row)
    return rows


def _first_numeric(point: dict[str, Any], preferred: str, exclude: set[str] | None = None) -> float | None:
    exclude = exclude or set()
    preferred_value = _number(point.get(preferred))
    if preferred_value is not None:
        return preferred_value
    for key, value in point.items():
        if key in exclude:
            continue
        numeric = _number(value)
        if numeric is not None:
            return numeric
    return None


def _format_value(value: float, visual: VisualSpec) -> str:
    if visual.value_format:
        try:
            return visual.value_format.format(value)
        except (IndexError, KeyError, ValueError):
            try:
                return format(value, visual.value_format)
            except ValueError:
                pass
    if "%" in (visual.y_label or "") or "rate" in (visual.y or "").lower():
        return f"{value:.1f}%"
    return f"{value:g}"


def _nice_max(values: list[float]) -> float:
    max_value = max(values) if values else 1
    if max_value <= 0:
        return 1
    magnitude = 10 ** math.floor(math.log10(max_value))
    return math.ceil(max_value / magnitude * 1.12) * magnitude


def _plot_area(
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    title: bool = False,
    bottom_label: bool = False,
    legend: bool = False,
) -> tuple[float, float, float, float]:
    chart_left = left + 0.38
    chart_top = top + (0.52 if title else 0.18)
    chart_w = width - 0.62
    chart_h = height - (0.98 if title else 0.54) - (0.28 if bottom_label else 0) - (0.28 if legend else 0)
    return chart_left, chart_top, chart_w, max(chart_h, 0.5)


def _chart_title(slide: Any, visual: VisualSpec, left: float, top: float, width: float, theme: ConsultingTheme) -> None:
    if visual.title:
        _text(slide, shorten(visual.title, 78), left + 0.05, top + 0.02, width - 0.1, 0.28, theme, size=10, bold=True, color=theme.navy)


def _draw_y_grid(slide: Any, left: float, top: float, width: float, height: float, max_value: float, theme: ConsultingTheme) -> None:
    for step in range(4):
        y = top + height * step / 3
        line = _rect(slide, left, y, width, 0.006, theme.line)
        line.name = "reconstructed_chart_gridline"
    _rect(slide, left, top + height, width, 0.012, theme.line).name = "reconstructed_chart_axis"


def _axis_labels(slide: Any, visual: VisualSpec, left: float, top: float, width: float, height: float, theme: ConsultingTheme) -> None:
    if visual.x_label:
        _text(slide, shorten(visual.x_label, 50), left + 0.45, top + height - 0.18, width - 0.9, 0.16, theme, size=7, color=theme.muted, align=PP_ALIGN.CENTER)
    if visual.y_label:
        _text(slide, shorten(visual.y_label, 38), left + 0.04, top + 0.36, 1.6, 0.16, theme, size=7, color=theme.muted)


def _legend(slide: Any, labels: list[str], left: float, top: float, width: float, theme: ConsultingTheme) -> None:
    cursor = left
    for index, label in enumerate(labels[:5]):
        _rect(slide, cursor, top + 0.03, 0.12, 0.08, theme.chart_palette[index % len(theme.chart_palette)]).name = "reconstructed_chart_legend_swatch"
        _text(slide, shorten(label, 18), cursor + 0.16, top, 1.05, 0.14, theme, size=6, color=theme.muted)
        cursor += min(1.35, max(0.75, len(label) * 0.055 + 0.28))


def _rect(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: tuple[int, int, int],
    *,
    line: tuple[int, int, int] | None = None,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    if line:
        shape.line.color.rgb = RGBColor(*line)
        shape.line.width = Pt(0.45)
    else:
        shape.line.fill.background()
    return shape


def _oval(slide: Any, left: float, top: float, width: float, height: float, fill: tuple[int, int, int]) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.fill.background()
    return shape


def _text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ConsultingTheme,
    *,
    size: int = 8,
    color: tuple[int, int, int] | None = None,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.01)
    frame.margin_right = Inches(0.01)
    frame.margin_top = Inches(0.0)
    frame.margin_bottom = Inches(0.0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = theme.font_family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*(color or theme.ink))
    return box


def resolve_image_fallback(visual: VisualSpec, theme: ConsultingTheme) -> str:
    for candidate in (visual.image_path, visual.fallback_path):
        if candidate and os.path.exists(candidate):
            return candidate
    return render_matplotlib_fallback(visual, theme)


def _rgb(value: tuple[int, int, int]) -> Any:
    return RGBColor(*value)
