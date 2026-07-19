from __future__ import annotations

import os
import re
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Any, Callable

try:
    from PIL import Image as PILImage
except ImportError:  # pragma: no cover
    PILImage = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False

from .chart_renderer import add_native_chart, add_structured_chart, resolve_image_fallback
from .deck_spec import DeckSpec, SlideSpec
from .theme import THEME, ConsultingTheme
from .text_refiner import shorten


def render_deck(deck: DeckSpec, output_path: str = "analytics_report.pptx") -> str:
    resolved_output_path = _preferred_output_path(output_path)
    if not PPTX_AVAILABLE:
        return _render_text_fallback(deck, resolved_output_path)

    prs = Presentation()
    prs.slide_width = Inches(THEME.width)
    prs.slide_height = Inches(THEME.height)
    renderer = PowerPointRenderer(prs, deck, THEME)
    for slide in deck.slides:
        renderer.render_slide(slide)
    prs.save(resolved_output_path)
    return resolved_output_path


class PowerPointRenderer:
    def __init__(self, presentation: Any, deck: DeckSpec, theme: ConsultingTheme) -> None:
        self.prs = presentation
        self.deck = deck
        self.theme = theme

    def render_slide(self, spec: SlideSpec) -> None:
        renderers: dict[str, Callable[[SlideSpec], None]] = {
            "title_cover": self._render_title_cover,
            "data_understanding_overview": self._render_data_understanding,
            "market_context_bullets": self._render_market_context,
            "chart_left_insight_right": self._render_chart_left,
            "chart_right_insight_left": self._render_chart_right,
            "full_width_chart_takeaway": self._render_full_width_chart,
            "metric_strip_plus_chart": self._render_metric_strip_chart,
            "small_multiples_with_takeaway": self._render_full_width_chart,
            "single_bar_chart_with_insight": self._render_chart_left,
            "horizontal_bar_ranking": self._render_chart_right,
            "metric_cards_with_chart": self._render_metric_strip_chart,
            "comparison_chart_with_interpretation": self._render_chart_left,
            "distribution_with_callout": self._render_chart_right,
            "segment_profile_cards": self._render_three_cards,
            "three_finding_cards": self._render_three_cards,
            "comparison_matrix": self._render_comparison_matrix,
            "recommendation_priority": self._render_recommendations,
            "limitations_professional": self._render_limitations,
            "executive_summary_closing": self._render_closing,
        }
        renderers.get(spec.template, self._render_three_cards)(spec)
        if self.prs.slides:
            self._add_speaker_notes(self.prs.slides[-1], spec)

    def _add_speaker_notes(self, slide: Any, spec: SlideSpec) -> None:
        note = str(spec.speaker_note or "").strip()
        if not note:
            return
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = note
        except (AttributeError, TypeError, ValueError):
            # Older python-pptx versions may expose notes as read-only. The evidence
            # metadata remains available in slide_plan.json in that environment.
            return

    def _slide(self, background: tuple[int, int, int] | None = None) -> Any:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(background or self.theme.background)
        return slide

    def _rect(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        fill: tuple[int, int, int],
        *,
        radius: bool = False,
        line: tuple[int, int, int] | None = None,
    ) -> Any:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        if line:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = Pt(0.6)
        else:
            shape.line.fill.background()
        return shape

    def _text(
        self,
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        size: int = 12,
        color: tuple[int, int, int] | None = None,
        bold: bool = False,
        italic: bool = False,
        align: Any = PP_ALIGN.LEFT,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.06)
        frame.margin_right = Inches(0.06)
        frame.margin_top = Inches(0.03)
        frame.margin_bottom = Inches(0.03)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.name = self.theme.font_family
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color or self.theme.ink)
        return box

    def _header(self, slide: Any, spec: SlideSpec) -> None:
        font_size = 23 if len(spec.headline) <= 58 else 20 if len(spec.headline) <= 82 else 18
        self._text(slide, spec.headline, 0.72, 0.32, 11.75, 0.72, size=font_size, color=self.theme.navy, bold=True)
        self._rect(slide, 0.75, 1.08, 1.0, 0.035, self.theme.gold)

    def _footer(self, slide: Any, spec: SlideSpec) -> None:
        self._rect(slide, 0.72, 6.88, 11.9, 0.012, self.theme.line)
        evidence = ", ".join(spec.evidence_ids[:3] or spec.source_ids[:2])
        footer_text = shorten(self.deck.deck_title, 58)
        if evidence:
            footer_text += f"  |  Evidence: {evidence}"
        self._text(slide, shorten(footer_text, 105), 0.72, 7.03, 9.4, 0.25, size=8, color=self.theme.muted)
        self._text(slide, f"{spec.slide_number} / {len(self.deck.slides)}", 10.6, 7.03, 1.9, 0.25, size=8, color=self.theme.muted, align=PP_ALIGN.RIGHT)

    def _bullets(
        self,
        slide: Any,
        items: list[str],
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        size: int = 11,
        marker_color: tuple[int, int, int] | None = None,
        text_color: tuple[int, int, int] | None = None,
        max_chars: int = 145,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.04)
        frame.margin_right = Inches(0.04)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        for index, item in enumerate(items[:5]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(7)
            marker = paragraph.add_run()
            marker.text = "- "
            marker.font.name = self.theme.font_family
            marker.font.size = Pt(size)
            marker.font.bold = True
            marker.font.color.rgb = _rgb(marker_color or self.theme.blue)
            body = paragraph.add_run()
            body.text = shorten(item, max_chars)
            body.font.name = self.theme.font_family
            body.font.size = Pt(size)
            body.font.color.rgb = _rgb(text_color or self.theme.ink)

    def _message(self, slide: Any, spec: SlideSpec, left: float, top: float, width: float, height: float = 0.58) -> None:
        if not spec.main_message:
            return
        self._rect(slide, left, top, width, height, self.theme.pale_blue, radius=True, line=self.theme.line)
        self._text(slide, shorten(spec.main_message, 170), left + 0.14, top + 0.08, width - 0.28, height - 0.12, size=11, color=self.theme.navy, bold=True)

    def _visual(self, slide: Any, spec: SlideSpec, left: float, top: float, width: float, height: float) -> bool:
        if not spec.visual:
            return False
        self._rect(slide, left, top, width, height, self.theme.panel, radius=True, line=self.theme.line)
        if spec.visual.type in {"structured_chart", "chart"}:
            if spec.visual.fallback_path or spec.visual.image_path:
                image_path = resolve_image_fallback(spec.visual, self.theme)
                if image_path and os.path.exists(image_path):
                    metrics_text = _decision_tree_metric_text(spec.visual.data)
                    reserved_metric_h = 0.28 if spec.visual.chart_type == "decision_tree" and metrics_text else 0.0
                    image_width, image_height = _fit_image_size(image_path, width - 0.3, height - 0.35 - reserved_metric_h)
                    x = left + max((width - image_width) / 2, 0)
                    y = top + max((height - reserved_metric_h - image_height) / 2, 0)
                    picture = slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(image_width), Inches(image_height))
                    picture.name = "shared_decision_tree_image" if spec.visual.chart_type == "decision_tree" else "shared_python_eda_image"
                    if spec.visual.chart_type == "decision_tree" and metrics_text:
                        self._text(
                            slide,
                            metrics_text,
                            left + 0.18,
                            top + height - 0.3,
                            width - 0.36,
                            0.18,
                            size=7,
                            color=self.theme.muted,
                            align=PP_ALIGN.CENTER,
                        )
                    return True
            if add_structured_chart(slide, spec.visual, left + 0.15, top + 0.18, width - 0.3, height - 0.38, self.theme):
                return True
            self._chart_unavailable(slide, spec, left, top, width, height)
            return False
        if spec.visual.type == "native_chart" and add_native_chart(slide, spec.visual, left + 0.15, top + 0.18, width - 0.3, height - 0.38, self.theme):
            return True
        image_path = resolve_image_fallback(spec.visual, self.theme)
        if image_path and os.path.exists(image_path):
            metrics_text = _decision_tree_metric_text(spec.visual.data) if spec.visual.chart_type == "decision_tree" else ""
            reserved_metric_h = 0.28 if metrics_text else 0.0
            image_width, image_height = _fit_image_size(image_path, width - 0.3, height - 0.35 - reserved_metric_h)
            x = left + max((width - image_width) / 2, 0)
            y = top + max((height - reserved_metric_h - image_height) / 2, 0)
            slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(image_width), Inches(image_height))
            if metrics_text:
                self._text(
                    slide,
                    metrics_text,
                    left + 0.18,
                    top + height - 0.3,
                    width - 0.36,
                    0.18,
                    size=7,
                    color=self.theme.muted,
                    align=PP_ALIGN.CENTER,
                )
            return True
        return False

    def _chart_unavailable(self, slide: Any, spec: SlideSpec, left: float, top: float, width: float, height: float) -> None:
        self._rect(slide, left + 0.28, top + 0.3, width - 0.56, height - 0.6, self.theme.pale_blue, radius=True, line=self.theme.line)
        self._text(
            slide,
            "Structured chart data was unavailable for rendering",
            left + 0.55,
            top + height / 2 - 0.18,
            width - 1.1,
            0.32,
            size=11,
            color=self.theme.navy,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _render_title_cover(self, spec: SlideSpec) -> None:
        slide = self._slide(self.theme.navy)
        self._rect(slide, 0, 0, 0.46, self.theme.height, self.theme.gold)
        self._rect(slide, 9.9, 0, 3.43, self.theme.height, self.theme.blue)
        self._text(slide, "MULTI-AGENT ANALYTICS SYSTEM", 0.82, 0.58, 5.8, 0.3, size=9, color=self.theme.pale_blue, bold=True)
        self._rect(slide, 0.86, 2.34, 1.26, 0.04, self.theme.gold)
        title_size = 33 if len(spec.headline) <= 55 else 28
        self._text(slide, spec.headline, 0.82, 2.6, 8.45, 1.45, size=title_size, color=self.theme.white, bold=True)
        self._text(slide, spec.subtitle or self.deck.subtitle, 0.84, 4.55, 7.9, 0.85, size=14, color=self.theme.pale_blue)
        meta = self.deck.metadata.get("generated_on") or datetime.now().strftime("%B %d, %Y")
        self._text(slide, meta, 0.84, 6.72, 4.3, 0.28, size=9, color=self.theme.pale_blue)
        self._text(slide, "EXECUTIVE\nDECISION DECK", 10.18, 2.78, 2.55, 1.0, size=17, color=self.theme.white, bold=True, align=PP_ALIGN.CENTER)

    def _render_data_understanding(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.32, 11.88)
        metrics = spec.metrics[:4] or [{"label": "Dataset", "value": self.deck.dataset_context.get("name", "Available")}]
        for index, metric in enumerate(metrics[:4]):
            left = 0.82 + index * 3.0
            self._rect(slide, left, 2.08, 2.55, 1.05, self.theme.panel, radius=True, line=self.theme.line)
            self._text(slide, shorten(metric.get("value", ""), 26), left + 0.12, 2.28, 2.3, 0.34, size=18, color=self.theme.blue, bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, shorten(metric.get("label", ""), 34), left + 0.12, 2.66, 2.3, 0.28, size=9, color=self.theme.muted, align=PP_ALIGN.CENTER)
        self._rect(slide, 0.82, 3.55, 7.25, 2.55, self.theme.panel, radius=True, line=self.theme.line)
        self._text(slide, "What the data represents", 1.05, 3.82, 4.6, 0.3, size=13, color=self.theme.navy, bold=True)
        self._bullets(slide, spec.bullets, 1.02, 4.28, 6.72, 1.45, size=10)
        self._rect(slide, 8.45, 3.55, 3.95, 2.55, self.theme.pale_gold, radius=True, line=self.theme.line)
        self._text(slide, "Quality note", 8.7, 3.82, 3.35, 0.3, size=13, color=self.theme.navy, bold=True)
        quality = self.deck.dataset_context.get("data_quality_notes", [])
        self._bullets(slide, quality or ["No major structured quality notes were captured."], 8.68, 4.28, 3.35, 1.45, size=10)
        self._footer(slide, spec)

    def _render_market_context(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.32, 11.88)
        self._rect(slide, 0.88, 2.22, 7.1, 3.85, self.theme.panel, radius=True, line=self.theme.line)
        self._bullets(slide, spec.bullets, 1.14, 2.58, 6.55, 2.9, size=11, max_chars=210)
        self._rect(slide, 8.35, 2.22, 4.02, 3.85, self.theme.pale_blue, radius=True, line=self.theme.line)
        self._text(slide, "Why it matters", 8.62, 2.58, 3.4, 0.32, size=14, color=self.theme.navy, bold=True)
        self._text(slide, shorten(spec.main_message, 260), 8.62, 3.08, 3.42, 1.5, size=12, color=self.theme.ink)
        self._footer(slide, spec)

    def _render_chart_left(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._visual(slide, spec, 0.82, 1.5, 7.2, 4.95)
        self._rect(slide, 8.35, 1.5, 4.0, 4.95, self.theme.panel, radius=True, line=self.theme.line)
        self._text(slide, "Insight", 8.62, 1.82, 3.35, 0.35, size=14, color=self.theme.navy, bold=True)
        self._text(slide, shorten(spec.main_message, 190), 8.62, 2.28, 3.34, 0.92, size=12, color=self.theme.ink, bold=True)
        self._bullets(slide, spec.bullets, 8.58, 3.42, 3.42, 2.1, size=10)
        self._footer(slide, spec)

    def _render_chart_right(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._rect(slide, 0.82, 1.5, 4.0, 4.95, self.theme.panel, radius=True, line=self.theme.line)
        self._text(slide, "Readout", 1.08, 1.82, 3.35, 0.35, size=14, color=self.theme.navy, bold=True)
        self._text(slide, shorten(spec.main_message, 190), 1.08, 2.28, 3.34, 0.92, size=12, color=self.theme.ink, bold=True)
        self._bullets(slide, spec.bullets, 1.04, 3.42, 3.42, 2.1, size=10)
        self._visual(slide, spec, 5.15, 1.5, 7.2, 4.95)
        self._footer(slide, spec)

    def _render_full_width_chart(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._visual(slide, spec, 0.9, 1.42, 11.45, 4.58)
        takeaway = spec.visual.takeaway if spec.visual and spec.visual.takeaway else spec.main_message
        self._rect(slide, 1.05, 6.14, 11.1, 0.48, self.theme.pale_blue, radius=True, line=self.theme.line)
        self._text(slide, shorten(takeaway, 180), 1.22, 6.24, 10.7, 0.24, size=11, color=self.theme.navy, bold=True, align=PP_ALIGN.CENTER)
        self._footer(slide, spec)

    def _render_metric_strip_chart(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        metrics = spec.metrics[:4]
        for index, metric in enumerate(metrics):
            left = 0.86 + index * 3.0
            self._rect(slide, left, 1.35, 2.55, 0.82, self.theme.panel, radius=True, line=self.theme.line)
            self._text(slide, shorten(metric.get("value", ""), 24), left + 0.12, 1.49, 2.28, 0.28, size=16, color=self.theme.blue, bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, shorten(metric.get("label", ""), 32), left + 0.12, 1.8, 2.28, 0.2, size=8, color=self.theme.muted, align=PP_ALIGN.CENTER)
        chart_top = 2.35 if metrics else 1.5
        chart_height = 4.1 if metrics else 4.95
        self._visual(slide, spec, 0.92, chart_top, 11.38, chart_height)
        self._footer(slide, spec)

    def _render_three_cards(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.32, 11.88)
        items = spec.bullets[:3] or [spec.main_message]
        finding_start = 4 if spec.slide_role == "business_translation" and spec.template == "three_finding_cards" else 1
        for index, item in enumerate(items[:3]):
            left = 0.88 + index * 3.9
            self._rect(slide, left, 2.45, 3.45, 3.35, self.theme.panel, radius=True, line=self.theme.line)
            self._rect(slide, left, 2.45, 0.08, 3.35, self.theme.chart_palette[index % len(self.theme.chart_palette)])
            self._text(slide, f"Finding {finding_start + index}", left + 0.25, 2.82, 2.85, 0.28, size=11, color=self.theme.muted, bold=True)
            self._text(slide, shorten(item, 170), left + 0.25, 3.28, 2.85, 1.4, size=12, color=self.theme.ink, bold=index == 0)
        self._footer(slide, spec)

    def _render_comparison_matrix(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.32, 11.88)
        items = spec.bullets[:4] or [spec.main_message]
        labels = (
            ["Decision use", "Tradeoff", "Operating rule", "Validation"]
            if spec.slide_role == "business_translation"
            else ["Evidence", "Implication", "Signal", "Next check"]
        )
        for index in range(4):
            left = 0.9 + (index % 2) * 5.85
            top = 2.25 + (index // 2) * 1.82
            self._rect(slide, left, top, 5.28, 1.35, self.theme.panel, radius=True, line=self.theme.line)
            self._text(slide, labels[index], left + 0.2, top + 0.18, 4.8, 0.24, size=10, color=self.theme.blue, bold=True)
            item = items[index] if index < len(items) else ""
            item = _strip_leading_label(item, labels[index])
            self._text(slide, shorten(item, 118), left + 0.2, top + 0.55, 4.8, 0.55, size=11, color=self.theme.ink)
        self._footer(slide, spec)

    def _render_recommendations(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.22, 11.88, 0.52)
        items = spec.bullets[:3] or [spec.main_message]
        labels = ["Priority 1", "Priority 2", "Priority 3"]
        for index, item in enumerate(items[:3]):
            left = 0.88 + index * 3.9
            self._rect(slide, left, 2.2, 3.45, 3.65, self.theme.panel, radius=True, line=self.theme.line)
            self._rect(slide, left, 2.2, 3.45, 0.12, self.theme.chart_palette[index % len(self.theme.chart_palette)])
            self._text(slide, labels[index], left + 0.22, 2.55, 2.95, 0.28, size=11, color=self.theme.blue, bold=True)
            pieces = [piece for piece in item.split(". ") if piece.strip()]
            action = pieces[0] if pieces else item
            rationale = ". ".join(pieces[1:]) if len(pieces) > 1 else self._recommendation_validation_fallback(spec, item)
            fields = _structured_recommendation_fields(item)
            if fields:
                action = f"Action: {fields.get('action', action)}"
                rationale = " | ".join(
                    part
                    for part in (
                        f"Owner: {shorten(fields.get('owner'), 30, placeholder='')}" if fields.get("owner") else "",
                        f"Trigger: {shorten(fields.get('trigger'), 52, placeholder='')}" if fields.get("trigger") else "",
                        f"Guardrail: {shorten(fields.get('guardrail'), 55, placeholder='')}" if fields.get("guardrail") else "",
                        f"Risk: {shorten(fields.get('risk') or fields.get('caveat'), 42, placeholder='')}" if fields.get("risk") or fields.get("caveat") else "",
                    )
                    if part
                ) or rationale
            self._text(slide, shorten(action, 88, placeholder=""), left + 0.22, 3.0, 2.95, 0.62, size=12, color=self.theme.navy, bold=True)
            self._text(slide, shorten(rationale, 170, placeholder=""), left + 0.22, 3.82, 2.95, 1.22, size=9, color=self.theme.ink)
            footer = "Trigger + guardrail required" if "trigger" in item.lower() or "guardrail" in item.lower() else "Validation required"
            self._text(slide, footer, left + 0.22, 5.35, 2.95, 0.24, size=8, color=self.theme.muted, bold=True)
        self._footer(slide, spec)

    def _recommendation_validation_fallback(self, spec: SlideSpec, item: str) -> str:
        context = " ".join(
            [
                spec.headline,
                spec.main_message,
                spec.subtitle,
                item,
                " ".join(spec.bullets),
            ]
        ).lower()
        if any(token in context for token in ("attrition", "retention", "employee", "workload", "overtime")):
            return "Track attrition, workload, engagement, and manager capacity before scaling."
        if any(token in context for token in ("loan", "credit", "default", "borrower", "approval")):
            return "Track default capture, false-positive review rate, approval impact, and fairness before scaling."
        if any(token in context for token in ("mental health", "depression", "student", "screening")):
            return "Track support uptake, false positives, escalation quality, and human review before scaling."
        if any(token in context for token in ("stock", "price", "return", "drawdown", "volatility")):
            return "Track forward return, drawdown, volatility, and liquidity before acting."
        return "Track pilot outcome, adoption, guardrails, and stop/go criteria before scaling."

    def _render_limitations(self, spec: SlideSpec) -> None:
        slide = self._slide()
        self._header(slide, spec)
        self._message(slide, spec, 0.72, 1.32, 11.88)
        items = spec.bullets[:4]
        for index, item in enumerate(items):
            left = 0.9 + (index % 2) * 5.85
            top = 2.32 + (index // 2) * 1.72
            self._rect(slide, left, top, 5.28, 1.22, self.theme.panel, radius=True, line=self.theme.line)
            if spec.slide_role == "sources":
                label = "Source" if str(item).lower().startswith(("dataset:", "[")) else "Evidence"
                label_color = self.theme.blue
            elif str(item).lstrip().lower().startswith("mitigation:"):
                label = "Mitigation"
                label_color = self.theme.teal
            else:
                label = "Limitation"
                label_color = self.theme.red
            self._text(slide, label, left + 0.2, top + 0.17, 1.4, 0.23, size=9, color=label_color, bold=True)
            self._text(slide, shorten(item, 118, placeholder=""), left + 0.2, top + 0.52, 4.8, 0.46, size=10, color=self.theme.ink)
        self._footer(slide, spec)

    def _render_closing(self, spec: SlideSpec) -> None:
        slide = self._slide(self.theme.navy)
        self._rect(slide, 0, 0, 0.46, self.theme.height, self.theme.gold)
        headline = shorten(spec.headline, 88, placeholder="")
        self._text(
            slide, headline, 0.82, 1.08, 11.1, 1.22,
            size=27 if len(headline) > 58 else 31, color=self.theme.white, bold=True,
        )
        self._rect(slide, 0.86, 2.48, 10.9, 0.04, self.theme.gold)
        self._text(slide, shorten(spec.main_message, 210), 0.84, 2.76, 10.4, 1.0, size=17, color=self.theme.pale_blue)
        self._bullets(
            slide,
            spec.bullets,
            0.9,
            4.2,
            9.8,
            1.4,
            size=12,
            marker_color=self.theme.gold,
            text_color=self.theme.white,
        )
        self._text(slide, "Generated by the Multi-Agent Analytics System", 0.84, 6.82, 6.8, 0.25, size=9, color=self.theme.pale_blue)
        self._text(slide, f"{spec.slide_number} / {len(self.deck.slides)}", 10.65, 6.82, 1.8, 0.25, size=9, color=self.theme.pale_blue, align=PP_ALIGN.RIGHT)


def _strip_leading_label(text: str, label: str) -> str:
    clean = shorten(text, 180, placeholder="")
    prefix = f"{label}:"
    if clean.lower().startswith(prefix.lower()):
        return clean[len(prefix):].strip()
    return clean


def _structured_recommendation_fields(text: str) -> dict[str, str]:
    labels = ("Action", "Owner", "Trigger", "Target", "Timeline", "Impact", "Guardrail", "Risk", "Caveat", "Evidence", "Validate")
    pattern = r"\b(" + "|".join(labels) + r"):\s*"
    matches = list(re.finditer(pattern, text, flags=re.IGNORECASE))
    fields: dict[str, str] = {}
    for index, match in enumerate(matches):
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        value = text[start:end].strip(" .;|")
        if value:
            key = match.group(1).lower()
            if key == "validate":
                key = "guardrail"
            fields[key] = value
    return fields


def _preferred_output_path(output_path: str) -> str:
    path = Path(output_path)
    if not path.exists():
        return str(path)
    try:
        with open(path, "ab"):
            return str(path)
    except OSError:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return str(path.with_name(f"{path.stem}_{timestamp}{path.suffix}"))


def _fit_image_size(path: str, max_width: float, max_height: float) -> tuple[float, float]:
    if PILImage is None:
        return max_width, max_height
    try:
        with PILImage.open(path) as image:
            width, height = image.size
    except OSError:
        return max_width, max_height
    if not width or not height:
        return max_width, max_height
    ratio = min(max_width / width, max_height / height)
    return width * ratio, height * ratio


def _decision_tree_metric_text(data: Any) -> str:
    if not isinstance(data, dict):
        return ""
    metric_bits = []
    for key in (
        "train_accuracy",
        "test_accuracy",
        "train_score",
        "test_score",
        "accuracy",
        "baseline_accuracy",
        "baseline_score",
        "precision",
        "recall",
        "f1",
        "f1_score",
        "support",
        "positive_class_rate",
        "train_r2",
        "test_r2",
        "r2",
        "train_mae",
        "test_mae",
        "mae",
    ):
        value = data.get(key)
        if value not in (None, ""):
            metric_bits.append(f"{_decision_tree_metric_label(key, data)}: {_decision_tree_metric_value(key, value)}")
    return " | ".join(metric_bits[:4])


def _decision_tree_metric_label(key: str, data: dict[str, Any]) -> str:
    model_type = str(data.get("model_type", "")).lower()
    if key == "train_score" and not model_type.startswith("reg"):
        return "Train Accuracy"
    if key == "test_score" and not model_type.startswith("reg"):
        return "Test Accuracy"
    if key == "baseline_score" and not model_type.startswith("reg"):
        return "Baseline Accuracy"
    return key.replace("_", " ").title()


def _decision_tree_metric_value(key: str, value: Any) -> str:
    if key in {"train_score", "test_score", "baseline_score", "train_accuracy", "test_accuracy", "baseline_accuracy", "accuracy", "precision", "recall", "f1", "f1_score", "positive_class_rate"}:
        try:
            number = float(value)
        except (TypeError, ValueError):
            return str(value)
        if 0 <= number <= 1:
            return f"{number * 100:.1f}%"
    return str(value)


def _render_text_fallback(deck: DeckSpec, output_path: str) -> str:
    path = Path(output_path)
    fallback_path = _preferred_output_path(str(path.with_name(f"{path.stem}_slide_deck_fallback.txt")))
    lines = [deck.deck_title, deck.subtitle, ""]
    for slide in deck.slides:
        lines.append(f"Slide {slide.slide_number}: {slide.headline}")
        lines.append(f"Role: {slide.slide_role}")
        lines.append(f"Template: {slide.template}")
        lines.append(f"Message: {slide.main_message}")
        for bullet in slide.bullets:
            lines.append(f"- {textwrap.fill(bullet, width=100)}")
        if slide.visual:
            lines.append(f"Visual: {slide.visual.type} {slide.visual.chart_type or slide.visual.image_path}")
        lines.append("")
    Path(fallback_path).write_text("\n".join(lines), encoding="utf-8")
    return fallback_path


def _rgb(value: tuple[int, int, int]) -> Any:
    return RGBColor(*value)
