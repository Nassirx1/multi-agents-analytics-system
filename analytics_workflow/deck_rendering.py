from __future__ import annotations

import json
import logging
import shutil
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation  # noqa: F401

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    PPTX_AVAILABLE = False

from .slides import build_deck_spec
from .slides.validators import fatal_deck_issues
from .slides.templates import LEGACY_LAYOUT_TO_TEMPLATE, TEMPLATE_REGISTRY
from .slides.text_refiner import refine_bullets, refine_headline, shorten


ALLOWED_LAYOUT_TYPES = set(LEGACY_LAYOUT_TO_TEMPLATE) | set(TEMPLATE_REGISTRY)
MAX_TITLE_CHARS = 96
MAX_BULLETS = 5
MAX_DETAIL_CHARS = 150


def validate_slide(slide: dict[str, Any], saved_figures: list[str] | None = None) -> list[str]:
    issues: list[str] = []
    template_or_layout = slide.get("template") or slide.get("layout_type")
    if template_or_layout not in ALLOWED_LAYOUT_TYPES:
        issues.append("invalid_layout_type")
    if len(str(slide.get("headline") or slide.get("title") or "")) > MAX_TITLE_CHARS:
        issues.append("title_too_long")
    if not str(slide.get("main_message", "")).strip():
        issues.append("missing_main_message")
    details = slide.get("details", []) or []
    if isinstance(details, list) and len(details) > MAX_BULLETS:
        issues.append("too_many_bullets")
    if isinstance(details, list) and any(len(str(detail)) > MAX_DETAIL_CHARS for detail in details):
        issues.append("details_too_verbose")
    visual_path = str(slide.get("visual_path") or slide.get("visual_element") or "").strip()
    if visual_path and saved_figures is not None and visual_path not in saved_figures:
        issues.append("missing_visual_file")
    return issues


def repair_slide(slide: dict[str, Any], issues: list[str], index: int) -> dict[str, Any]:
    repaired = dict(slide)
    if "invalid_layout_type" in issues:
        repaired["layout_type"] = "three_finding_cards"
    title = repaired.get("headline") or repaired.get("title") or f"Slide {index + 1}"
    repaired["title"] = refine_headline(title, f"Slide {index + 1}", width=MAX_TITLE_CHARS)
    repaired["headline"] = repaired["title"]
    if "missing_main_message" in issues:
        details = repaired.get("details", []) or []
        repaired["main_message"] = shorten(details[0] if details else repaired["title"], 170)
    if "too_many_bullets" in issues or "details_too_verbose" in issues:
        repaired["details"] = refine_bullets(repaired.get("details", []), max_items=MAX_BULLETS, max_chars=MAX_DETAIL_CHARS)
    if "missing_visual_file" in issues:
        repaired["visual_path"] = ""
        repaired["visual_element"] = ""
    return repaired


def normalize_slide_plan(workflow_state: dict[str, Any]) -> list[dict[str, Any]]:
    deck = build_deck_spec(workflow_state)
    # Legacy callers historically received planned content slides, not the cover.
    return [slide.to_legacy_dict() for slide in deck.slides if slide.slide_role != "title"]


def build_consulting_deck(
    workflow_state: dict[str, Any],
    output_path: str = "analytics_report.pptx",
    *,
    runtime_config: Any | None = None,
) -> str:
    deck = build_deck_spec(workflow_state)
    for warning in workflow_state.get("slide_validation_warnings", []) or []:
        if isinstance(warning, dict) and warning.get("scope") in {"slide", "deck"}:
            try:
                logging.getLogger("SlideDeckGenerator").warning("Slide validation warning: %s", warning)
            except OSError:
                pass
    fatal_issues = fatal_deck_issues(list(workflow_state.get("slide_unresolved_warnings", []) or []))
    enforce_quality_gate = bool(
        workflow_state.get(
            "enforce_deck_quality_gate",
            (workflow_state.get("analysis_results", {}) or {}).get("execution_status") == "success",
        )
    )
    if fatal_issues and enforce_quality_gate:
        raise RuntimeError(
            "Slide deck failed the export quality gate: "
            + "; ".join(
                f"slide {issue.get('slide', '?')}: {issue.get('issue', 'unknown_issue')}"
                for issue in fatal_issues[:12]
            )
        )
    from .presentation_backends import (
        _apply_hybrid_story_contract,
        PowerPointMCPBackend,
        PythonPresentationBackend,
        backend_from_config,
        enrich_deck_executive_copy,
        inspect_presentation,
    )
    from .runtime_config import get_active_runtime_config

    config = runtime_config if runtime_config is not None else get_active_runtime_config()
    openrouter_client = workflow_state.get("_presentation_openrouter_client")
    backend = backend_from_config(config, openrouter_client) if config is not None else PythonPresentationBackend()
    enrich_deck_executive_copy(deck, workflow_state.get("agent_outputs", {}) or {})
    evidence_hash = str((workflow_state.get("evidence_bundle", {}) or {}).get("bundle_hash", ""))
    if evidence_hash or str((deck.metadata or {}).get("evidence_bundle_hash", "")):
        _apply_hybrid_story_contract(deck)
    deck.renumber()
    workflow_state["deck_spec"] = deck.to_dict()
    backend_log = workflow_state.setdefault("presentation_backend_log", [])
    try:
        deck_path = backend.render(deck, output_path, workflow_state=workflow_state)
        deck_path = _materialize_deck_in_run_directory(deck_path, output_path)
        inspection = inspect_presentation(deck_path, expected_deck=deck)
        inspection.rendered_files.extend(getattr(backend, "last_rendered_files", []) or [])
        if not inspection.valid:
            raise RuntimeError(
                "Presentation backend QA failed: "
                + "; ".join(str(issue.get("issue")) for issue in inspection.issues[:12])
            )
        backend_log.append({"backend": backend.name, "status": "success", "qa": inspection.to_dict()})
        workflow_state["presentation_backend_used"] = backend.name
    except Exception as exc:
        if not isinstance(backend, PowerPointMCPBackend):
            raise
        error_detail = _exception_detail(exc)
        logging.getLogger("SlideDeckGenerator").warning(
            "PowerPoint MCP failed; switching automatically to Python fallback: %s", error_detail
        )
        backend_log.append({"backend": backend.name, "status": "failed", "error": error_detail})
        fallback = PythonPresentationBackend()
        deck_path = fallback.render(deck, output_path, workflow_state=workflow_state)
        deck_path = _materialize_deck_in_run_directory(deck_path, output_path)
        inspection = inspect_presentation(deck_path, expected_deck=deck)
        if not inspection.valid:
            raise RuntimeError(
                "Python presentation fallback also failed QA: "
                + "; ".join(str(issue.get("issue")) for issue in inspection.issues[:12])
            )
        backend_log.append({"backend": fallback.name, "status": "success", "qa": inspection.to_dict()})
        workflow_state["presentation_backend_used"] = fallback.name
        workflow_state.setdefault("run_manifest", {}).setdefault("warnings", []).append(
            f"PowerPoint MCP failed; Python presentation fallback used: {error_detail}"
        )
    _save_deck_artifacts(workflow_state, deck_path, deck.to_dict())
    return deck_path


def _exception_detail(exc: BaseException) -> str:
    messages: list[str] = []

    def visit(item: BaseException) -> None:
        nested = getattr(item, "exceptions", None)
        if isinstance(nested, (list, tuple)):
            for child in nested:
                if isinstance(child, BaseException):
                    visit(child)
            return
        message = str(item).strip()
        if message and message not in messages:
            messages.append(message)

    visit(exc)
    return " | ".join(messages) or f"{type(exc).__name__}: no detail"


def _materialize_deck_in_run_directory(deck_path: str, requested_output_path: str) -> str:
    """Guarantee that every presentation backend leaves its artifact in the run folder."""
    source = Path(deck_path).resolve()
    requested = Path(requested_output_path).resolve()
    requested.parent.mkdir(parents=True, exist_ok=True)
    if not source.is_file():
        raise RuntimeError(f"Presentation backend returned a missing deck: {source}")
    if source == requested:
        return str(source)

    destination = requested
    try:
        shutil.copy2(source, destination)
    except PermissionError:
        # Match the deterministic renderer's locked-file behavior while keeping
        # the artifact inside the active run directory.
        from datetime import datetime

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        destination = requested.with_name(f"{requested.stem}_{timestamp}{requested.suffix}")
        shutil.copy2(source, destination)
    return str(destination.resolve())


def _save_deck_artifacts(workflow_state: dict[str, Any], deck_path: str, deck_payload: dict[str, Any]) -> None:
    output_dir = Path(deck_path).parent
    slide_plan_path = output_dir / "slide_plan.json"
    chart_specs_path = output_dir / "chart_specs.json"
    analysis_results = workflow_state.get("analysis_results", {}) or {}
    chart_payload = {
        "artifact_type": "chart_specs",
        "deck_path": deck_path,
        "chart_specs": analysis_results.get("chart_specs", []) or [],
        "analysis_artifacts": analysis_results.get("analysis_artifacts", []) or [],
        "saved_figures": workflow_state.get("saved_figures", []) or [],
        "slide_visuals": [
            {
                "slide_number": slide.get("slide_number"),
                "slide_role": slide.get("slide_role"),
                "template": slide.get("template"),
                "visual": slide.get("visual"),
            }
            for slide in deck_payload.get("slides", [])
            if slide.get("visual")
        ],
    }
    slide_plan_path.write_text(json.dumps(deck_payload, indent=2, default=str), encoding="utf-8")
    chart_specs_path.write_text(json.dumps(chart_payload, indent=2, default=str), encoding="utf-8")
    workflow_state.setdefault("generated_reports", {})["slide_plan"] = str(slide_plan_path)
    workflow_state.setdefault("generated_reports", {})["chart_specs"] = str(chart_specs_path)
